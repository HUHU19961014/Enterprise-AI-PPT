import datetime
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from sie_autoppt.config import DEFAULT_HTML, DEFAULT_OUTPUT_DIR


def _extract_single(html: str, class_name: str) -> str:
    pattern = rf'<div class="{class_name}">(.*?)</div>'
    match = re.search(pattern, html, flags=re.S)
    return re.sub(r"<.*?>", "", match.group(1)).strip() if match else ""


def _extract_list(html: str, class_name: str) -> list[str]:
    pattern = rf'<div class="{class_name}">(.*?)</div>'
    return [re.sub(r"<.*?>", "", m).strip() for m in re.findall(pattern, html, flags=re.S)]


def _extract_phases(html: str) -> list[dict]:
    blocks = re.findall(r'<div class="phase">(.*?)</div>\s*</div>?', html, flags=re.S)
    phases = []
    for b in blocks:
        phases.append(
            {
                "time": _extract_single(b, "phase-time"),
                "name": _extract_single(b, "phase-name"),
                "code": _extract_single(b, "phase-code"),
                "func": _extract_single(b, "phase-func"),
                "owner": _extract_single(b, "phase-owner"),
            }
        )
    return [p for p in phases if p["name"]]


def _set_text(shape, text, size=12, bold=False, color=RGBColor(44, 62, 55), align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def build_ppt(html_path: Path) -> Path:
    html = html_path.read_text(encoding="utf-8")

    title = _extract_single(html, "title") or "UAT测试计划与责任分工"
    subtitle = _extract_single(html, "subtitle")
    phases = _extract_phases(html)
    scenarios = _extract_list(html, "scenario")
    notes = _extract_list(html, "note")
    footer = _extract_single(html, "footer")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 251, 250)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(5.6), Inches(0.35))
    _set_text(title_box, title, size=20, bold=True, color=RGBColor(20, 79, 66))
    sub_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.53), Inches(8.8), Inches(0.26))
    _set_text(sub_box, subtitle, size=10, color=RGBColor(108, 127, 120))

    # 四阶段卡片（PPT化优化：略压缩高度提升可读性）
    start_x, start_y = 0.4, 0.85
    gap = 0.14
    card_w = (10 - 0.8 - gap * 3) / 4
    card_h = 1.58
    for i, p in enumerate(phases[:4]):
        x = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(start_y), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = RGBColor(220, 233, 228)

        y0 = start_y + 0.08
        t = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y0), Inches(card_w - 0.16), Inches(0.18))
        _set_text(t, p["time"], size=8, bold=True, color=RGBColor(13, 107, 87))
        n = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y0 + 0.2), Inches(card_w - 0.16), Inches(0.24))
        _set_text(n, p["name"], size=13, bold=True)
        c = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y0 + 0.44), Inches(card_w - 0.16), Inches(0.2))
        _set_text(c, p["code"], size=8, color=RGBColor(91, 111, 104))
        f = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y0 + 0.62), Inches(card_w - 0.16), Inches(0.44))
        _set_text(f, p["func"], size=8, color=RGBColor(47, 79, 70))
        o = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y0 + 1.11), Inches(card_w - 0.16), Inches(0.28))
        _set_text(o, p["owner"], size=8, color=RGBColor(47, 79, 70))

    # 下半区
    left_x, right_x = 0.4, 6.35
    bottom_y = 2.56
    left_w, right_w = 5.75, 3.25
    panel_h = 2.28
    for x, w in [(left_x, left_w), (right_x, right_w)]:
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(bottom_y), Inches(w), Inches(panel_h))
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        panel.line.color.rgb = RGBColor(220, 233, 228)

    ltitle = slide.shapes.add_textbox(Inches(left_x + 0.12), Inches(bottom_y + 0.08), Inches(1.5), Inches(0.22))
    _set_text(ltitle, "测试范围", size=11, bold=True, color=RGBColor(13, 107, 87))
    rtitle = slide.shapes.add_textbox(Inches(right_x + 0.12), Inches(bottom_y + 0.08), Inches(1.7), Inches(0.22))
    _set_text(rtitle, "测试关注点", size=11, bold=True, color=RGBColor(13, 107, 87))

    # 左侧 2x2 场景卡
    sx, sy = left_x + 0.12, bottom_y + 0.35
    sgap_x, sgap_y = 0.12, 0.1
    sc_w = (left_w - 0.24 - sgap_x) / 2
    sc_h = 0.72
    for i, txt in enumerate(scenarios[:4]):
        row, col = divmod(i, 2)
        x = sx + col * (sc_w + sgap_x)
        y = sy + row * (sc_h + sgap_y)
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(sc_w), Inches(sc_h))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(246, 251, 249)
        s.line.color.rgb = RGBColor(226, 239, 234)
        tb = slide.shapes.add_textbox(Inches(x + 0.08), Inches(y + 0.2), Inches(sc_w - 0.16), Inches(0.3))
        _set_text(tb, txt, size=9)

    # 右侧关注点
    ny = bottom_y + 0.35
    for txt in notes[:4]:
        nbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(right_x + 0.12), Inches(ny), Inches(right_w - 0.24), Inches(0.42))
        nbox.fill.solid()
        nbox.fill.fore_color.rgb = RGBColor(245, 250, 248)
        nbox.line.color.rgb = RGBColor(226, 239, 234)
        tb = slide.shapes.add_textbox(Inches(right_x + 0.2), Inches(ny + 0.11), Inches(right_w - 0.4), Inches(0.2))
        _set_text(tb, txt, size=9)
        ny += 0.5

    footer_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.01), Inches(9.2), Inches(0.38))
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = RGBColor(234, 245, 241)
    footer_shape.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(5.01), Inches(0.04), Inches(0.38))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(13, 107, 87)
    accent.line.fill.background()
    footer_tb = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(9.0), Inches(0.2))
    _set_text(footer_tb, footer, size=9, color=RGBColor(47, 79, 70))

    DEFAULT_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
    out_path = DEFAULT_OUTPUT_DIR / f"SIE_UAT_HTML_Rebuild_{timestamp}.pptx"
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    input_html = DEFAULT_HTML
    output = build_ppt(input_html)
    print(str(output))
