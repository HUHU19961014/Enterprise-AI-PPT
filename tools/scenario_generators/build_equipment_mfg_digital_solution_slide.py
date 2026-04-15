from __future__ import annotations

import argparse
import datetime
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

try:
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest
except ModuleNotFoundError:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest


PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_PPT = PROJECT_ROOT / "output" / "equipment_mfg_digital_solution_onepage.pptx"

ACCENT = (173, 5, 61)
INK = (35, 41, 46)
MUTED = (98, 109, 121)
LIGHT_TEXT = (146, 154, 161)
LINE = (224, 228, 232)
WHITE = (255, 255, 255)


@dataclass(frozen=True)
class CardSpec:
    index: str
    key: str
    title: str
    english: str
    accent: tuple[int, int, int]
    strip: tuple[int, int, int]
    summary: str
    bullets: tuple[str, str, str]
    closing: str


TITLE_SEGMENTS = [
    ("装备制造行业", INK, True),
    ("数字化转型", ACCENT, True),
    ("解决方案", INK, True),
]
SUBTITLE = "围绕非标定制、复杂协同与交付压力，数字化建设的关键不是单点上系统，而是打通从设计到售后的全链路经营能力。"
HEADLINE = "装备制造企业真正要解决的，不只是生产效率问题，而是如何在复杂产品与定制化需求下，实现设计、计划、供应链、制造和服务的协同闭环。"
SUPPORT = "建议从核心痛点出发，构建数据底座与经营蓝图，再按阶段推进业务数字化、数字工厂、服务化与智能柔性制造。"
FOOTNOTE = "一句话总结：装备制造业数字化转型的本质，是通过数据驱动打通设计到售后全链路，实现研发协同化、生产精益化、管控实时化、服务智能化。"

CARDS = (
    CardSpec(
        index="01",
        key="痛",
        title="核心痛点",
        english="Pain Points",
        accent=(118, 86, 167),
        strip=(248, 243, 250),
        summary="问题根源在于复杂产品模式下的多线并行协同。",
        bullets=(
            "研产协同复杂：边设计边采购边生产，变更频繁且影响面广。",
            "计划与物料难控：计划动态变化快，库存、齐套与供应协同压力大。",
            "制造与服务脱节：数据孤岛多，质量闭环慢，售后响应成本高。"
        ),
        closing="判断：如果链路不通，计划波动、质量问题和交付压力会持续放大。",
    ),
    CardSpec(
        index="02",
        key="图",
        title="方案蓝图",
        english="Blueprint",
        accent=ACCENT,
        strip=(252, 241, 245),
        summary="先建设统一数据底座，再形成经营协同与决策闭环。",
        bullets=(
            "以数据采集与集成为底座，打破设计、计划、制造和服务之间的数据壁垒。",
            "围绕供应链、制造、财务管控、运营管控四大领域形成一体化经营视图。",
            "让数据实时共享、业务在线协同、管理过程可视、经营决策可跟踪。"
        ),
        closing="目标：通过数据驱动业务，支撑企业战略落地与集约化经营提升。",
    ),
    CardSpec(
        index="03",
        key="路",
        title="落地路径",
        english="Roadmap",
        accent=(53, 116, 136),
        strip=(241, 247, 249),
        summary="转型应分阶段推进，而不是一次性铺开所有能力。",
        bullets=(
            "第一阶段做核心业务数字化，打通销售、计划、供应链、生产到售后流程。",
            "第二阶段推进 IT/OT 融合数字工厂，提升车间协同、透明度与执行效率。",
            "第三阶段延伸到服务化与智能柔性制造，形成新的利润增长点与竞争壁垒。"
        ),
        closing="路径：先夯基础，再建工厂，再做服务化，最终走向智能柔性制造。",
    ),
)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def set_run_style(run, *, size: float, color: tuple[int, int, int], bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: float,
    color: tuple[int, int, int] = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=font_size, color=color, bold=bold)
    return shape


def add_shape(
    slide,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width: float = 1,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def keep_only_body_template(prs: Presentation) -> None:
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    keep_idx = manifest.slide_roles.body_template
    for index in range(len(prs.slides) - 1, -1, -1):
        if index != keep_idx:
            remove_slide(prs, index)


def render_title(slide) -> None:
    title_shape = slide.shapes[0]
    tf = title_shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    for text, color, bold in TITLE_SEGMENTS:
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size=24, color=color, bold=bold)

    add_text(slide, 720000, 545000, 8600000, 220000, SUBTITLE, font_size=11.8, color=MUTED)


def render_judgement(slide, slide_width: int) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 1140000, 70000, 300000, fill=ACCENT)
    add_text(slide, 860000, 1130000, 1900000, 220000, "核心判断", font_size=13.2, color=ACCENT, bold=True)
    add_text(slide, 720000, 1425000, slide_width - 1440000, 320000, HEADLINE, font_size=17.2, color=INK, bold=True)
    add_text(slide, 720000, 1815000, slide_width - 1440000, 180000, SUPPORT, font_size=10.8, color=MUTED)


def card_geometries(slide_width: int) -> list[tuple[int, int, int, int]]:
    left = 720000
    gap = 240000
    usable_width = slide_width - left * 2
    width = int((usable_width - gap * 2) / 3)
    top = 2200000
    height = 2600000
    return [(left + idx * (width + gap), top, width, height) for idx in range(3)]


def render_card(slide, geom: tuple[int, int, int, int], spec: CardSpec) -> None:
    left, top, width, height = geom
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=WHITE, line=LINE, line_width=1)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, 65000, fill=spec.accent)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + 180000, top + 170000, 360000, 360000, fill=spec.accent)
    add_text(slide, left + 180000, top + 240000, 360000, 120000, spec.key, font_size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, left + 640000, top + 150000, 1350000, 170000, spec.title, font_size=17.2, color=INK, bold=True)
    add_text(slide, left + 640000, top + 360000, 1200000, 110000, spec.english, font_size=8.5, color=spec.accent, bold=True)
    add_text(slide, left + width - 450000, top + 185000, 280000, 120000, spec.index, font_size=10.5, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.RIGHT)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 180000, top + 640000, width - 360000, 320000, fill=spec.strip)
    add_text(slide, left + 235000, top + 705000, width - 470000, 190000, spec.summary, font_size=12.1, color=spec.accent, bold=True)

    bullet_y = top + 1070000
    for bullet in spec.bullets:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + 185000, bullet_y + 65000, 70000, 70000, fill=spec.accent)
        add_text(slide, left + 315000, bullet_y, width - 470000, 265000, bullet, font_size=10.0, color=MUTED)
        bullet_y += 450000

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + 180000, top + height - 420000, width - 360000, 1, fill=LINE)
    add_text(slide, left + 180000, top + height - 325000, width - 360000, 180000, spec.closing, font_size=9.8, color=INK, bold=True)


def render_footer(slide, slide_width: int) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 5030000, slide_width - 1440000, 1, fill=LINE)
    add_text(slide, 720000, 5150000, slide_width - 1440000, 180000, FOOTNOTE, font_size=10.2, color=MUTED)


def build_slide(output_path: Path) -> Path:
    prs = Presentation(str(DEFAULT_TEMPLATE))
    keep_only_body_template(prs)
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)

    render_title(slide)
    render_judgement(slide, slide_width)
    for geom, spec in zip(card_geometries(slide_width), CARDS):
        render_card(slide, geom, spec)
    render_footer(slide, slide_width)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        timestamp = datetime.datetime.now().strftime("%H%M%S_%f")[:9]
        fallback = output_path.with_stem(f"{output_path.stem}_{timestamp}")
        prs.save(str(fallback))
        return fallback


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a one-page SIE slide for equipment manufacturing digital transformation.")
    parser.add_argument("--output", default=str(OUTPUT_PPT), help="Output PPTX path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = build_slide(Path(args.output).resolve())
    print(output)


if __name__ == "__main__":
    main()
