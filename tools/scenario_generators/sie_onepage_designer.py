from __future__ import annotations

import json
import math
from dataclasses import asdict, dataclass, replace
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

try:
    from sie_autoppt.config import COLOR_ACTIVE, DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.llm_openai import (
        OpenAIConfigurationError,
        OpenAIResponsesClient,
        OpenAIResponsesError,
        load_openai_responses_config,
    )
    from sie_autoppt.models import StructureArgument, StructureSection, StructureSpec
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import TemplateManifest, load_template_manifest
    from sie_autoppt.text_ops import write_text
except ModuleNotFoundError:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
    from sie_autoppt.config import COLOR_ACTIVE, DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.llm_openai import (
        OpenAIConfigurationError,
        OpenAIResponsesClient,
        OpenAIResponsesError,
        load_openai_responses_config,
    )
    from sie_autoppt.models import StructureArgument, StructureSection, StructureSpec
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import TemplateManifest, load_template_manifest
    from sie_autoppt.text_ops import write_text

try:
    from .review_onepage_slide import ReviewReport, review_onepage_slide, write_review_report
except ImportError:
    from review_onepage_slide import ReviewReport, review_onepage_slide, write_review_report


ACCENT = COLOR_ACTIVE
INK = (15, 23, 42)
TEXT_DARK = (30, 41, 59)
TEXT_MED = (71, 85, 105)
TEXT_LIGHT = (100, 116, 139)
FOOTER_TEXT = (148, 163, 184)
LINE = (203, 213, 225)
LIGHT_BG = (248, 250, 252)
WHITE = (255, 255, 255)
STRATEGY_BG = (30, 41, 59)
STRATEGY_TEXT = (148, 163, 184)
WATERMARK = (235, 239, 243)
BADGE_SOFT = (241, 245, 249)
BADGE_RED_SOFT = (255, 241, 242)
BADGE_RED_TEXT = (159, 18, 57)
BADGE_RED_LINE = (254, 205, 211)
COPY_ALLOWED_KEYWORDS = ("保持一致", "复刻", "一比一", "完全一致", "照着来", "same style", "replicate")
SUPPORTED_VARIANTS = (
    "asymmetric_focus",
    "balanced_dual_panel",
    "signal_band",
    "summary_board",
    "comparison_split",
    "timeline_vertical",
)
AUTO_LAYOUT_STRATEGY = "auto"


@dataclass(frozen=True)
class OnePageLayoutStrategy:
    strategy_id: str
    label: str
    business_use_case: str
    layout_variant: str
    cues: tuple[str, ...]


@dataclass(frozen=True)
class StrategySelectionResult:
    strategy_id: str
    layout_variant: str
    rationale: str
    source: str


ONEPAGE_STRATEGIES: tuple[OnePageLayoutStrategy, ...] = (
    OnePageLayoutStrategy(
        strategy_id="executive_summary_board",
        label="Executive Summary Board",
        business_use_case="单页结论汇报、老板决策页、管理层快扫页",
        layout_variant="summary_board",
        cues=("结论先行", "一句话结论", "三段支撑", "决策建议"),
    ),
    OnePageLayoutStrategy(
        strategy_id="status_dashboard",
        label="Status Dashboard",
        business_use_case="周报、月报、责任分工、状态跟踪、风险与动作同步",
        layout_variant="summary_board",
        cues=("状态", "责任人", "时限", "动作", "跟进"),
    ),
    OnePageLayoutStrategy(
        strategy_id="comparison_decision",
        label="Comparison Decision",
        business_use_case="方案对比、选型页、竞品比较、取舍判断",
        layout_variant="comparison_split",
        cues=("对比", "差异", "选项", "优劣", "判断"),
    ),
    OnePageLayoutStrategy(
        strategy_id="evidence_dense_brief",
        label="Evidence Dense Brief",
        business_use_case="高信息密度分析页、研究论证页、证据链说明页",
        layout_variant="balanced_dual_panel",
        cues=("分析", "论证", "依据", "证据", "洞察"),
    ),
    OnePageLayoutStrategy(
        strategy_id="strategy_blueprint",
        label="Strategy Blueprint",
        business_use_case="战略页、阶段规划页、路径蓝图页",
        layout_variant="balanced_dual_panel",
        cues=("战略", "蓝图", "路径", "规划", "阶段"),
    ),
    OnePageLayoutStrategy(
        strategy_id="process_storyline",
        label="Process Storyline",
        business_use_case="流程页、机制页、操作闭环页、执行路径页",
        layout_variant="signal_band",
        cues=("流程", "步骤", "提交", "审批", "闭环"),
    ),
    OnePageLayoutStrategy(
        strategy_id="roadmap_milestones",
        label="Roadmap Milestones",
        business_use_case="路线图、里程碑、分阶段推进页",
        layout_variant="timeline_vertical",
        cues=("路线图", "里程碑", "季度", "年度", "节点"),
    ),
    OnePageLayoutStrategy(
        strategy_id="solution_story",
        label="Solution Story",
        business_use_case="问题-方案-价值页、方案说明页、能力价值页",
        layout_variant="asymmetric_focus",
        cues=("问题", "方案", "价值", "抓手", "收益"),
    ),
)
STRATEGY_BY_ID = {strategy.strategy_id: strategy for strategy in ONEPAGE_STRATEGIES}


@dataclass(frozen=True)
class TextFragment:
    text: str
    bold: bool = False
    color: tuple[int, int, int] | None = None
    new_paragraph: bool = False


@dataclass(frozen=True)
class LawRow:
    number: str
    title: str
    badge: str
    badge_red: bool
    runs: tuple[TextFragment, ...]


@dataclass(frozen=True)
class BulletItem:
    label: str
    body: str


@dataclass(frozen=True)
class OnePageBrief:
    title: str
    kicker: str
    summary_fragments: tuple[TextFragment, ...]
    law_rows: tuple[LawRow, ...]
    right_kicker: str
    right_title: str
    process_steps: tuple[str, ...]
    right_bullets: tuple[BulletItem, ...]
    strategy_title: str
    strategy_fragments: tuple[TextFragment, ...]
    footer: str
    page_no: str
    required_terms: tuple[str, ...]
    variant: str = "asymmetric_focus"
    layout_strategy: str = ""
    reference_request: str = ""
    banned_phrases: tuple[str, ...] = ()
    layout_overrides: dict[str, float] | None = None
    typography_overrides: dict[str, float] | None = None


@dataclass(frozen=True)
class ScoreResult:
    total: int
    level: str
    template_fidelity: int
    title_fidelity_to_sie: int
    layout_originality: int
    self_check: int
    heuristic_review: int
    content_coverage: int
    reference_policy: str
    selected_strategy_id: str
    selected_variant: str
    strategy_selection_source: str
    strategy_selection_rationale: str
    issues: tuple[str, ...]
    review_findings: tuple[dict[str, str], ...]

    def to_dict(self) -> dict[str, object]:
        payload = asdict(self)
        payload["issues"] = list(self.issues)
        payload["review_findings"] = list(self.review_findings)
        return payload


def list_onepage_strategies() -> tuple[OnePageLayoutStrategy, ...]:
    return ONEPAGE_STRATEGIES


def _normalize_text(text: str) -> str:
    return "".join(str(text).lower().split())


def _brief_text_blob(brief: OnePageBrief) -> str:
    parts = [
        brief.title,
        brief.kicker,
        brief.right_kicker,
        brief.right_title,
        brief.strategy_title,
        " ".join(fragment.text for fragment in brief.summary_fragments),
        " ".join(row.title for row in brief.law_rows),
        " ".join(fragment.text for row in brief.law_rows for fragment in row.runs),
        " ".join(brief.process_steps),
        " ".join(item.label + item.body for item in brief.right_bullets),
        " ".join(fragment.text for fragment in brief.strategy_fragments),
    ]
    return "\n".join(part for part in parts if part).strip()


def _strategy_prompt_catalog() -> str:
    return "\n".join(
        f"- {strategy.strategy_id}: {strategy.label}; use for {strategy.business_use_case}; "
        f"default variant={strategy.layout_variant}; cues={', '.join(strategy.cues)}"
        for strategy in ONEPAGE_STRATEGIES
    )


def _build_strategy_selection_schema() -> dict[str, object]:
    return {
        "type": "object",
        "properties": {
            "strategy_id": {"type": "string", "enum": [strategy.strategy_id for strategy in ONEPAGE_STRATEGIES]},
            "rationale": {"type": "string", "minLength": 12, "maxLength": 180},
        },
        "required": ["strategy_id", "rationale"],
        "additionalProperties": False,
    }


def select_onepage_strategy_heuristically(brief: OnePageBrief) -> StrategySelectionResult:
    text_blob = _normalize_text(_brief_text_blob(brief))
    process_steps = tuple(step for step in brief.process_steps if step.strip())

    def score(*keywords: str) -> int:
        return sum(1 for keyword in keywords if _normalize_text(keyword) in text_blob)

    comparison_score = score("对比", "比较", "差异", "选型", "方案a", "方案b", "优劣", "竞品", "vs", "取舍")
    roadmap_score = score("路线图", "里程碑", "阶段", "季度", "年度", "节点", "排期", "规划")
    status_score = score("状态", "责任人", "时限", "周报", "月报", "风险", "问题", "跟进", "完成率", "上传")
    solution_score = score("问题", "方案", "价值", "抓手", "收益", "落地")
    strategy_score = score("战略", "蓝图", "路径", "规划", "举措", "目标")
    analysis_score = score("分析", "论证", "依据", "证据", "洞察", "判断")

    if comparison_score >= 2:
        strategy = STRATEGY_BY_ID["comparison_decision"]
        rationale = "内容包含明显对比/选型信号，优先用对称对比型单页承载差异与判断。"
    elif roadmap_score >= 2:
        strategy = STRATEGY_BY_ID["roadmap_milestones"]
        rationale = "内容带有阶段、里程碑或规划节奏，优先用路线图/里程碑型单页。"
    elif status_score >= 2:
        strategy = STRATEGY_BY_ID["status_dashboard"]
        rationale = "内容以责任人、时限和跟进为核心，更适合状态看板型单页。"
    elif len(process_steps) >= 4:
        strategy = STRATEGY_BY_ID["process_storyline"]
        rationale = "内容自带完整步骤链，优先用流程叙事型单页保持阅读顺序。"
    elif strategy_score >= 2:
        strategy = STRATEGY_BY_ID["strategy_blueprint"]
        rationale = "内容偏战略规划和路径设计，更适合蓝图型单页。"
    elif solution_score >= 2:
        strategy = STRATEGY_BY_ID["solution_story"]
        rationale = "内容结构接近问题-方案-价值，更适合方案叙事型单页。"
    elif analysis_score >= 2 or len(brief.law_rows) >= 4:
        strategy = STRATEGY_BY_ID["evidence_dense_brief"]
        rationale = "内容信息密度较高且有分析/论证属性，更适合高密度分析型单页。"
    else:
        strategy = STRATEGY_BY_ID["executive_summary_board"]
        rationale = "内容更像单页管理汇报，默认采用结论看板型单页。"

    return StrategySelectionResult(
        strategy_id=strategy.strategy_id,
        layout_variant=strategy.layout_variant,
        rationale=rationale,
        source="heuristic",
    )


def select_onepage_strategy_with_ai(brief: OnePageBrief, model: str | None = None) -> StrategySelectionResult:
    client = OpenAIResponsesClient(load_openai_responses_config(model=model))
    developer_prompt = (
        "You are selecting the best business one-page PPT strategy for a SIE-style slide.\n"
        "Choose the single best strategy based on content structure, not personal preference.\n"
        "Available strategies:\n"
        f"{_strategy_prompt_catalog()}\n"
        "Selection rules:\n"
        "- Prefer comparison_decision for explicit alternatives, trade-offs, or option selection.\n"
        "- Prefer roadmap_milestones for phases, milestones, or time-based planning.\n"
        "- Prefer process_storyline for 4+ ordered steps or operational handoffs.\n"
        "- Prefer status_dashboard for owners, deadlines, risks, or recurring status tracking.\n"
        "- Prefer evidence_dense_brief for analysis-heavy or proof-heavy content.\n"
        "- Prefer executive_summary_board when the slide is mainly a leadership summary.\n"
        "Return JSON only."
    )
    user_prompt = (
        "Select the best one-page strategy for this content.\n\n"
        f"Title: {brief.title}\n"
        f"Layout strategy hint: {brief.layout_strategy or 'auto'}\n"
        f"Right title: {brief.right_title}\n"
        f"Section count: {len(brief.law_rows)}\n"
        f"Process step count: {len(tuple(step for step in brief.process_steps if step.strip()))}\n"
        f"Content:\n{_brief_text_blob(brief)}"
    )
    payload = client.create_structured_json(
        developer_prompt=developer_prompt,
        user_prompt=user_prompt,
        schema_name="onepage_layout_strategy",
        schema=_build_strategy_selection_schema(),
    )
    strategy_id = str(payload.get("strategy_id", "")).strip()
    strategy = STRATEGY_BY_ID.get(strategy_id)
    if strategy is None:
        raise ValueError(f"Unknown one-page strategy selected by AI: {strategy_id}")
    rationale = str(payload.get("rationale", "")).strip() or "AI selected this strategy based on content structure."
    return StrategySelectionResult(
        strategy_id=strategy.strategy_id,
        layout_variant=strategy.layout_variant,
        rationale=rationale,
        source="ai",
    )


def resolve_onepage_strategy(
    brief: OnePageBrief,
    model: str | None = None,
    require_ai: bool = False,
) -> tuple[OnePageBrief, StrategySelectionResult]:
    layout_strategy = brief.layout_strategy.strip().lower()
    if layout_strategy and layout_strategy != AUTO_LAYOUT_STRATEGY:
        manual_strategy = STRATEGY_BY_ID.get(layout_strategy)
        if manual_strategy is not None:
            selection = StrategySelectionResult(
                strategy_id=manual_strategy.strategy_id,
                layout_variant=manual_strategy.layout_variant,
                rationale="Layout strategy was explicitly set by the caller.",
                source="manual",
            )
            return replace(brief, variant=manual_strategy.layout_variant), selection
        if layout_strategy in SUPPORTED_VARIANTS:
            return replace(brief, variant=layout_strategy), StrategySelectionResult(
                strategy_id="manual_variant",
                layout_variant=layout_strategy,
                rationale="Layout variant was explicitly set by the caller.",
                source="manual",
            )

    if brief.variant in SUPPORTED_VARIANTS and brief.layout_strategy.strip().lower() != AUTO_LAYOUT_STRATEGY:
        return brief, StrategySelectionResult(
            strategy_id="manual_variant",
            layout_variant=brief.variant,
            rationale="Layout variant was explicitly set by the caller.",
            source="manual",
        )

    if require_ai:
        try:
            selection = select_onepage_strategy_with_ai(brief, model=model)
        except (OpenAIConfigurationError, OpenAIResponsesError, ValueError):
            selection = select_onepage_strategy_heuristically(brief)
        return replace(brief, variant=selection.layout_variant), selection

    try:
        selection = select_onepage_strategy_with_ai(brief, model=model)
    except (OpenAIConfigurationError, OpenAIResponsesError, ValueError):
        selection = select_onepage_strategy_heuristically(brief)

    return replace(brief, variant=selection.layout_variant), selection


def _clip(text: str, limit: int) -> str:
    normalized = " ".join(str(text).strip().split())
    if len(normalized) <= limit:
        return normalized
    return normalized[: max(0, limit - 1)].rstrip() + "…"


def _format_argument(argument: StructureArgument) -> str:
    point = " ".join(argument.point.strip().split())
    evidence = " ".join(argument.evidence.strip().split())
    if point and evidence:
        return f"{point}：{evidence}"
    return point or evidence


def _section_to_law_row(section: StructureSection, index: int) -> LawRow:
    body_runs: list[TextFragment] = [TextFragment(_clip(section.key_message, 40), bold=True, color=TEXT_DARK)]
    for argument in section.arguments[:3]:
        text = _format_argument(argument)
        if not text:
            continue
        body_runs.append(TextFragment(" " + _clip(text, 50)))
    badge = "重点事项" if index == 1 else "核心支撑"
    return LawRow(
        number=f"{index:02d}",
        title=_clip(section.title, 20),
        badge=badge,
        badge_red=index == 1,
        runs=tuple(body_runs),
    )


def build_onepage_brief_from_structure(
    structure: StructureSpec,
    *,
    topic: str,
    footer: str = "STRICTLY CONFIDENTIAL | 2026 SIE One-page Brief",
    page_no: str = "01",
    layout_strategy: str = AUTO_LAYOUT_STRATEGY,
) -> OnePageBrief:
    sections = structure.sections[:3] or [
        StructureSection(
            title="核心结论",
            key_message=structure.core_message or topic,
            arguments=[StructureArgument(point="补充业务信息后可生成更完整单页", evidence="")],
        )
    ]
    law_rows = tuple(_section_to_law_row(section, index + 1) for index, section in enumerate(sections))
    summary_intro = _clip(structure.core_message or topic, 72)
    strongest_section = sections[0]
    right_bullets: list[BulletItem] = []
    for section in sections:
        bullet_body = _clip(
            "；".join(
                _format_argument(argument)
                for argument in section.arguments[:2]
                if _format_argument(argument)
            )
            or section.key_message,
            72,
        )
        right_bullets.append(BulletItem(f"{_clip(section.title, 10)}：", bullet_body))

    process_steps = tuple(_clip(section.title, 8) for section in sections)
    if len(process_steps) < 4:
        process_steps = process_steps + ("行动落地",)

    required_terms = tuple(
        term for term in [topic.strip(), structure.core_message.strip(), *(section.title.strip() for section in sections)] if term
    )
    return OnePageBrief(
        title=_clip(topic.strip() or structure.core_message or "SIE 单页汇报", 24),
        kicker="",
        summary_fragments=(
            TextFragment(summary_intro),
            TextFragment(
                _clip(f"当前页围绕“{strongest_section.title}”展开，适合用于商务汇报中的单页正文表达。", 72),
                bold=True,
                color=ACCENT,
                new_paragraph=True,
            ),
        ),
        law_rows=law_rows,
        right_kicker="EXECUTION VIEW",
        right_title=_clip(structure.core_message or strongest_section.key_message or topic, 24),
        process_steps=process_steps,
        right_bullets=tuple(right_bullets[:3]),
        strategy_title="行动建议：围绕核心结论组织表达与下一步动作",
        strategy_fragments=(
            TextFragment(_clip(f"建议优先突出“{strongest_section.title}”的业务判断，再用 2-3 个支撑点解释原因和动作。", 86)),
            TextFragment("需要更强管理层风格时，可进一步压缩成结论、动作、风险三块。", bold=True, color=WHITE, new_paragraph=True),
        ),
        footer=footer,
        page_no=page_no,
        required_terms=required_terms[:8] or ("SIE",),
        variant="auto",
        layout_strategy=layout_strategy,
    )


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def resolve_reference_policy(reference_request: str) -> str:
    normalized = reference_request.strip().lower()
    if normalized and any(keyword.lower() in normalized for keyword in COPY_ALLOWED_KEYWORDS):
        return "replicate_allowed"
    return "extract_style_only"


def set_run_style(run, *, size: float, color: tuple[int, int, int], bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_shape(
    slide,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width: float = 0.8,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    text: str = "",
    font_size: float = 12.0,
    color: tuple[int, int, int] = TEXT_MED,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    if text:
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size=font_size, color=color, bold=bold)
    return shape


def write_fragments(shape, fragments: tuple[TextFragment, ...], *, font_size: float, default_color: tuple[int, int, int], align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    first = True
    for fragment in fragments:
        if not first and fragment.new_paragraph:
            paragraph = tf.add_paragraph()
            paragraph.alignment = align
        run = paragraph.add_run()
        run.text = fragment.text
        set_run_style(run, size=font_size, color=fragment.color or default_color, bold=fragment.bold)
        first = False


def estimate_lines(text: str, width: int, font_size: float) -> int:
    weighted_chars = 0.0
    for char in text:
        if char.isspace():
            weighted_chars += 0.35
        elif ord(char) < 128:
            weighted_chars += 0.58
        else:
            weighted_chars += 1.0
    chars_per_line = max(1.0, width / max(font_size * 7000, 1))
    return max(1, math.ceil(weighted_chars / chars_per_line))


def choose_title_font_size(text: str) -> int:
    length = len(text.strip())
    if length > 28:
        return 18
    if length > 22:
        return 20
    if length > 16:
        return 22
    return 24


def keep_only_body_template(prs: Presentation) -> None:
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    keep_idx = manifest.slide_roles.body_template
    for index in range(len(prs.slides) - 1, -1, -1):
        if index != keep_idx:
            remove_slide(prs, index)


def _pick_text_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def _clear_render_area(slide, manifest: TemplateManifest, protected_shapes: list[object]) -> None:
    removable = []
    for shape in slide.shapes:
        if any(shape is protected for protected in protected_shapes):
            continue
        if manifest.selectors.body_render_area.matches(shape):
            removable.append(shape)
    for shape in removable:
        element = shape._element
        element.getparent().remove(element)


def _find_title_shape(slide, manifest: TemplateManifest):
    texts = sorted(_pick_text_shapes(slide), key=lambda shape: (shape.top, shape.left))
    title_candidates = [shape for shape in texts if manifest.selectors.body_title.matches(shape)]
    return title_candidates[0] if title_candidates else None


def _find_subtitle_shape(slide, manifest: TemplateManifest):
    texts = sorted(_pick_text_shapes(slide), key=lambda shape: (shape.top, shape.left))
    title_candidates = [shape for shape in texts if manifest.selectors.body_title.matches(shape)]
    subtitle_candidates = [
        shape
        for shape in texts
        if manifest.selectors.body_subtitle.matches(shape) and shape not in title_candidates
    ]
    return subtitle_candidates[0] if subtitle_candidates else None


def apply_sie_header(slide, manifest: TemplateManifest, brief: OnePageBrief) -> list[object]:
    title_shape = _find_title_shape(slide, manifest)
    subtitle_shape = _find_subtitle_shape(slide, manifest)
    if title_shape is None:
        fallback = manifest.fallback_boxes.body_title
        title_shape = add_textbox(
            slide,
            fallback.left,
            fallback.top,
            fallback.width,
            fallback.height,
            text="",
        )
    write_text(
        title_shape,
        brief.title,
        color=ACCENT,
        size_pt=choose_title_font_size(brief.title),
        bold=True,
        preserve_runs=True,
    )

    kicker_box = add_textbox(slide, 650000, 250000, 3500000, 170000, text=brief.kicker, font_size=9.6, color=TEXT_LIGHT, bold=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 3200000, 332000, 420000, 12000, fill=LINE)
    if subtitle_shape is not None:
        subtitle_shape.text_frame.word_wrap = True
        subtitle_shape.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        write_fragments(subtitle_shape, brief.summary_fragments, font_size=12.2, default_color=TEXT_MED)
    return [title_shape, subtitle_shape, kicker_box]


def render_law_rows(slide, rows: tuple[LawRow, ...], *, left: int, top: int, title_width: int, desc_width: int, badge_left: int, row_gap: int, font_size: float) -> None:
    current_y = top
    for item in rows:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, current_y + 120000, 30000, 760000, fill=LINE)
        add_textbox(slide, left - 220000, current_y - 160000, 520000, 520000, text=item.number, font_size=42, color=WATERMARK, bold=True)
        add_textbox(slide, left + 140000, current_y, title_width, 220000, text=item.title, font_size=14.0, color=TEXT_DARK, bold=True)

        badge_fill = BADGE_RED_SOFT if item.badge_red else BADGE_SOFT
        badge_line = BADGE_RED_LINE if item.badge_red else LINE
        badge_text = BADGE_RED_TEXT if item.badge_red else TEXT_MED
        badge = add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            badge_left,
            current_y + 8000,
            1020000,
            220000,
            fill=badge_fill,
            line=badge_line,
        )
        badge.text_frame.clear()
        paragraph = badge.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = item.badge
        set_run_style(run, size=9.0, color=badge_text, bold=True)

        desc = add_textbox(slide, left + 140000, current_y + 270000, desc_width, 780000)
        write_fragments(desc, item.runs, font_size=font_size, default_color=TEXT_MED)
        current_y += row_gap


def render_process_flow(slide, steps: tuple[str, ...], *, left: int, top: int, width: int, active_last: bool = True) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, 420000, fill=LIGHT_BG, line=LINE, line_width=0.9)
    step_w = max(520000, int((width - 200000 - (len(steps) - 1) * 105000) / len(steps)))
    gap_w = 105000
    current_x = left + 100000
    for index, step in enumerate(steps):
        is_last = active_last and index == len(steps) - 1
        fill = (255, 250, 251) if is_last else WHITE
        line = ACCENT if is_last else LINE
        text_color = ACCENT if is_last else TEXT_DARK
        step_box = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, current_x, top + 75000, step_w, 250000, fill=fill, line=line)
        step_box.text_frame.clear()
        paragraph = step_box.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = step
        set_run_style(run, size=8.8, color=text_color, bold=True)
        if index != len(steps) - 1:
            add_textbox(slide, current_x + step_w + 10000, top + 82000, gap_w - 20000, 240000, text=">", font_size=11.5, color=TEXT_LIGHT, bold=True, align=PP_ALIGN.CENTER)
        current_x += step_w + gap_w


def render_bullets(
    slide,
    bullets: tuple[BulletItem, ...],
    *,
    left: int,
    top: int,
    width: int,
    row_gap: int = 490000,
    marker_color: tuple[int, int, int] = ACCENT,
    font_size: float = 10.6,
) -> None:
    for index, item in enumerate(bullets):
        bullet_y = top + index * row_gap
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, bullet_y + 26000, 56000, 56000, fill=marker_color)
        box = add_textbox(slide, left + 120000, bullet_y, width - 120000, 340000)
        paragraph = box.text_frame.paragraphs[0]
        head = paragraph.add_run()
        head.text = item.label
        set_run_style(head, size=font_size, color=INK, bold=True)
        tail = paragraph.add_run()
        tail.text = item.body
        set_run_style(tail, size=font_size, color=TEXT_MED, bold=False)


def render_strategy_box(slide, title: str, fragments: tuple[TextFragment, ...], *, left: int, top: int, width: int, height: int) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=STRATEGY_BG)
    title_box = add_textbox(slide, left + 140000, top + 90000, width - 280000, 160000)
    title_box.text_frame.clear()
    paragraph = title_box.text_frame.paragraphs[0]
    star = paragraph.add_run()
    star.text = "* "
    set_run_style(star, size=10.0, color=ACCENT, bold=True)
    title_run = paragraph.add_run()
    title_run.text = title
    set_run_style(title_run, size=10.8, color=WHITE, bold=True)

    body_box = add_textbox(slide, left + 140000, top + 260000, width - 280000, height - 330000)
    write_fragments(body_box, fragments, font_size=9.5, default_color=STRATEGY_TEXT)


def render_footer(slide, slide_width: int, *, footer_text: str, page_no: str) -> None:
    footer_top = 6400000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 650000, footer_top, slide_width - 1300000, 12000, fill=LINE)
    add_textbox(slide, 650000, footer_top + 90000, 8600000, 180000, text=footer_text, font_size=8.8, color=FOOTER_TEXT, bold=True)
    add_textbox(slide, slide_width - 1100000, footer_top + 90000, 420000, 180000, text=page_no, font_size=9.4, color=FOOTER_TEXT, bold=True, align=PP_ALIGN.RIGHT)


def render_asymmetric_focus(slide, brief: OnePageBrief) -> None:
    render_law_rows(
        slide,
        brief.law_rows,
        left=760000,
        top=1980000,
        title_width=3500000,
        desc_width=4780000,
        badge_left=4460000,
        row_gap=1260000,
        font_size=10.6,
    )
    panel_left = 6500000
    panel_top = 1940000
    panel_width = 5050000
    panel_height = 4030000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, panel_left, panel_top, panel_width, panel_height, fill=WHITE, line=LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, panel_left, panel_top, panel_width, 72000, fill=ACCENT)
    add_textbox(slide, panel_left + 260000, panel_top + 220000, panel_width - 520000, 170000, text=brief.right_kicker, font_size=9.8, color=ACCENT, bold=True)
    add_textbox(slide, panel_left + 260000, panel_top + 420000, panel_width - 520000, 250000, text=brief.right_title, font_size=16.8, color=INK, bold=True)
    render_process_flow(slide, brief.process_steps, left=panel_left + 260000, top=panel_top + 860000, width=panel_width - 520000)
    render_bullets(slide, brief.right_bullets, left=panel_left + 280000, top=panel_top + 1420000, width=panel_width - 700000)
    render_strategy_box(slide, brief.strategy_title, brief.strategy_fragments, left=panel_left + 260000, top=panel_top + 3080000, width=panel_width - 520000, height=760000)


def render_balanced_dual_panel(slide, brief: OnePageBrief) -> None:
    panel_top = 1980000
    panel_width = 5150000
    panel_height = 2680000
    left_panel = 720000
    right_panel = 6170000
    for left, title in ((left_panel, "Regulation Stack"), (right_panel, brief.right_kicker)):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, panel_top, panel_width, panel_height, fill=WHITE, line=LINE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, panel_top, panel_width, 72000, fill=ACCENT)
        add_textbox(slide, left + 220000, panel_top + 190000, panel_width - 440000, 180000, text=title, font_size=9.4, color=ACCENT, bold=True)

    render_law_rows(
        slide,
        brief.law_rows,
        left=880000,
        top=panel_top + 520000,
        title_width=2820000,
        desc_width=3840000,
        badge_left=3500000,
        row_gap=700000,
        font_size=9.6,
    )
    add_textbox(slide, right_panel + 220000, panel_top + 390000, panel_width - 440000, 220000, text=brief.right_title, font_size=15.4, color=INK, bold=True)
    render_process_flow(slide, brief.process_steps, left=right_panel + 220000, top=panel_top + 760000, width=panel_width - 440000)
    render_bullets(slide, brief.right_bullets, left=right_panel + 260000, top=panel_top + 1320000, width=panel_width - 520000, row_gap=360000)
    render_strategy_box(slide, brief.strategy_title, brief.strategy_fragments, left=720000, top=4880000, width=10600000, height=980000)


def render_signal_band(slide, brief: OnePageBrief) -> None:
    band_left = 720000
    band_top = 1900000
    band_width = 10600000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, band_left, band_top, band_width, 520000, fill=LIGHT_BG, line=LINE)
    band_box = add_textbox(slide, band_left + 180000, band_top + 110000, band_width - 360000, 250000)
    write_fragments(band_box, brief.summary_fragments, font_size=12.0, default_color=TEXT_MED)

    card_top = 2600000
    card_width = 3200000
    gap = 220000
    for index, item in enumerate(brief.law_rows[:3]):
        left = 720000 + index * (card_width + gap)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, card_top, card_width, 1580000, fill=WHITE, line=LINE)
        add_textbox(slide, left + 140000, card_top + 120000, 520000, 220000, text=item.number, font_size=26, color=WATERMARK, bold=True)
        badge_fill = BADGE_RED_SOFT if item.badge_red else BADGE_SOFT
        badge_line = BADGE_RED_LINE if item.badge_red else LINE
        badge_text = BADGE_RED_TEXT if item.badge_red else TEXT_MED
        badge = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 1960000, card_top + 120000, 900000, 220000, fill=badge_fill, line=badge_line)
        badge.text_frame.clear()
        paragraph = badge.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = item.badge
        set_run_style(run, size=8.8, color=badge_text, bold=True)
        add_textbox(slide, left + 140000, card_top + 390000, 2780000, 210000, text=item.title, font_size=13.0, color=TEXT_DARK, bold=True)
        desc = add_textbox(slide, left + 140000, card_top + 660000, 2900000, 700000)
        write_fragments(desc, item.runs, font_size=9.4, default_color=TEXT_MED)

    process_box_left = 720000
    process_box_top = 4420000
    process_box_width = 5000000
    process_box_height = 1440000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, process_box_left, process_box_top, process_box_width, process_box_height, fill=WHITE, line=LINE)
    add_textbox(slide, process_box_left + 180000, process_box_top + 160000, process_box_width - 360000, 190000, text=brief.right_kicker, font_size=9.4, color=ACCENT, bold=True)
    add_textbox(slide, process_box_left + 180000, process_box_top + 420000, process_box_width - 360000, 220000, text=brief.right_title, font_size=15.0, color=INK, bold=True)
    render_process_flow(slide, brief.process_steps, left=process_box_left + 180000, top=process_box_top + 780000, width=process_box_width - 360000)

    action_box_left = 5920000
    action_box_width = 5380000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, action_box_left, process_box_top, action_box_width, process_box_height, fill=WHITE, line=LINE)
    render_bullets(slide, brief.right_bullets, left=action_box_left + 220000, top=process_box_top + 200000, width=action_box_width - 440000, row_gap=340000)
    render_strategy_box(slide, brief.strategy_title, brief.strategy_fragments, left=720000, top=5940000, width=10600000, height=420000)


def render_summary_board(slide, brief: OnePageBrief) -> None:
    layout = brief.layout_overrides or {}
    typo = brief.typography_overrides or {}

    y_offset = int(layout.get("summary_y_offset", -420000))
    hero_height = int(layout.get("hero_height", 560000))
    card_height = int(layout.get("card_height", 1520000))
    process_panel_height = int(layout.get("process_panel_height", 1320000))
    strategy_height = int(layout.get("strategy_height", 420000))
    right_title_font = float(typo.get("right_title_font_size", 17.0))
    panel_title_font = float(typo.get("panel_title_font_size", 15.2))
    bullet_font = float(typo.get("bullet_font_size", 10.6))

    hero_left = 720000
    hero_top = 1880000 + y_offset
    hero_width = 10600000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, hero_left, hero_top, hero_width, hero_height, fill=LIGHT_BG, line=LINE)
    add_textbox(slide, hero_left + 180000, hero_top + 90000, 3600000, 150000, text=brief.right_kicker or "EXECUTIVE VIEW", font_size=9.4, color=ACCENT, bold=True)
    add_textbox(slide, hero_left + 180000, hero_top + 220000, 4600000, 220000, text=brief.right_title, font_size=right_title_font, color=INK, bold=True)
    hero_summary = add_textbox(slide, hero_left + 5100000, hero_top + 120000, 5200000, 280000)
    write_fragments(hero_summary, brief.summary_fragments, font_size=10.8, default_color=TEXT_MED)

    card_top = 2620000 + y_offset
    card_width = 3200000
    gap = 220000
    for index, item in enumerate(brief.law_rows[:3]):
        left = 720000 + index * (card_width + gap)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, card_top, card_width, card_height, fill=WHITE, line=LINE)
        add_textbox(slide, left + 130000, card_top + 90000, 420000, 220000, text=item.number, font_size=24, color=WATERMARK, bold=True)
        badge_fill = BADGE_RED_SOFT if item.badge_red else BADGE_SOFT
        badge_line = BADGE_RED_LINE if item.badge_red else LINE
        badge_text = BADGE_RED_TEXT if item.badge_red else TEXT_MED
        badge = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 1800000, card_top + 100000, 980000, 210000, fill=badge_fill, line=badge_line)
        badge.text_frame.clear()
        paragraph = badge.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = item.badge
        set_run_style(run, size=8.8, color=badge_text, bold=True)
        add_textbox(slide, left + 130000, card_top + 380000, 2880000, 210000, text=item.title, font_size=12.4, color=TEXT_DARK, bold=True)
        desc = add_textbox(slide, left + 130000, card_top + 640000, 2880000, 650000)
        write_fragments(desc, item.runs, font_size=9.2, default_color=TEXT_MED)

    process_panel_left = 720000
    process_panel_top = 4380000 + y_offset
    process_panel_width = 5000000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, process_panel_left, process_panel_top, process_panel_width, process_panel_height, fill=WHITE, line=LINE)
    add_textbox(slide, process_panel_left + 180000, process_panel_top + 180000, process_panel_width - 360000, 170000, text="KEY FLOW", font_size=9.4, color=ACCENT, bold=True)
    add_textbox(slide, process_panel_left + 180000, process_panel_top + 390000, process_panel_width - 360000, 220000, text="关键流程与交接节奏", font_size=panel_title_font, color=INK, bold=True)
    render_process_flow(slide, brief.process_steps, left=process_panel_left + 180000, top=process_panel_top + 760000, width=process_panel_width - 360000)

    action_panel_left = 5940000
    action_panel_top = 4380000 + y_offset
    action_panel_width = 5380000
    action_panel_height = process_panel_height
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, action_panel_left, action_panel_top, action_panel_width, action_panel_height, fill=WHITE, line=LINE)
    add_textbox(slide, action_panel_left + 220000, action_panel_top + 180000, action_panel_width - 440000, 170000, text="ACTION POINTS", font_size=9.4, color=ACCENT, bold=True)
    add_textbox(slide, action_panel_left + 220000, action_panel_top + 390000, action_panel_width - 440000, 180000, text="关键动作与跟进重点", font_size=panel_title_font, color=INK, bold=True)
    render_bullets(slide, brief.right_bullets, left=action_panel_left + 220000, top=action_panel_top + 700000, width=action_panel_width - 440000, row_gap=250000, font_size=bullet_font)

    render_strategy_box(
        slide,
        brief.strategy_title,
        brief.strategy_fragments,
        left=720000,
        top=5860000 + y_offset,
        width=10600000,
        height=strategy_height,
    )


def render_comparison_split(slide, brief: OnePageBrief) -> None:
    panel_top = 1940000
    panel_width = 5000000
    panel_height = 3060000
    left_panel = 720000
    right_panel = 6200000

    left_rows = brief.law_rows[: max(1, math.ceil(len(brief.law_rows) / 2))]
    right_rows = brief.law_rows[max(1, math.ceil(len(brief.law_rows) / 2)) :] or brief.law_rows[-1:]

    for left, title in ((left_panel, "LEFT VIEW"), (right_panel, "RIGHT VIEW")):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, panel_top, panel_width, panel_height, fill=WHITE, line=LINE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, panel_top, panel_width, 72000, fill=ACCENT)
        add_textbox(slide, left + 180000, panel_top + 180000, panel_width - 360000, 160000, text=title, font_size=9.2, color=ACCENT, bold=True)

    render_law_rows(
        slide,
        left_rows,
        left=860000,
        top=panel_top + 430000,
        title_width=2500000,
        desc_width=3200000,
        badge_left=3050000,
        row_gap=1150000,
        font_size=9.6,
    )
    render_law_rows(
        slide,
        right_rows,
        left=6340000,
        top=panel_top + 430000,
        title_width=2500000,
        desc_width=3200000,
        badge_left=8530000,
        row_gap=1150000,
        font_size=9.6,
    )
    add_textbox(slide, right_panel + 180000, panel_top + 2420000, panel_width - 360000, 170000, text=brief.right_kicker, font_size=9.0, color=ACCENT, bold=True)
    render_bullets(slide, brief.right_bullets, left=right_panel + 180000, top=panel_top + 2640000, width=panel_width - 360000, row_gap=250000)

    flow_top = 5160000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 720000, flow_top, 10600000, 520000, fill=LIGHT_BG, line=LINE)
    render_process_flow(slide, brief.process_steps, left=900000, top=flow_top + 50000, width=10240000)
    render_strategy_box(slide, brief.strategy_title, brief.strategy_fragments, left=720000, top=5780000, width=10600000, height=420000)


def render_vertical_steps(slide, steps: tuple[str, ...], *, left: int, top: int, width: int, height: int) -> None:
    count = max(1, len(steps))
    usable_height = height - 160000
    step_h = min(420000, max(260000, int((usable_height - (count - 1) * 100000) / count)))
    current_y = top + 80000
    for index, step in enumerate(steps):
        is_last = index == count - 1
        fill = (255, 250, 251) if is_last else WHITE
        line = ACCENT if is_last else LINE
        text_color = ACCENT if is_last else TEXT_DARK
        step_box = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 140000, current_y, width - 280000, step_h, fill=fill, line=line)
        step_box.text_frame.clear()
        paragraph = step_box.text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = step
        set_run_style(run, size=10.0, color=text_color, bold=True)
        if index != count - 1:
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + width // 2 - 10000, current_y + step_h, 20000, 100000, fill=LINE)
        current_y += step_h + 100000


def render_timeline_vertical(slide, brief: OnePageBrief) -> None:
    timeline_left = 720000
    timeline_top = 1940000
    timeline_width = 2420000
    timeline_height = 3820000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, timeline_left, timeline_top, timeline_width, timeline_height, fill=WHITE, line=LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, timeline_left, timeline_top, timeline_width, 72000, fill=ACCENT)
    add_textbox(slide, timeline_left + 170000, timeline_top + 180000, timeline_width - 340000, 160000, text=brief.right_kicker or "ROADMAP", font_size=9.2, color=ACCENT, bold=True)
    add_textbox(slide, timeline_left + 170000, timeline_top + 390000, timeline_width - 340000, 220000, text=brief.right_title, font_size=14.6, color=INK, bold=True)
    render_vertical_steps(slide, brief.process_steps, left=timeline_left, top=timeline_top + 760000, width=timeline_width, height=2800000)

    summary_left = 3400000
    summary_top = 1940000
    summary_width = 7900000
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, summary_left, summary_top, summary_width, 740000, fill=LIGHT_BG, line=LINE)
    summary_box = add_textbox(slide, summary_left + 200000, summary_top + 120000, summary_width - 400000, 360000)
    write_fragments(summary_box, brief.summary_fragments, font_size=11.2, default_color=TEXT_MED)

    render_law_rows(
        slide,
        brief.law_rows,
        left=3560000,
        top=2920000,
        title_width=2500000,
        desc_width=3700000,
        badge_left=6800000,
        row_gap=930000,
        font_size=9.5,
    )
    render_bullets(slide, brief.right_bullets, left=8300000, top=2880000, width=2800000, row_gap=390000)
    render_strategy_box(slide, brief.strategy_title, brief.strategy_fragments, left=720000, top=5900000, width=10600000, height=420000)


def render_variant(slide, brief: OnePageBrief) -> None:
    if brief.variant == "summary_board":
        render_summary_board(slide, brief)
        return
    if brief.variant == "comparison_split":
        render_comparison_split(slide, brief)
        return
    if brief.variant == "timeline_vertical":
        render_timeline_vertical(slide, brief)
        return
    if brief.variant == "balanced_dual_panel":
        render_balanced_dual_panel(slide, brief)
        return
    if brief.variant == "signal_band":
        render_signal_band(slide, brief)
        return
    render_asymmetric_focus(slide, brief)


def collect_slide_text(prs: Presentation) -> str:
    slide = prs.slides[0]
    return "\n".join(shape.text_frame.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)).strip()


def _title_shape_is_sie(prs: Presentation, brief: OnePageBrief) -> bool:
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    slide = prs.slides[0]
    title_shape = _find_title_shape(slide, manifest)
    if title_shape is None:
        return False
    if title_shape.text_frame.text.strip() != brief.title.strip():
        return False
    for paragraph in title_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                if run.font.color is None or run.font.color.rgb is None:
                    return False
                if tuple(int(value) for value in run.font.color.rgb) != ACCENT:
                    return False
                return True
    return False


def self_check_layout(prs: Presentation, brief: OnePageBrief) -> list[str]:
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    issues: list[str] = []
    if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
        issues.append("slide contains picture shapes; expected editable-only output")
    for index, shape in enumerate(slide.shapes, start=1):
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_width or shape.top + shape.height > slide_height:
            issues.append(f"shape {index} exceeds slide bounds")
    if estimate_lines(brief.title, 7800000, choose_title_font_size(brief.title)) > 2:
        issues.append("title is likely to wrap to more than two lines")
    summary_text = "".join(fragment.text for fragment in brief.summary_fragments)
    if estimate_lines(summary_text, 10300000, 12.2) > 3:
        issues.append("summary is too dense")
    for row in brief.law_rows:
        body_text = "".join(fragment.text for fragment in row.runs)
        if estimate_lines(body_text, 4200000, 10.0) > 5:
            issues.append(f"{row.title} description is too dense")
    for bullet in brief.right_bullets:
        if estimate_lines(bullet.label + bullet.body, 4400000, 10.6) > 4:
            issues.append(f"{bullet.label} bullet is too dense")
    strategy_text = "".join(fragment.text for fragment in brief.strategy_fragments)
    if estimate_lines(strategy_text, 9000000, 9.5) > 4:
        issues.append("strategy box is too dense")
    if brief.variant not in SUPPORTED_VARIANTS:
        issues.append(f"unsupported one-page variant: {brief.variant}")
    return issues


def score_slide(
    prs: Presentation,
    brief: OnePageBrief,
    issues: list[str],
    review_report: ReviewReport,
    selection: StrategySelectionResult,
) -> ScoreResult:
    slide_text = collect_slide_text(prs)
    reference_policy = resolve_reference_policy(brief.reference_request)
    template_fidelity = 18
    title_fidelity_to_sie = 15 if _title_shape_is_sie(prs, brief) else 6
    if reference_policy == "replicate_allowed":
        layout_originality = 6
    elif brief.variant == "asymmetric_focus":
        layout_originality = 8
    else:
        layout_originality = 10
    self_check_score = max(0, 20 - len(issues) * 5)
    blockers = [finding for finding in review_report.findings if finding.level == "blocker"]
    review_warnings = [finding for finding in review_report.findings if finding.level == "warning"]
    heuristic_score = 20 - len(blockers) * 12
    if review_report.preview_note == "preview export skipped":
        heuristic_score -= 2
    heuristic_score -= len(review_warnings) * 3
    heuristic_score = max(0, heuristic_score)
    hits = sum(1 for term in brief.required_terms if term in slide_text)
    content_coverage = round(17 * hits / max(1, len(brief.required_terms)))
    total = max(
        0,
        min(
            100,
            template_fidelity
            + title_fidelity_to_sie
            + layout_originality
            + self_check_score
            + heuristic_score
            + content_coverage,
        ),
    )
    if total >= 95:
        level = "优秀"
    elif total >= 88:
        level = "良好"
    elif total >= 80:
        level = "合格"
    else:
        level = "需继续修正"
    return ScoreResult(
        total=total,
        level=level,
        template_fidelity=template_fidelity,
        title_fidelity_to_sie=title_fidelity_to_sie,
        layout_originality=layout_originality,
        self_check=self_check_score,
        heuristic_review=heuristic_score,
        content_coverage=content_coverage,
        reference_policy=reference_policy,
        selected_strategy_id=selection.strategy_id,
        selected_variant=selection.layout_variant,
        strategy_selection_source=selection.source,
        strategy_selection_rationale=selection.rationale,
        issues=tuple(issues),
        review_findings=tuple(finding.to_dict() for finding in review_report.findings),
    )


def build_onepage_slide(
    brief: OnePageBrief,
    *,
    output_path: Path,
    export_review: bool = True,
    template_path: Path = DEFAULT_TEMPLATE,
    model: str | None = None,
    require_ai_strategy: bool = False,
) -> tuple[Path, Path | None, Path, ScoreResult]:
    resolved_brief, selection = resolve_onepage_strategy(brief, model=model, require_ai=require_ai_strategy)
    prs = Presentation(str(template_path))
    keep_only_body_template(prs)
    manifest = load_template_manifest(template_path=template_path)
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)

    protected = apply_sie_header(slide, manifest, resolved_brief)
    _clear_render_area(slide, manifest, [shape for shape in protected if shape is not None])
    render_variant(slide, resolved_brief)
    render_footer(slide, slide_width, footer_text=resolved_brief.footer, page_no=resolved_brief.page_no)

    issues = self_check_layout(prs, resolved_brief)
    if issues:
        raise ValueError("layout self-check failed: " + " | ".join(issues))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    if export_review:
        review_report = review_onepage_slide(
            output_path,
            expected_card_count=0,
            banned_phrases=resolved_brief.banned_phrases,
            export_previews=False,
        )
        review_path = output_path.with_suffix(".review.json")
        write_review_report(review_report, review_path)
    else:
        review_report = ReviewReport(passed=True, findings=(), preview_paths=(), preview_note="review skipped")
        review_path = None

    scored_prs = Presentation(str(output_path))
    score = score_slide(scored_prs, resolved_brief, issues, review_report, selection)
    score_path = output_path.with_suffix(".score.json")
    score_path.write_text(json.dumps(score.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
    return output_path, review_path, score_path, score
