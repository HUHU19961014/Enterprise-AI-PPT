from __future__ import annotations

import argparse
import datetime
from pathlib import Path
from statistics import mean
import sys
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

try:
    from sie_autoppt.config import COLOR_ACTIVE, DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide
except ModuleNotFoundError:
    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
    from sie_autoppt.config import COLOR_ACTIVE, DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide


ACCENT = COLOR_ACTIVE
NAVY = (60, 76, 96)
SLATE = (91, 112, 135)
MUTED = (132, 147, 162)
WHITE = (255, 255, 255)
LIGHT_BG = (247, 249, 251)
LIGHT_BORDER = (221, 226, 232)
QUADRANT_FILL = (251, 242, 245)
LOW_FILL = (242, 246, 249)
SECTION_COLORS = [
    (243, 230, 235),
    (234, 240, 245),
    (245, 236, 240),
    (234, 239, 243),
]
HEAT_COLORS = {
    0: (243, 246, 249),
    1: (248, 232, 237),
    2: (242, 210, 220),
    3: (232, 177, 193),
    4: (214, 110, 142),
    5: ACCENT,
}

VENDORS = [
    ("一道新能", "一道"),
    ("正泰新能", "正泰"),
    ("晶科能源", "晶科"),
    ("晶澳科技", "晶澳"),
    ("横店东磁", "横店"),
    ("无锡博达", "博达"),
    ("天合光能", "天合"),
    ("英发睿能", "英发"),
]

SECTIONS = [
    {
        "name": "追溯能力",
        "short": "追溯能力",
        "items": [
            ("1.1", "体系完整性"),
            ("1.2", "颗粒度&数据链"),
            ("1.3", "采集及时性"),
            ("1.4", "混料隔离能力"),
            ("1.5", "改造成本/难度"),
        ],
    },
    {
        "name": "追溯要求",
        "short": "追溯要求",
        "items": [
            ("2.1", "目标市场业务量"),
            ("2.2", "目标市场合规复杂度"),
        ],
    },
    {
        "name": "认证通过率",
        "short": "认证通过率",
        "items": [
            ("3.1", "机构认证/三方审核"),
            ("3.2", "历史审核通过率"),
        ],
    },
    {
        "name": "供应链管控力",
        "short": "供应链管控力",
        "items": [
            ("4.1", "供应商系统能力"),
            ("4.2", "高风险环节可控"),
            ("4.3", "供应商协同能力"),
            ("4.4", "多级供应商覆盖"),
        ],
    },
]

RAW_SCORES: dict[str, dict[str, int | None]] = {
    "1.1": {"一道新能": 2, "正泰新能": 4, "晶科能源": 5, "晶澳科技": 4, "横店东磁": 1, "无锡博达": 1, "天合光能": 4, "英发睿能": 2},
    "1.2": {"一道新能": 3, "正泰新能": 4, "晶科能源": 5, "晶澳科技": 4, "横店东磁": 3, "无锡博达": 2, "天合光能": 5, "英发睿能": 3},
    "1.3": {"一道新能": 4, "正泰新能": 5, "晶科能源": 5, "晶澳科技": 4, "横店东磁": 2, "无锡博达": 2, "天合光能": 5, "英发睿能": 3},
    "1.4": {"一道新能": None, "正泰新能": 3, "晶科能源": 4, "晶澳科技": 3, "横店东磁": 4, "无锡博达": 5, "天合光能": 3, "英发睿能": None},
    "1.5": {"一道新能": 3, "正泰新能": 2, "晶科能源": 5, "晶澳科技": 5, "横店东磁": 2, "无锡博达": 1, "天合光能": 5, "英发睿能": 2},
    "2.1": {"一道新能": 2, "正泰新能": 3, "晶科能源": 5, "晶澳科技": 5, "横店东磁": 2, "无锡博达": 2, "天合光能": 5, "英发睿能": 2},
    "2.2": {"一道新能": 3, "正泰新能": 5, "晶科能源": 3, "晶澳科技": 5, "横店东磁": 4, "无锡博达": 4, "天合光能": 4, "英发睿能": 3},
    "3.1": {"一道新能": 1, "正泰新能": 4, "晶科能源": 4, "晶澳科技": 4, "横店东磁": 2, "无锡博达": 1, "天合光能": 4, "英发睿能": 1},
    "3.2": {"一道新能": 1, "正泰新能": 5, "晶科能源": 4, "晶澳科技": 4, "横店东磁": 1, "无锡博达": 1, "天合光能": 4, "英发睿能": 1},
    "4.1": {"一道新能": 2, "正泰新能": 4, "晶科能源": 4, "晶澳科技": 4, "横店东磁": 2, "无锡博达": 2, "天合光能": 4, "英发睿能": 2},
    "4.2": {"一道新能": 2, "正泰新能": 4, "晶科能源": 3, "晶澳科技": 3, "横店东磁": 2, "无锡博达": 2, "天合光能": 3, "英发睿能": 2},
    "4.3": {"一道新能": 2, "正泰新能": 4, "晶科能源": 4, "晶澳科技": 4, "横店东磁": 2, "无锡博达": 2, "天合光能": 4, "英发睿能": 2},
    "4.4": {"一道新能": 3, "正泰新能": 4, "晶科能源": 3, "晶澳科技": 3, "横店东磁": 4, "无锡博达": 4, "天合光能": 3, "英发睿能": 3},
}


def rgb(value: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*value)


def score_raw(metric_key: str, vendor: str) -> int | None:
    return RAW_SCORES[metric_key][vendor]


def metric_order() -> list[str]:
    return [metric_key for section in SECTIONS for metric_key, _ in section["items"]]


def average_non_null(metric_keys: list[str], vendor: str) -> float:
    values = [score_raw(key, vendor) for key in metric_keys]
    valid_values = [value for value in values if value is not None]
    return round(mean(valid_values), 2) if valid_values else 0.0


def vendor_stats() -> dict[str, dict[str, float]]:
    capability_keys = ["1.1", "1.2", "1.3", "1.4", "1.5", "4.1", "4.2", "4.3", "4.4"]
    cert_keys = ["2.1", "2.2", "3.1", "3.2"]
    all_keys = metric_order()
    stats: dict[str, dict[str, float]] = {}
    for vendor, _ in VENDORS:
        missing_count = sum(1 for key in all_keys if score_raw(key, vendor) is None)
        coverage = round((len(all_keys) - missing_count) / len(all_keys), 2)
        stats[vendor] = {
            "overall": average_non_null(all_keys, vendor),
            "capability": average_non_null(capability_keys, vendor),
            "cert_market": average_non_null(cert_keys, vendor),
            "missing_count": float(missing_count),
            "coverage": coverage,
        }
    return stats


def vendor_short_name(vendor: str) -> str:
    for full_name, short_name in VENDORS:
        if full_name == vendor:
            return short_name
    return vendor


def ranked_vendors(stats: dict[str, dict[str, float]]) -> list[str]:
    return sorted(
        [vendor for vendor, _ in VENDORS],
        key=lambda vendor: (stats[vendor]["overall"], stats[vendor]["capability"], stats[vendor]["cert_market"]),
        reverse=True,
    )


def quadrant_groups(stats: dict[str, dict[str, float]], threshold: float = 3.0) -> dict[str, list[str]]:
    groups = {
        "优先合作": [],
        "可切入": [],
        "观察池": [],
        "重点补强": [],
    }
    for vendor in ranked_vendors(stats):
        x_value = stats[vendor]["capability"]
        y_value = stats[vendor]["cert_market"]
        if x_value >= threshold and y_value >= threshold:
            groups["优先合作"].append(vendor)
        elif x_value < threshold and y_value >= threshold:
            groups["可切入"].append(vendor)
        elif x_value >= threshold and y_value < threshold:
            groups["观察池"].append(vendor)
        else:
            groups["重点补强"].append(vendor)
    return groups


def data_gap_vendors(stats: dict[str, dict[str, float]]) -> list[str]:
    return [vendor for vendor in ranked_vendors(stats) if stats[vendor]["missing_count"] > 0]


def executive_takeaway(stats: dict[str, dict[str, float]]) -> str:
    groups = quadrant_groups(stats)
    top_names = "、".join(vendor_short_name(vendor) for vendor in groups["优先合作"]) or "暂无"
    gap_names = "、".join(vendor_short_name(vendor) for vendor in data_gap_vendors(stats)) or "无"
    return f"当前样本中，{top_names}位于优先合作象限；{gap_names}存在数据待补项，建议先补齐追溯隔离与认证证据后再进入同口径比较。"


def format_vendor_list(vendors: list[str], stats: dict[str, dict[str, float]], limit: int | None = None) -> str:
    render_vendors = vendors[:limit] if limit else vendors
    if not render_vendors:
        return "暂无"
    return " / ".join(f"{vendor_short_name(vendor)} {stats[vendor]['overall']:.1f}" for vendor in render_vendors)


def add_box(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    text: str = "",
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width_pt: float = 1,
    shape_type: MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE.RECTANGLE,
    font_size: float = 10,
    color: tuple[int, int, int] = NAVY,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.MIDDLE,
    margin_left: float = 0.06,
    margin_right: float = 0.06,
    margin_top: float = 0.03,
    margin_bottom: float = 0.03,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width_pt)
    shape.text_frame.clear()
    shape.text_frame.vertical_anchor = valign
    shape.text_frame.margin_left = Inches(margin_left)
    shape.text_frame.margin_right = Inches(margin_right)
    shape.text_frame.margin_top = Inches(margin_top)
    shape.text_frame.margin_bottom = Inches(margin_bottom)
    if text:
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return shape


def add_text(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    font_size: float,
    color: tuple[int, int, int],
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return textbox


def add_connector(slide, begin_x: Emu, begin_y: Emu, end_x: Emu, end_y: Emu, color: tuple[int, int, int], width_pt: float = 0.8):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, begin_x, begin_y, end_x, end_y)
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width_pt)
    return connector


def set_title(slide, title: str):
    placeholder = slide.shapes.title
    placeholder.text = ""
    text_frame = placeholder.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = FONT_NAME
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = rgb(ACCENT)
    paragraph.alignment = PP_ALIGN.LEFT


def prune_to_single_body_slide(prs: Presentation, body_slide_index: int = 3) -> None:
    for slide_index in range(len(prs.slides) - 1, -1, -1):
        if slide_index != body_slide_index:
            remove_slide(prs, slide_index)


def draw_matrix(slide, left: Emu, top: Emu, width: Emu, height: Emu, stats: dict[str, dict[str, float]]) -> None:
    add_box(
        slide,
        left,
        top,
        width,
        height,
        fill=WHITE,
        line=LIGHT_BORDER,
        line_width_pt=1.2,
        shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    )
    add_text(slide, left + Inches(0.18), top + Inches(0.12), Inches(3.0), Inches(0.22), "评分矩阵", font_size=13, color=NAVY, bold=True)
    add_text(slide, left + width - Inches(2.55), top + Inches(0.12), Inches(2.2), Inches(0.22), "口径：星级折算 0-5 分，空白项不计入均值", font_size=8.5, color=MUTED, align=PP_ALIGN.RIGHT)

    chip_left = left + Inches(0.18)
    chip_top = top + Inches(0.42)
    for index, section in enumerate(SECTIONS):
        chip_width = Inches(1.0 if index < 3 else 1.15)
        add_box(
            slide,
            chip_left,
            chip_top,
            chip_width,
            Inches(0.24),
            text=section["short"],
            fill=SECTION_COLORS[index],
            line=None,
            shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            font_size=8,
            color=NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        chip_left += chip_width + Inches(0.08)

    table_left = left + Inches(0.18)
    table_top = top + Inches(0.80)
    label_width = Inches(2.05)
    col_width = Inches(0.64)
    row_height = Inches(0.28)
    header_height = Inches(0.34)
    avg_height = Inches(0.34)

    add_box(slide, table_left, table_top, label_width, header_height, fill=(246, 248, 250), line=LIGHT_BORDER, text="评估维度", font_size=9, color=SLATE, bold=True)
    current_left = table_left + label_width
    for _, short_name in VENDORS:
        add_box(
            slide,
            current_left,
            table_top,
            col_width,
            header_height,
            fill=(246, 248, 250),
            line=LIGHT_BORDER,
            text=short_name,
            font_size=9,
            color=SLATE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        current_left += col_width

    current_top = table_top + header_height
    section_starts = {section["items"][0][0]: index for index, section in enumerate(SECTIONS)}
    for section_index, section in enumerate(SECTIONS):
        for metric_key, metric_name in section["items"]:
            if metric_key in section_starts:
                add_box(slide, table_left, current_top, label_width + col_width * len(VENDORS), Pt(1.1), fill=LIGHT_BORDER, line=None)
            label_fill = SECTION_COLORS[section_index] if metric_key == section["items"][0][0] else WHITE
            add_box(
                slide,
                table_left,
                current_top,
                label_width,
                row_height,
                fill=label_fill,
                line=LIGHT_BORDER,
                text=f"{metric_key}  {metric_name}",
                font_size=8.5,
                color=NAVY,
                bold=metric_key == section["items"][0][0],
            )
            current_left = table_left + label_width
            for vendor, _ in VENDORS:
                raw_value = score_raw(metric_key, vendor)
                numeric = raw_value if raw_value is not None else 0
                add_box(
                    slide,
                    current_left,
                    current_top,
                    col_width,
                    row_height,
                    fill=HEAT_COLORS[numeric],
                    line=WHITE,
                    line_width_pt=1,
                    text="待补" if raw_value is None else str(raw_value),
                    font_size=7.8 if raw_value is None else 9,
                    color=WHITE if numeric >= 4 else NAVY,
                    bold=True,
                    align=PP_ALIGN.CENTER,
                )
                current_left += col_width
            current_top += row_height

    add_box(slide, table_left, current_top + Inches(0.04), label_width, avg_height, fill=(244, 246, 249), line=LIGHT_BORDER, text="总评均值", font_size=9, color=NAVY, bold=True)
    current_left = table_left + label_width
    for vendor, _ in VENDORS:
        score = stats[vendor]["overall"]
        bucket = min(5, max(0, round(score)))
        add_box(
            slide,
            current_left,
            current_top + Inches(0.04),
            col_width,
            avg_height,
            fill=HEAT_COLORS[bucket],
            line=WHITE,
            text=f"{score:.1f}",
            font_size=8.5,
            color=WHITE if bucket >= 4 else NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        current_left += col_width


def draw_summary(slide, left: Emu, top: Emu, width: Emu, stats: dict[str, dict[str, float]]) -> None:
    groups = quadrant_groups(stats)
    gaps = data_gap_vendors(stats)
    add_box(
        slide,
        left,
        top,
        width,
        Inches(1.60),
        fill=WHITE,
        line=LIGHT_BORDER,
        line_width_pt=1.2,
        shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    )
    add_text(slide, left + Inches(0.18), top + Inches(0.10), Inches(1.8), Inches(0.22), "结论摘要", font_size=13, color=NAVY, bold=True)
    rows = [
        ("优先合作", format_vendor_list(groups["优先合作"], stats), ACCENT, WHITE),
        ("重点补强", format_vendor_list(groups["重点补强"], stats), (237, 241, 244), NAVY),
        ("数据待补", " / ".join(vendor_short_name(vendor) for vendor in gaps) or "无", (247, 234, 238), ACCENT),
    ]
    current_top = top + Inches(0.40)
    for label, content, fill_color, text_color in rows:
        add_box(
            slide,
            left + Inches(0.18),
            current_top,
            Inches(0.86),
            Inches(0.26),
            text=label,
            fill=fill_color,
            line=None,
            shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            font_size=8.5,
            color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, left + Inches(1.10), current_top, width - Inches(1.28), Inches(0.26), content, font_size=9.3, color=NAVY)
        current_top += Inches(0.33)
    add_text(slide, left + Inches(0.18), top + Inches(1.38), width - Inches(0.36), Inches(0.12), "注：均值按非空项计算，避免把信息缺失直接当作低分。", font_size=7.2, color=MUTED)


def rects_overlap(rect1: tuple[Emu, Emu, Emu, Emu], rect2: tuple[Emu, Emu, Emu, Emu]) -> bool:
    left1, top1, width1, height1 = rect1
    left2, top2, width2, height2 = rect2
    return not (
        left1 + width1 <= left2
        or left2 + width2 <= left1
        or top1 + height1 <= top2
        or top2 + height2 <= top1
    )


def choose_label_rect(
    plot_rect: tuple[Emu, Emu, Emu, Emu],
    center_x: Emu,
    center_y: Emu,
    bubble_size: Emu,
    label_width: Emu,
    label_height: Emu,
    existing_rects: list[tuple[Emu, Emu, Emu, Emu]],
) -> tuple[Emu, Emu, Emu, Emu]:
    plot_left, plot_top, plot_width, plot_height = plot_rect
    candidates = [
        (center_x + bubble_size // 2 + Inches(0.05), center_y - label_height - Inches(0.01)),
        (center_x - label_width - bubble_size // 2 - Inches(0.05), center_y - label_height - Inches(0.01)),
        (center_x + bubble_size // 2 + Inches(0.05), center_y + Inches(0.01)),
        (center_x - label_width - bubble_size // 2 - Inches(0.05), center_y + Inches(0.01)),
        (center_x - label_width // 2, center_y - label_height - Inches(0.12)),
        (center_x - label_width // 2, center_y + Inches(0.08)),
    ]
    best_rect = (candidates[0][0], candidates[0][1], label_width, label_height)
    best_penalty = float("inf")
    for label_left, label_top in candidates:
        rect = (label_left, label_top, label_width, label_height)
        penalty = 0.0
        if label_left < plot_left:
            penalty += 10 + (plot_left - label_left) / 10000
        if label_top < plot_top:
            penalty += 10 + (plot_top - label_top) / 10000
        if label_left + label_width > plot_left + plot_width:
            penalty += 10 + (label_left + label_width - (plot_left + plot_width)) / 10000
        if label_top + label_height > plot_top + plot_height:
            penalty += 10 + (label_top + label_height - (plot_top + plot_height)) / 10000
        penalty += abs(label_left - center_x) / 100000 + abs(label_top - center_y) / 100000
        for existing_rect in existing_rects:
            if rects_overlap(rect, existing_rect):
                penalty += 25
        if penalty < best_penalty:
            best_penalty = penalty
            best_rect = rect
    return best_rect


def draw_quadrant(slide, left: Emu, top: Emu, width: Emu, height: Emu, stats: dict[str, dict[str, float]]) -> None:
    add_box(
        slide,
        left,
        top,
        width,
        height,
        fill=WHITE,
        line=LIGHT_BORDER,
        line_width_pt=1.2,
        shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    )
    add_text(slide, left + Inches(0.18), top + Inches(0.10), Inches(1.8), Inches(0.22), "象限分布", font_size=13, color=NAVY, bold=True)
    add_text(slide, left + width - Inches(1.72), top + Inches(0.10), Inches(1.48), Inches(0.22), "横轴: 追溯基础  纵轴: 认证适配", font_size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)

    plot_left = left + Inches(0.28)
    plot_top = top + Inches(0.48)
    plot_width = width - Inches(0.46)
    plot_height = height - Inches(0.88)
    plot_rect = (plot_left, plot_top, plot_width, plot_height)

    add_box(slide, plot_left, plot_top, plot_width, plot_height, fill=WHITE, line=LIGHT_BORDER)
    add_box(slide, plot_left + plot_width // 2, plot_top, plot_width // 2, plot_height // 2, fill=QUADRANT_FILL, line=None)
    add_box(slide, plot_left, plot_top + plot_height // 2, plot_width // 2, plot_height // 2, fill=LOW_FILL, line=None)

    for ratio in (0.0, 0.5, 1.0):
        x = plot_left + int(plot_width * ratio)
        y = plot_top + int(plot_height * (1 - ratio))
        add_box(slide, plot_left, y, plot_width, Pt(0.8), fill=(233, 236, 240), line=None)
        add_box(slide, x, plot_top, Pt(0.8), plot_height, fill=(233, 236, 240), line=None)

    threshold = 3.0
    threshold_x = plot_left + int(plot_width * (threshold / 5))
    threshold_y = plot_top + int(plot_height * (1 - threshold / 5))
    add_box(slide, threshold_x, plot_top, Pt(1.2), plot_height, fill=ACCENT, line=None)
    add_box(slide, plot_left, threshold_y, plot_width, Pt(1.2), fill=ACCENT, line=None)

    corner_labels = [
        ("可切入", plot_left + Inches(0.10), plot_top + Inches(0.05), NAVY, (242, 246, 249)),
        ("优先合作", plot_left + plot_width - Inches(0.85), plot_top + Inches(0.05), ACCENT, (248, 232, 237)),
        ("重点补强", plot_left + Inches(0.10), plot_top + plot_height - Inches(0.28), NAVY, (242, 246, 249)),
        ("观察池", plot_left + plot_width - Inches(0.70), plot_top + plot_height - Inches(0.28), NAVY, (242, 246, 249)),
    ]
    for label, label_left, label_top, text_color, fill_color in corner_labels:
        add_box(
            slide,
            label_left,
            label_top,
            Inches(0.58 if len(label) == 3 else 0.66),
            Inches(0.18),
            text=label,
            fill=fill_color,
            line=None,
            shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            font_size=7.3,
            color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    add_text(slide, plot_left - Inches(0.04), plot_top + plot_height + Inches(0.02), Inches(0.20), Inches(0.16), "0", font_size=7, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, threshold_x - Inches(0.03), plot_top + plot_height + Inches(0.02), Inches(0.22), Inches(0.16), "3", font_size=7, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, plot_left + plot_width - Inches(0.12), plot_top + plot_height + Inches(0.02), Inches(0.22), Inches(0.16), "5", font_size=7, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, plot_left - Inches(0.18), plot_top - Inches(0.02), Inches(0.14), Inches(0.16), "5", font_size=7, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, plot_left - Inches(0.18), threshold_y - Inches(0.07), Inches(0.14), Inches(0.16), "3", font_size=7, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(slide, plot_left - Inches(0.18), plot_top + plot_height - Inches(0.16), Inches(0.14), Inches(0.16), "0", font_size=7, color=MUTED, align=PP_ALIGN.RIGHT)

    label_rects: list[tuple[Emu, Emu, Emu, Emu]] = []
    for vendor in ranked_vendors(stats):
        x_value = stats[vendor]["capability"]
        y_value = stats[vendor]["cert_market"]
        overall = stats[vendor]["overall"]
        center_x = plot_left + int(plot_width * (x_value / 5))
        center_y = plot_top + int(plot_height * (1 - y_value / 5))
        bubble_size = Inches(0.18 + max(0.0, overall - 2.0) * 0.02)
        bubble_left = center_x - bubble_size // 2
        bubble_top = center_y - bubble_size // 2
        bubble_fill = ACCENT if overall >= 3.8 else (136, 149, 165)
        add_box(
            slide,
            bubble_left,
            bubble_top,
            bubble_size,
            bubble_size,
            text="",
            fill=bubble_fill,
            line=WHITE,
            line_width_pt=1.2,
            shape_type=MSO_AUTO_SHAPE_TYPE.OVAL,
        )

        label_width = Inches(0.44)
        label_height = Inches(0.16)
        label_rect = choose_label_rect(plot_rect, center_x, center_y, bubble_size, label_width, label_height, label_rects)
        label_rects.append(label_rect)
        label_left, label_top, _, _ = label_rect
        label_center_x = label_left + label_width // 2
        label_center_y = label_top + label_height // 2
        add_connector(slide, center_x, center_y, label_center_x, label_center_y, color=(180, 188, 197), width_pt=0.8)
        add_text(slide, label_left, label_top, label_width, label_height, vendor_short_name(vendor), font_size=7.5, color=NAVY, bold=True)


def add_footer_note(slide, stats: dict[str, dict[str, float]]) -> None:
    gap_names = "、".join(vendor_short_name(vendor) for vendor in data_gap_vendors(stats)) or "无"
    note = f"方法注释：空白项按“待补”处理，不直接折算低分；当前存在数据待补的厂商为 {gap_names}。"
    add_text(slide, Inches(0.52), Inches(6.88), Inches(12.0), Inches(0.20), note, font_size=7.5, color=MUTED)


def build_slide(output_path: Path, template_path: Path) -> Path:
    prs = Presentation(str(template_path))
    prune_to_single_body_slide(prs, body_slide_index=3)
    slide = prs.slides[0]
    stats = vendor_stats()

    set_title(slide, "光伏组件企业追溯能力与认证可行性对标")
    add_text(slide, Inches(0.52), Inches(0.60), Inches(3.9), Inches(0.20), "SIE TRACEABILITY BENCHMARK  |  8家企业 · 13项指标 · 单页筛选视图", font_size=8.5, color=SLATE, bold=True)
    add_box(slide, Inches(0.52), Inches(0.92), Inches(0.74), Inches(0.22), fill=ACCENT, line=None, shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, text="核心判断", font_size=8.3, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.34), Inches(0.92), Inches(11.2), Inches(0.22), executive_takeaway(stats), font_size=8.5, color=NAVY)

    draw_matrix(slide, Inches(0.50), Inches(1.22), Inches(7.95), Inches(5.48), stats)
    draw_summary(slide, Inches(8.62), Inches(1.22), Inches(3.70), stats)
    draw_quadrant(slide, Inches(8.62), Inches(2.96), Inches(3.70), Inches(3.74), stats)
    add_footer_note(slide, stats)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        fallback_output = output_path.with_stem(f"{output_path.stem}_{timestamp}")
        prs.save(str(fallback_output))
        return fallback_output


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a one-page SIE-style traceability benchmark slide.")
    parser.add_argument("--output", default="output/sie_traceability_benchmark_onepage.pptx", help="Output PPTX path.")
    parser.add_argument("--template", default=str(DEFAULT_TEMPLATE), help="Template PPTX path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output_path = Path(args.output).resolve()
    template_path = Path(args.template).resolve()
    build_slide(output_path, template_path)
    print(output_path)


if __name__ == "__main__":
    main()
