from __future__ import annotations

import argparse
import json
import platform
import subprocess
from dataclasses import asdict, dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE


DEFAULT_BANNED_PHRASES = (
    "讲解重点",
    "关键用户在UAT前先看懂",
    "先看懂三件事",
    "建议优先验证",
    "UAT介绍",
    "讲解目标",
)


@dataclass(frozen=True)
class ReviewFinding:
    level: str
    dimension: str
    message: str

    def to_dict(self) -> dict[str, str]:
        return asdict(self)


@dataclass(frozen=True)
class ReviewReport:
    passed: bool
    findings: tuple[ReviewFinding, ...]
    preview_paths: tuple[str, ...]
    preview_note: str

    def to_dict(self) -> dict[str, object]:
        return {
            "passed": self.passed,
            "findings": [finding.to_dict() for finding in self.findings],
            "preview_paths": list(self.preview_paths),
            "preview_note": self.preview_note,
        }


def _iter_text_shapes(slide):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if text:
            yield shape, text


def _estimate_lines(text: str, width: int, font_size_pt: float) -> int:
    weighted_chars = 0.0
    for char in text:
        if char.isspace():
            weighted_chars += 0.35
        elif ord(char) < 128:
            weighted_chars += 0.58
        else:
            weighted_chars += 1.0
    chars_per_line = max(1.0, width / max(font_size_pt * 7000, 1))
    return max(1, round(weighted_chars / chars_per_line + 0.499))


def _font_size_points(shape) -> float:
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return float(run.font.size.pt)
    return 12.0


def _line_capacity(shape, font_size_pt: float) -> int:
    emu_per_line = max(font_size_pt * 12700 * 1.65, 1)
    return max(1, int(shape.height / emu_per_line))


def _review_bounds(prs: Presentation) -> list[ReviewFinding]:
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    findings: list[ReviewFinding] = []
    for idx, shape in enumerate(slide.shapes, start=1):
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_width or shape.top + shape.height > slide_height:
            findings.append(ReviewFinding("blocker", "layout", f"shape {idx} exceeds slide bounds"))
    return findings


def _review_pictures(prs: Presentation) -> list[ReviewFinding]:
    slide = prs.slides[0]
    if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
        return [ReviewFinding("warning", "deliverability", "slide contains picture shapes; editable output is preferred")]
    return []


def _review_meta_phrases(prs: Presentation, banned_phrases: tuple[str, ...]) -> list[ReviewFinding]:
    slide = prs.slides[0]
    slide_text = "\n".join(text for _, text in _iter_text_shapes(slide))
    findings: list[ReviewFinding] = []
    for phrase in banned_phrases:
        if phrase and phrase in slide_text:
            findings.append(ReviewFinding("warning", "content", f"audience-facing slide still contains meta guidance phrase: {phrase}"))
    return findings


def _review_density(prs: Presentation) -> list[ReviewFinding]:
    slide = prs.slides[0]
    findings: list[ReviewFinding] = []
    for _, text in _iter_text_shapes(slide):
        pass
    for index, (shape, text) in enumerate(_iter_text_shapes(slide), start=1):
        font_size = _font_size_points(shape)
        estimated_lines = _estimate_lines(text, int(shape.width), font_size)
        capacity = _line_capacity(shape, font_size)
        if estimated_lines > capacity + 1:
            findings.append(
                ReviewFinding(
                    "warning",
                    "content_density",
                    f"text box {index} may be too dense for its area ({estimated_lines} estimated lines vs {capacity} capacity)",
                )
            )
    return findings


def _review_card_balance(prs: Presentation, expected_card_count: int) -> list[ReviewFinding]:
    if expected_card_count <= 1:
        return []

    slide = prs.slides[0]
    candidates = []
    for shape in slide.shapes:
        try:
            auto_shape_type = shape.auto_shape_type
        except ValueError:
            continue
        if auto_shape_type != MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
            continue
        if shape.width < 2000000 or shape.height < 1800000:
            continue
        candidates.append(shape)

    findings: list[ReviewFinding] = []
    if len(candidates) < expected_card_count:
        findings.append(
            ReviewFinding(
                "warning",
                "layout",
                f"expected around {expected_card_count} main cards, but only found {len(candidates)} large rounded rectangles",
            )
        )
        return findings

    cards = sorted(candidates, key=lambda item: (item.top, item.left))[:expected_card_count]
    widths = [int(card.width) for card in cards]
    heights = [int(card.height) for card in cards]
    gaps = [int(cards[idx + 1].left - (cards[idx].left + cards[idx].width)) for idx in range(len(cards) - 1)]
    if max(widths) - min(widths) > 80000:
        findings.append(ReviewFinding("warning", "layout", "main card widths are not visually balanced"))
    if max(heights) - min(heights) > 80000:
        findings.append(ReviewFinding("warning", "layout", "main card heights are not visually balanced"))
    if gaps and max(gaps) - min(gaps) > 80000:
        findings.append(ReviewFinding("warning", "layout", "horizontal gaps between main cards are inconsistent"))

    slide_width = int(prs.slide_width)
    left_margin = cards[0].left
    right_margin = slide_width - (cards[-1].left + cards[-1].width)
    if abs(left_margin - right_margin) > 100000:
        findings.append(ReviewFinding("warning", "layout", "page margins around the main card group are not symmetric"))
    return findings


def export_slide_previews(pptx_path: Path, output_dir: Path) -> tuple[tuple[str, ...], str]:
    output_dir.mkdir(parents=True, exist_ok=True)
    system = platform.system().lower()
    try:
        if system == "windows":
            script = (
                "$pptPath = Resolve-Path $args[0];"
                "$outDir = $args[1];"
                "New-Item -ItemType Directory -Path $outDir -Force | Out-Null;"
                "$ppt = New-Object -ComObject PowerPoint.Application;"
                "$ppt.Visible = -1;"
                "$pres = $ppt.Presentations.Open($pptPath.Path, $false, $false, $false);"
                "$idx = 1;"
                "foreach ($slide in $pres.Slides) {"
                "  $target = Join-Path $outDir ('slide' + $idx + '.png');"
                "  $slide.Export($target, 'PNG', 1920, 1080);"
                "  $idx += 1"
                "};"
                "$pres.Close();"
                "$ppt.Quit();"
            )
            result = subprocess.run(
                ["powershell", "-NoProfile", "-Command", script, str(pptx_path), str(output_dir)],
                text=True,
                capture_output=True,
                check=False,
            )
        else:
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "png", "--outdir", str(output_dir), str(pptx_path)],
                text=True,
                capture_output=True,
                check=False,
            )
        if result.returncode != 0:
            note = (result.stderr or result.stdout or "").strip() or "preview export failed"
            return (), note
        previews = tuple(str(path) for path in sorted(output_dir.glob("slide*.png")))
        if not previews:
            return (), "no preview images were generated"
        return previews, "preview export succeeded"
    except Exception as exc:  # pragma: no cover - defensive best effort
        return (), f"preview export unavailable: {exc}"


def review_onepage_slide(
    pptx_path: Path,
    *,
    expected_card_count: int = 3,
    banned_phrases: tuple[str, ...] = DEFAULT_BANNED_PHRASES,
    export_previews: bool = False,
    preview_dir: Path | None = None,
) -> ReviewReport:
    prs = Presentation(str(pptx_path))
    findings = [
        *_review_bounds(prs),
        *_review_pictures(prs),
        *_review_meta_phrases(prs, banned_phrases),
        *_review_density(prs),
        *_review_card_balance(prs, expected_card_count),
    ]

    preview_paths: tuple[str, ...] = ()
    preview_note = "preview export skipped"
    if export_previews:
        preview_paths, preview_note = export_slide_previews(pptx_path, preview_dir or pptx_path.with_suffix(""))
        if not preview_paths:
            findings.append(ReviewFinding("warning", "deliverability", f"preview export unavailable: {preview_note}"))

    passed = not any(finding.level == "blocker" for finding in findings)
    return ReviewReport(
        passed=passed,
        findings=tuple(findings),
        preview_paths=preview_paths,
        preview_note=preview_note,
    )


def write_review_report(report: ReviewReport, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(report.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
    return output_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run heuristic review checks for a one-page PPT.")
    parser.add_argument("pptx", help="PPTX file to review.")
    parser.add_argument("--review-json", default="", help="Optional JSON output path for the review report.")
    parser.add_argument("--expected-card-count", type=int, default=3, help="Expected number of main cards in the slide.")
    parser.add_argument("--export-previews", action="store_true", help="Export PNG previews when supported by the local environment.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    pptx_path = Path(args.pptx).resolve()
    report = review_onepage_slide(
        pptx_path,
        expected_card_count=args.expected_card_count,
        export_previews=args.export_previews,
        preview_dir=pptx_path.parent / f"{pptx_path.stem}_previews",
    )
    if args.review_json:
        write_review_report(report, Path(args.review_json).resolve())
    print(json.dumps(report.to_dict(), ensure_ascii=False, indent=2))
    return 0 if report.passed else 1


if __name__ == "__main__":
    raise SystemExit(main())
