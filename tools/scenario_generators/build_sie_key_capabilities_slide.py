from __future__ import annotations

import argparse
import datetime
from dataclasses import dataclass
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

try:
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.onepage_layout_presets import DEFAULT_ONEPAGE_PRESET_ID, get_onepage_layout_preset
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest
except ModuleNotFoundError:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.onepage_layout_presets import DEFAULT_ONEPAGE_PRESET_ID, get_onepage_layout_preset
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest


PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_PPT = PROJECT_ROOT / "output" / "sie_supply_chain_key_capabilities.pptx"

ACCENT = (173, 5, 61)
INK = (35, 41, 46)
MUTED = (98, 109, 121)
LIGHT_TEXT = (145, 152, 160)
LINE = (224, 228, 232)
WHITE = (255, 255, 255)
CARD_STYLES = (
    {"accent": (118, 86, 167), "soft": (247, 242, 250)},
    {"accent": ACCENT, "soft": (252, 242, 246)},
    {"accent": (53, 116, 136), "soft": (241, 247, 249)},
)


@dataclass(frozen=True)
class CapabilityCard:
    index: str
    english: str
    chinese: str
    definition: str
    bullets: tuple[str, str, str]


@dataclass(frozen=True)
class Rect:
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    def overlaps(self, other: "Rect") -> bool:
        return not (
            self.right <= other.left
            or other.right <= self.left
            or self.bottom <= other.top
            or other.bottom <= self.top
        )


TITLE_SEGMENTS = [
    ("关键能力：", ACCENT, True),
    ("追溯能力不止于“可查”，", INK, True),
    ("更在于“可证”“可控”", ACCENT, True),
]
SUBTITLE = "真正有效的供应链追溯能力，不止于信息可查询，更在于链路可还原、证据可支撑、风险可响应。"
SUMMARY_LABEL = "核心判断"
SUMMARY_LINE_1 = "面向外部审查的供应链能力，本质上要同时回答三个问题："
SUMMARY_LINE_2 = "链路能否还原、证据能否支撑、风险能否响应。"
FOOTNOTE = "建议作为对外沟通的统一框架：先说明链路可追，再证明证据可证，最后展示风险可控。"

CARDS = (
    CapabilityCard(
        index="01",
        english="TRACEABILITY",
        chinese="链路可追",
        definition="能力定义：实现关键供应链链路的端到端可追溯与可还原。",
        bullets=(
            "打通从矿源、原材料、生产过程到成品交付的关键链路。",
            "支撑按订单、SN、批次等维度进行精准追溯与路径还原。",
            "形成跨供应商、跨工厂、跨系统的一体化追溯视图。",
        ),
    ),
    CapabilityCard(
        index="02",
        english="PROOF",
        chinese="证据可证",
        definition="能力定义：形成支撑外部审查的数据与证明材料闭环体系。",
        bullets=(
            "将业务数据、业务单据、合规文件建立一一对应与相互勾稽关系。",
            "保证数据与证明材料具备一致性、可校验性与可审计性。",
            "支撑客户尽调、海关检查、审计核查等对外证明场景。",
        ),
    ),
    CapabilityCard(
        index="03",
        english="CONTROL",
        chinese="风险可控",
        definition="能力定义：建立面向合规风险的快速识别与响应机制。",
        bullets=(
            "快速定位风险物料、风险批次与风险供应商，缩小排查范围。",
            "在客户抽查、审查升级时，能够快速响应、快速说明、快速出具材料。",
            "支撑风险隔离、经营决策与海外市场准入过程中的稳定交付。",
        ),
    ),
)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def set_run_style(run, *, size: float, color: tuple[int, int, int], bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def estimate_lines(text: str, width: int, font_size: float) -> int:
    weighted_chars = 0.0
    for char in text:
        if char.isspace():
            weighted_chars += 0.35
        elif ord(char) < 128:
            weighted_chars += 0.58
        else:
            weighted_chars += 1.0
    chars_per_line = max(1.0, width / max(font_size * 7000, 1))
    return max(1, math.ceil(weighted_chars / chars_per_line))


def add_text(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: float,
    color: tuple[int, int, int] = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=font_size, color=color, bold=bold)
    return shape


def add_shape(
    slide,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width: float = 1,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def keep_only_body_template(prs: Presentation) -> None:
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    keep_idx = manifest.slide_roles.body_template
    for index in range(len(prs.slides) - 1, -1, -1):
        if index != keep_idx:
            remove_slide(prs, index)


def render_title(slide, preset_id: str) -> None:
    preset = get_onepage_layout_preset(preset_id)
    title_shape = slide.shapes[0]
    tf = title_shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    for text, color, bold in TITLE_SEGMENTS:
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size=float(preset.renderer_hints["title_font_size"]), color=color, bold=bold)

    add_text(
        slide,
        720000,
        540000,
        8500000,
        190000,
        SUBTITLE,
        font_size=float(preset.renderer_hints["subtitle_font_size"]),
        color=MUTED,
    )


def render_summary(slide, slide_width: int, preset_id: str) -> None:
    preset = get_onepage_layout_preset(preset_id)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 1140000, 70000, 270000, fill=ACCENT)
    add_text(
        slide,
        850000,
        1135000,
        1500000,
        180000,
        SUMMARY_LABEL,
        font_size=float(preset.renderer_hints["summary_label_font_size"]),
        color=ACCENT,
        bold=True,
    )
    add_text(
        slide,
        720000,
        1400000,
        slide_width - 1440000,
        240000,
        SUMMARY_LINE_1,
        font_size=float(preset.renderer_hints["summary_intro_font_size"]),
        color=MUTED,
    )
    add_text(
        slide,
        720000,
        1600000,
        slide_width - 1440000,
        340000,
        SUMMARY_LINE_2,
        font_size=float(preset.renderer_hints["summary_headline_font_size"]),
        color=INK,
        bold=True,
    )
    if bool(preset.renderer_hints.get("show_support_line", False)):
        add_text(slide, 720000, 1880000, slide_width - 1440000, 170000, "换句话说，页面应以结论先行、论据跟进的方式组织阅读。", font_size=10.8, color=LIGHT_TEXT)


def card_geometries(slide_width: int, preset_id: str) -> list[tuple[int, int, int, int]]:
    preset = get_onepage_layout_preset(preset_id)
    left = 720000
    gap = int(preset.renderer_hints["card_gap"])
    usable_width = slide_width - left * 2
    width = int((usable_width - gap * 2) / 3)
    top = int(preset.renderer_hints["card_top"])
    height = int(preset.renderer_hints["card_height"])
    return [(left + idx * (width + gap), top, width, height) for idx in range(3)]


def render_card(slide, geom: tuple[int, int, int, int], spec: CapabilityCard, style: dict[str, tuple[int, int, int]], preset_id: str) -> None:
    preset = get_onepage_layout_preset(preset_id)
    left, top, width, height = geom
    accent = style["accent"]
    soft = style["soft"]

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=WHITE, line=LINE, line_width=1)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, 70000, fill=accent)
    add_text(slide, left + 180000, top + 150000, 280000, 120000, spec.index, font_size=10.5, color=LIGHT_TEXT, bold=True)
    add_text(
        slide,
        left + width - 1100000,
        top + 150000,
        900000,
        120000,
        spec.english,
        font_size=float(preset.renderer_hints["card_english_font_size"]),
        color=accent,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + 180000, top + 320000, 350000, 350000, fill=accent)
    add_text(slide, left + 180000, top + 395000, 350000, 110000, spec.index[-1], font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        left + 620000,
        top + 325000,
        1200000,
        150000,
        spec.chinese,
        font_size=float(preset.renderer_hints["card_title_font_size"]),
        color=INK,
        bold=True,
    )

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 180000, top + 770000, width - 360000, 380000, fill=soft)
    add_text(
        slide,
        left + 230000,
        top + 860000,
        width - 460000,
        180000,
        spec.definition,
        font_size=float(preset.renderer_hints["card_definition_font_size"]),
        color=accent,
        bold=True,
    )

    bullet_y = top + 1320000
    for bullet in spec.bullets:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + 190000, bullet_y + 65000, 65000, 65000, fill=accent)
        add_text(
            slide,
            left + 310000,
            bullet_y,
            width - 470000,
            300000,
            bullet,
            font_size=float(preset.renderer_hints["card_bullet_font_size"]),
            color=MUTED,
        )
        bullet_y += 520000


def render_footer(slide, slide_width: int, preset_id: str) -> None:
    preset = get_onepage_layout_preset(preset_id)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 5230000, slide_width - 1440000, 1, fill=LINE)
    add_text(
        slide,
        720000,
        5350000,
        slide_width - 1440000,
        160000,
        FOOTNOTE,
        font_size=float(preset.renderer_hints["footer_font_size"]),
        color=MUTED,
    )


def self_check_layout(prs: Presentation, preset_id: str = DEFAULT_ONEPAGE_PRESET_ID) -> list[str]:
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    preset = get_onepage_layout_preset(preset_id)
    issues: list[str] = []

    if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
        issues.append("slide contains picture shapes; expected editable-only output")

    for idx, shape in enumerate(slide.shapes, 1):
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_width or shape.top + shape.height > slide_height:
            issues.append(f"shape {idx} exceeds slide bounds")

    cards = [Rect(*geom) for geom in card_geometries(slide_width, preset_id)]
    for i, current in enumerate(cards):
        for j, other in enumerate(cards[i + 1 :], start=i + 1):
            if current.overlaps(other):
                issues.append(f"card {i + 1} overlaps card {j + 1}")

    summary = Rect(720000, 1140000, slide_width - 1440000, 800000)
    for idx, card in enumerate(cards, 1):
        if summary.overlaps(card):
            issues.append(f"summary region overlaps card {idx}")

    footer = Rect(720000, 5230000, slide_width - 1440000, 320000)
    for idx, card in enumerate(cards, 1):
        if footer.overlaps(card):
            issues.append(f"footer region overlaps card {idx}")

    title_text = "".join(segment[0] for segment in TITLE_SEGMENTS)
    if estimate_lines(title_text, 10034086, float(preset.renderer_hints["title_font_size"])) > 1:
        issues.append("title is likely to wrap to more than one line")

    if estimate_lines(SUMMARY_LINE_2, slide_width - 1440000, float(preset.renderer_hints["summary_headline_font_size"])) > 2:
        issues.append("headline is too dense for the allocated summary box")

    for spec in CARDS:
        if estimate_lines(spec.definition, 2500000, float(preset.renderer_hints["card_definition_font_size"])) > 2:
            issues.append(f"{spec.chinese} definition is too dense")
        for bullet in spec.bullets:
            if estimate_lines(bullet, 2500000, float(preset.renderer_hints["card_bullet_font_size"])) > 2:
                issues.append(f"{spec.chinese} bullet is too dense")

    return issues


def build_slide(output_path: Path, preset_id: str = DEFAULT_ONEPAGE_PRESET_ID) -> Path:
    prs = Presentation(str(DEFAULT_TEMPLATE))
    keep_only_body_template(prs)
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)

    render_title(slide, preset_id)
    render_summary(slide, slide_width, preset_id)
    for geom, spec, style in zip(card_geometries(slide_width, preset_id), CARDS, CARD_STYLES):
        render_card(slide, geom, spec, style, preset_id)
    render_footer(slide, slide_width, preset_id)

    issues = self_check_layout(prs, preset_id=preset_id)
    if issues:
        raise ValueError("layout self-check failed: " + " | ".join(issues))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        timestamp = datetime.datetime.now().strftime("%H%M%S_%f")[:9]
        fallback = output_path.with_stem(f"{output_path.stem}_{timestamp}")
        prs.save(str(fallback))
        return fallback


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a consulting-style SIE slide for supply-chain traceability key capabilities.")
    parser.add_argument("--output", default=str(OUTPUT_PPT), help="Output PPTX path.")
    parser.add_argument("--layout-preset", default=DEFAULT_ONEPAGE_PRESET_ID, help="One-page layout preset id.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    get_onepage_layout_preset(args.layout_preset)
    output = build_slide(Path(args.output).resolve(), preset_id=args.layout_preset)
    print(output)


if __name__ == "__main__":
    main()
