from __future__ import annotations

import argparse
import datetime
import shutil
import subprocess
import tempfile
from pathlib import Path
from urllib.parse import quote

from pptx import Presentation
from pptx.util import Inches


PROJECT_ROOT = Path(__file__).resolve().parents[2]
DEFAULT_OUTPUT_DIR = PROJECT_ROOT / "output"


def resolve_browser_path(explicit_path: str = "") -> Path:
    candidates = [Path(explicit_path)] if explicit_path else []
    candidates.extend(
        [
            Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
            Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
            Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
        ]
    )
    for candidate in candidates:
        if candidate and candidate.exists():
            return candidate
    raise FileNotFoundError("Could not find Edge/Chrome. Please pass --browser with a valid path.")


def html_file_to_url(html_path: Path) -> str:
    resolved = html_path.resolve().as_posix()
    return f"file:///{quote(resolved, safe=':/()_-.,')}"


def capture_html_screenshot(
    *,
    browser_path: Path,
    html_path: Path,
    screenshot_path: Path,
    window_width: int = 1280,
    window_height: int = 760,
) -> Path:
    screenshot_path.parent.mkdir(parents=True, exist_ok=True)
    url = html_file_to_url(html_path)
    command = [
        str(browser_path),
        "--headless",
        f"--screenshot={screenshot_path}",
        f"--window-size={window_width},{window_height}",
        "--hide-scrollbars",
        "--disable-gpu",
        url,
    ]
    result = subprocess.run(command, capture_output=True, text=True, check=False)
    if result.returncode != 0 or not screenshot_path.exists():
        detail = (result.stderr or result.stdout or "").strip()
        raise RuntimeError(f"Failed to capture HTML screenshot. {detail or 'No output from browser.'}")
    return screenshot_path


def build_single_page_ppt(
    *,
    image_path: Path,
    output_path: Path,
) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def default_output_path(output_dir: Path, stem: str) -> Path:
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
    safe_stem = "".join(char if char not in '<>:"/\\|?*' else "_" for char in stem).strip(" ._") or "html_snapshot"
    return output_dir / f"{safe_stem}_{timestamp}.pptx"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Render a local HTML page as a high-fidelity single-slide PPT.")
    parser.add_argument("--html", required=True, help="Path to local HTML file.")
    parser.add_argument("--output", default="", help="Optional output PPTX path.")
    parser.add_argument("--browser", default="", help="Optional Edge/Chrome executable path.")
    parser.add_argument("--width", type=int, default=1280, help="Browser viewport width for screenshot.")
    parser.add_argument("--height", type=int, default=760, help="Browser viewport height for screenshot.")
    parser.add_argument("--keep-image", action="store_true", help="Keep the intermediate screenshot PNG.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    html_path = Path(args.html).resolve()
    if not html_path.exists():
        raise FileNotFoundError(f"HTML file not found: {html_path}")

    browser_path = resolve_browser_path(args.browser)
    output_path = Path(args.output).resolve() if args.output else default_output_path(DEFAULT_OUTPUT_DIR, html_path.stem)

    with tempfile.TemporaryDirectory() as temp_dir:
        temp_png = Path(temp_dir) / f"{html_path.stem}.png"
        capture_html_screenshot(
            browser_path=browser_path,
            html_path=html_path,
            screenshot_path=temp_png,
            window_width=args.width,
            window_height=args.height,
        )
        build_single_page_ppt(image_path=temp_png, output_path=output_path)
        if args.keep_image:
            shutil.copy2(temp_png, output_path.with_suffix(".png"))

    print(str(output_path))


if __name__ == "__main__":
    main()
