from __future__ import annotations

import argparse
import datetime
import math
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

try:
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.onepage_layout_presets import get_onepage_layout_preset
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest
except ModuleNotFoundError:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.onepage_layout_presets import get_onepage_layout_preset
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest

try:
    from .review_onepage_slide import review_onepage_slide, write_review_report
except ImportError:
    from review_onepage_slide import review_onepage_slide, write_review_report


PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_PPT = PROJECT_ROOT / "output" / "internal_traceability_uat_logic_onepage.pptx"
DEFAULT_PRESET_ID = "process_narrative"

ACCENT = (173, 5, 61)
INK = (35, 41, 46)
MUTED = (98, 109, 121)
LIGHT_TEXT = (145, 152, 160)
LINE = (224, 228, 232)
WHITE = (255, 255, 255)
PURPLE = (118, 86, 167)
TEAL = (53, 116, 136)
CARD_STYLES = (
    {"accent": PURPLE, "soft": (247, 242, 250)},
    {"accent": ACCENT, "soft": (252, 242, 246)},
    {"accent": TEAL, "soft": (241, 247, 249)},
)
BANNED_VISIBLE_PHRASES = (
    "讲解重点",
    "关键用户在UAT前先看懂",
    "先看懂三件事",
    "建议优先验证",
    "UAT介绍",
    "讲解目标",
)


@dataclass(frozen=True)
class NarrativeCard:
    index: str
    chinese: str
    stage: str
    objective: str
    bullets: tuple[str, str, str]
    output_label: str


@dataclass(frozen=True)
class Rect:
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    def overlaps(self, other: "Rect") -> bool:
        return not (
            self.right <= other.left
            or other.right <= self.left
            or self.bottom <= other.top
            or other.bottom <= self.top
        )


TITLE_SEGMENTS = [
    ("内部追溯模块", INK, True),
    ("业务逻辑", ACCENT, True),
]
SUBTITLE = "通过数据驱动文件生成与关联，形成从任务协同到监督预警的内部追溯闭环。"
VALUE_STATEMENT = "追溯请求发起后，系统不是人工翻找文件，而是沿着既定业务链路自动定位数据、匹配文件、推进任务并输出结果。"
FLOW_STATEMENT = "数据找文件  ->  任务带进度  ->  规则做预警"
BOTTOM_LINE = "最终形成：目标文件可快速定位、任务进度可实时联动、报告与异常可标准化输出。"

CARDS = (
    NarrativeCard(
        index="01",
        chinese="数据集成与建模",
        stage="把数据准备成可追溯底座",
        objective="先把外围系统数据转成可统一识别、可生成文件、可关联映射的结构化资产。",
        bullets=(
            "对外围数据做清洗和标准化建模，统一到业务域、业务主题和业务场景。",
            "有明确范围和印戳的文件直接集成；范围需按场景确认的内容即时生成。",
            "通过RPA补齐客观缺失单据，并维护文件清单及数据-文件映射关系。",
        ),
        output_label="输出：可定位的文件池",
    ),
    NarrativeCard(
        index="02",
        chinese="任务协作流程",
        stage="把找文件嵌入任务执行",
        objective="让任务发起、范围锁定、文件匹配、进度同步和报告输出在同一条业务链路上联动。",
        bullets=(
            "任务分为主任务和内部任务，主任务可继续派生内外协同动作。",
            "通过追溯元素和追溯参数锁定数据范围，自动识别对应业务场景。",
            "结合链路配置、TC参数、数据关系表和PPT模板，自动输出标准化报告。",
        ),
        output_label="输出：进度联动的追溯任务",
    ),
    NarrativeCard(
        index="03",
        chinese="监督预警机制",
        stage="把结果沉淀成可闭环规则",
        objective="将追溯过程中的文件、报告和数据继续结构化，为检查规则和预警闭环提供统一底座。",
        bullets=(
            "整合外围系统数据、建模加工数据、追溯文件和报告，统一为结构化数据资产。",
            "通过规则引擎比对结构化数据与业务检查规则，自动识别异常数据。",
            "按时间或条件触发预警，覆盖齐套、合理性和一致性等常见场景。",
        ),
        output_label="输出：可推送的异常闭环",
    ),
)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def set_run_style(run, *, size: float, color: tuple[int, int, int], bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def estimate_lines(text: str, width: int, font_size: float) -> int:
    weighted_chars = 0.0
    for char in text:
        if char.isspace():
            weighted_chars += 0.35
        elif ord(char) < 128:
            weighted_chars += 0.58
        else:
            weighted_chars += 1.0
    chars_per_line = max(1.0, width / max(font_size * 7000, 1))
    return max(1, math.ceil(weighted_chars / chars_per_line))


def add_text(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: float,
    color: tuple[int, int, int] = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=font_size, color=color, bold=bold)
    return shape


def add_shape(
    slide,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width: float = 1,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def keep_only_body_template(prs: Presentation) -> None:
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    keep_idx = manifest.slide_roles.body_template
    for index in range(len(prs.slides) - 1, -1, -1):
        if index != keep_idx:
            remove_slide(prs, index)


def render_title(slide, preset_id: str) -> None:
    preset = get_onepage_layout_preset(preset_id)
    title_shape = slide.shapes[0]
    tf = title_shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    for text, color, bold in TITLE_SEGMENTS:
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size=float(preset.renderer_hints["title_font_size"]), color=color, bold=bold)

    add_text(
        slide,
        720000,
        540000,
        9100000,
        190000,
        SUBTITLE,
        font_size=float(preset.renderer_hints["subtitle_font_size"]),
        color=MUTED,
    )


def render_value_statement(slide, slide_width: int) -> None:
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        720000,
        1130000,
        slide_width - 1440000,
        430000,
        fill=(250, 248, 248),
        line=(232, 222, 225),
        line_width=1,
    )
    add_text(
        slide,
        930000,
        1230000,
        slide_width - 1860000,
        220000,
        VALUE_STATEMENT,
        font_size=13.2,
        color=INK,
        bold=True,
    )
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 810000, 1200000, 65000, 240000, fill=ACCENT)


def render_flow_statement(slide, slide_width: int) -> None:
    add_text(
        slide,
        720000,
        1650000,
        slide_width - 1440000,
        170000,
        FLOW_STATEMENT,
        font_size=13.2,
        color=ACCENT,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def card_geometries(slide_width: int, preset_id: str) -> list[tuple[int, int, int, int]]:
    preset = get_onepage_layout_preset(preset_id)
    left = 720000
    gap = int(preset.renderer_hints["card_gap"])
    usable_width = slide_width - left * 2
    width = int((usable_width - gap * 2) / 3)
    top = 1980000
    height = 2830000
    return [(left + idx * (width + gap), top, width, height) for idx in range(3)]


def render_card(
    slide,
    geom: tuple[int, int, int, int],
    spec: NarrativeCard,
    style: dict[str, tuple[int, int, int]],
    preset_id: str,
) -> None:
    preset = get_onepage_layout_preset(preset_id)
    left, top, width, height = geom
    accent = style["accent"]
    soft = style["soft"]

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=WHITE, line=LINE, line_width=1)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, 70000, fill=accent)

    add_text(slide, left + 180000, top + 150000, 320000, 120000, spec.index, font_size=10.5, color=LIGHT_TEXT, bold=True)
    add_text(
        slide,
        left + 180000,
        top + 315000,
        width - 360000,
        150000,
        spec.chinese,
        font_size=float(preset.renderer_hints["card_title_font_size"]),
        color=INK,
        bold=True,
    )

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 180000, top + 650000, width - 360000, 250000, fill=soft)
    add_text(
        slide,
        left + 230000,
        top + 720000,
        width - 460000,
        120000,
        spec.stage,
        font_size=11.0,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_text(
        slide,
        left + 180000,
        top + 990000,
        width - 360000,
        270000,
        spec.objective,
        font_size=float(preset.renderer_hints["card_definition_font_size"]),
        color=INK,
        bold=True,
    )

    bullet_y = top + 1350000
    for bullet in spec.bullets:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + 190000, bullet_y + 58000, 70000, 70000, fill=accent)
        add_text(
            slide,
            left + 320000,
            bullet_y,
            width - 500000,
            245000,
            bullet,
            font_size=float(preset.renderer_hints["card_bullet_font_size"]),
            color=MUTED,
        )
        bullet_y += 395000

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 180000, top + height - 360000, width - 360000, 190000, fill=soft)
    add_text(
        slide,
        left + 230000,
        top + height - 305000,
        width - 460000,
        100000,
        spec.output_label,
        font_size=9.8,
        color=accent,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def render_connectors(slide, slide_width: int, preset_id: str) -> None:
    cards = card_geometries(slide_width, preset_id)
    top = cards[0][1] + 1180000
    for idx, geom in enumerate(cards[:-1]):
        left, _, width, _ = geom
        next_left = cards[idx + 1][0]
        gap = next_left - (left + width)
        connector_left = left + width + int(gap * 0.25)
        connector_width = int(gap * 0.5)
        add_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.CHEVRON,
            connector_left,
            top,
            connector_width,
            170000,
            fill=(246, 247, 249),
            line=(226, 229, 233),
            line_width=1,
        )


def render_footer(slide, slide_width: int, preset_id: str) -> None:
    preset = get_onepage_layout_preset(preset_id)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 5160000, slide_width - 1440000, 1, fill=LINE)
    add_text(
        slide,
        720000,
        5290000,
        slide_width - 1440000,
        160000,
        BOTTOM_LINE,
        font_size=float(preset.renderer_hints["footer_font_size"]),
        color=MUTED,
    )


def self_check_layout(prs: Presentation, preset_id: str = DEFAULT_PRESET_ID) -> list[str]:
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    preset = get_onepage_layout_preset(preset_id)
    issues: list[str] = []

    if any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes):
        issues.append("slide contains picture shapes; expected editable-only output")

    for idx, shape in enumerate(slide.shapes, 1):
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_width or shape.top + shape.height > slide_height:
            issues.append(f"shape {idx} exceeds slide bounds")

    cards = [Rect(*geom) for geom in card_geometries(slide_width, preset_id)]
    for i, current in enumerate(cards):
        for j, other in enumerate(cards[i + 1 :], start=i + 1):
            if current.overlaps(other):
                issues.append(f"card {i + 1} overlaps card {j + 1}")

    value_box = Rect(720000, 1130000, slide_width - 1440000, 430000)
    for idx, card in enumerate(cards, 1):
        if value_box.overlaps(card):
            issues.append(f"value statement overlaps card {idx}")

    footer = Rect(720000, 5160000, slide_width - 1440000, 280000)
    for idx, card in enumerate(cards, 1):
        if footer.overlaps(card):
            issues.append(f"footer overlaps card {idx}")

    title_text = "".join(segment[0] for segment in TITLE_SEGMENTS)
    if estimate_lines(title_text, 9200000, float(preset.renderer_hints["title_font_size"])) > 1:
        issues.append("title is likely to wrap to more than one line")

    if estimate_lines(VALUE_STATEMENT, slide_width - 1860000, 13.2) > 2:
        issues.append("value statement is too dense")

    if estimate_lines(FLOW_STATEMENT, slide_width - 1440000, 13.2) > 1:
        issues.append("flow statement is too dense")

    for spec in CARDS:
        if estimate_lines(spec.objective, 2500000, float(preset.renderer_hints["card_definition_font_size"])) > 2:
            issues.append(f"{spec.chinese} objective is too dense")
        for bullet in spec.bullets:
            if estimate_lines(bullet, 2480000, float(preset.renderer_hints["card_bullet_font_size"])) > 2:
                issues.append(f"{spec.chinese} bullet is too dense")

    slide_text = "\n".join(shape.text_frame.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
    for phrase in BANNED_VISIBLE_PHRASES:
        if phrase in slide_text:
            issues.append(f"slide still contains banned visible phrase: {phrase}")

    return issues


def build_slide(output_path: Path, preset_id: str = DEFAULT_PRESET_ID, *, export_review: bool = True) -> Path:
    prs = Presentation(str(DEFAULT_TEMPLATE))
    keep_only_body_template(prs)
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)

    render_title(slide, preset_id)
    render_value_statement(slide, slide_width)
    render_flow_statement(slide, slide_width)
    render_connectors(slide, slide_width, preset_id)
    for geom, spec, style in zip(card_geometries(slide_width, preset_id), CARDS, CARD_STYLES):
        render_card(slide, geom, spec, style, preset_id)
    render_footer(slide, slide_width, preset_id)

    issues = self_check_layout(prs, preset_id=preset_id)
    if issues:
        raise ValueError("layout self-check failed: " + " | ".join(issues))

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
        saved_path = output_path
    except PermissionError:
        timestamp = datetime.datetime.now().strftime("%H%M%S_%f")[:9]
        saved_path = output_path.with_stem(f"{output_path.stem}_{timestamp}")
        prs.save(str(saved_path))

    if export_review:
        review_report = review_onepage_slide(
            saved_path,
            expected_card_count=3,
            banned_phrases=BANNED_VISIBLE_PHRASES,
            export_previews=True,
            preview_dir=saved_path.parent / f"{saved_path.stem}_previews",
        )
        review_path = saved_path.with_suffix(".review.json")
        write_review_report(review_report, review_path)
        blocker_messages = [finding.message for finding in review_report.findings if finding.level == "blocker"]
        if blocker_messages:
            raise ValueError("post-generation review failed: " + " | ".join(blocker_messages))

    return saved_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a one-page PPT that explains the internal traceability business logic for UAT key users.")
    parser.add_argument("--output", default=str(OUTPUT_PPT), help="Output PPTX path.")
    parser.add_argument("--layout-preset", default=DEFAULT_PRESET_ID, help="One-page layout preset id.")
    parser.add_argument("--skip-review", action="store_true", help="Skip post-generation heuristic review export.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    get_onepage_layout_preset(args.layout_preset)
    output = build_slide(Path(args.output).resolve(), preset_id=args.layout_preset, export_review=not args.skip_review)
    print(output)


if __name__ == "__main__":
    main()
