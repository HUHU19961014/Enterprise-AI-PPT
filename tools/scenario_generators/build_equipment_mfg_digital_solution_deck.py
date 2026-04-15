from __future__ import annotations

import argparse
import datetime
from pathlib import Path
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

try:
    from sie_autoppt.body_renderers import fill_body_slide, pick_text_shapes
    from sie_autoppt.config import COLOR_ACTIVE, COLOR_INACTIVE, DEFAULT_TEMPLATE
    from sie_autoppt.models import BodyPageSpec
    from sie_autoppt.slide_ops import ensure_last_slide, remove_slide
    from sie_autoppt.template_manifest import load_template_manifest
    from sie_autoppt.text_ops import add_textbox, write_text
except ModuleNotFoundError:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parents[1]))
    from sie_autoppt.body_renderers import fill_body_slide, pick_text_shapes
    from sie_autoppt.config import COLOR_ACTIVE, COLOR_INACTIVE, DEFAULT_TEMPLATE
    from sie_autoppt.models import BodyPageSpec
    from sie_autoppt.slide_ops import ensure_last_slide, remove_slide
    from sie_autoppt.template_manifest import load_template_manifest
    from sie_autoppt.text_ops import add_textbox, write_text


PROJECT_ROOT = Path(__file__).resolve().parents[2]
OUTPUT_PPT = PROJECT_ROOT / "output" / "装备制造行业数字化解决方案_SIE整套版.pptx"

ACCENT = COLOR_ACTIVE
INK = (35, 41, 46)
MUTED = (98, 109, 121)
LIGHT = COLOR_INACTIVE


def build_body_pages() -> list[BodyPageSpec]:
    return [
        BodyPageSpec(
            page_key="pain_points",
            title="行业核心痛点",
            subtitle="真正制约装备制造企业的，不是单个环节效率，而是跨部门、跨系统、跨供应链的复杂协同。",
            bullets=[],
            pattern_id="pain_cards",
            nav_title="核心痛点",
            payload={
                "lead": "六类问题可归并为三大经营痛点：研产协同复杂、计划供应链承压、制造服务不闭环。",
                "bottom_banner": "先澄清问题边界，再推进数字化建设，避免“多系统并存但协同仍失灵”。",
                "cards": [
                    {
                        "title": "研产协同复杂",
                        "detail": "产品结构、零部件与工艺链条长，设计变更会快速传导到采购和生产。",
                        "points": [
                            "边设计、边采购、边生产并行发生",
                            "设计变更频繁，影响面广",
                            "研发与制造协同成本高",
                        ],
                    },
                    {
                        "title": "计划供应链承压",
                        "detail": "计划动态变化快，物料与供应商协同难，库存与齐套性矛盾突出。",
                        "points": [
                            "生产计划频繁重排",
                            "物料需求难精准协同",
                            "询价、比价、核价与供应商质量管理复杂",
                        ],
                    },
                    {
                        "title": "制造服务不闭环",
                        "detail": "车间执行、质量闭环和售后数据分散，柔性生产与服务响应能力不足。",
                        "points": [
                            "设备与系统数据孤岛明显",
                            "质量问题响应慢、闭环弱",
                            "售后故障响应慢、维护成本高",
                        ],
                    },
                ],
            },
        ),
        BodyPageSpec(
            page_key="solution_blueprint",
            title="解决方案应用蓝图",
            subtitle="核心目标是通过数据采集和集成打破数据壁垒，实现数据共享、业务协同与经营可视。",
            bullets=[],
            pattern_id="solution_architecture",
            nav_title="应用蓝图",
            payload={
                "banner_text": "APPLICATION BLUEPRINT",
                "layers": [
                    {
                        "label": "L1",
                        "title": "数据采集与集成底座",
                        "detail": "统一连接设计、计划、供应链、制造、设备与售后数据，形成实时共享的数据底座。",
                    },
                    {
                        "label": "L2",
                        "title": "供应链协同",
                        "detail": "围绕采购、询比核价、齐套、供应商质量与交付协同，提升主机厂与配套体系联动能力。",
                    },
                    {
                        "label": "L3",
                        "title": "制造与质量闭环",
                        "detail": "打通生产准备、投产、投料、完工、入库与质量响应，实现按质按量按时生产。",
                    },
                    {
                        "label": "L4",
                        "title": "财务与运营管控",
                        "detail": "形成财务管控与经营战情中心，实现实时洞察、风险预警和数据驱动决策。",
                    },
                ],
            },
        ),
        BodyPageSpec(
            page_key="path_group_business",
            title="转型路径一：业务主链能力",
            subtitle="第一组路径聚焦设计、订单、项目与计划，目标是先打通业务主链的协同基础。",
            bullets=[],
            pattern_id="capability_ring",
            nav_title="业务主链",
            payload={
                "headline": "四类能力决定了装备制造企业能否把前端需求和后端交付真正串起来。",
                "items": [
                    {"title": "设计制造协同", "detail": "让设计、变更和工艺信息同步到生产、供应链与财务，减少并行协同损耗。"},
                    {"title": "MTO 按单制造", "detail": "用特征值和选配规则支撑个性化配置与快速报价，提高非标订单响应能力。"},
                    {"title": "ETO 按项目制造", "detail": "按项目统筹人财物和进度、成本、质量、风险，避免项目型交付失控。"},
                    {"title": "多层级计划体系", "detail": "打通销售计划、主生产计划和物料需求计划，形成产供销协同节奏。"},
                ],
            },
        ),
        BodyPageSpec(
            page_key="path_group_operations",
            title="转型路径二：运营执行能力",
            subtitle="第二组路径聚焦制造执行、质量闭环、服务运维和经营管控，目标是把执行结果真正闭环起来。",
            bullets=[],
            pattern_id="capability_ring",
            nav_title="运营执行",
            payload={
                "headline": "四类能力决定了企业能否把生产现场、质量响应、售后服务和经营决策连成闭环。",
                "items": [
                    {"title": "制造执行与用料", "detail": "覆盖生产准备到入库全流程，并支持跨部门异常联动，保障按质按量按时生产。"},
                    {"title": "全生命周期质量", "detail": "从产品开发到售后服务建立质量闭环，实现质量可视、可控、可决策。"},
                    {"title": "产品服务与运维", "detail": "围绕预防性维修、预见性维修和服务闭环，提升响应效率和客户体验。"},
                    {"title": "数字化运营管控", "detail": "用预测和模型洞察经营不确定性与风险，支持实时监控和经营决策。"},
                ],
            },
        ),
        BodyPageSpec(
            page_key="implementation_roadmap",
            title="实施方向与分阶段路线",
            subtitle="建议按“业务数字化打底、工厂能力升级、服务化延伸、柔性智能深化”的路径推进。",
            bullets=[],
            pattern_id="roadmap_timeline",
            nav_title="实施路线",
            payload={
                "headline": "四大方向可作为装备制造数字化转型的主实施序列。",
                "footer": "先打通核心业务，再推进 IT/OT 融合，随后延伸服务化，最终走向智能柔性制造。",
                "stages": [
                    {
                        "period": "阶段一",
                        "title": "核心业务数字化",
                        "detail": "贯穿销售、计划、供应链、生产、安装调试到售后全环节，形成经营管理战情中心。",
                    },
                    {
                        "period": "阶段二",
                        "title": "IT/OT 融合工厂",
                        "detail": "以 AI、大数据、云计算、5G 与物联网支撑钣金、CNC、装配等典型车间方案。",
                    },
                    {
                        "period": "阶段三",
                        "title": "服务化转型",
                        "detail": "通过增值服务摆脱低价竞争，形成新的利润增长点与客户粘性。",
                    },
                    {
                        "period": "阶段四",
                        "title": "智能柔性制造",
                        "detail": "打通智能设备与关键部件，推动装备智能化、产线柔性化和软硬件融合。",
                    },
                ],
            },
        ),
    ]


AGENDA_LINES = ["行业判断", "核心痛点", "解决蓝图", "业务主链", "运营执行", "实施路线"]


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def add_panel(slide, left: int, top: int, width: int, height: int, *, fill: tuple[int, int, int], line: tuple[int, int, int]) -> None:
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(fill)
    panel.line.color.rgb = rgb(line)


def remove_all_shapes(slide) -> None:
    for shape in list(slide.shapes):
        element = shape._element
        element.getparent().remove(element)


def remove_text_shapes(slide) -> None:
    removable = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
    for shape in removable:
        element = shape._element
        element.getparent().remove(element)


def dedupe_package_entries(pptx_path: Path) -> None:
    with zipfile.ZipFile(pptx_path, "r") as source_zip:
        latest_entries: dict[str, bytes] = {}
        ordered_names: list[str] = []
        for info in source_zip.infolist():
            if info.filename in latest_entries:
                ordered_names.remove(info.filename)
            latest_entries[info.filename] = source_zip.read(info)
            ordered_names.append(info.filename)

    rebuilt_path = pptx_path.with_name(f"{pptx_path.stem}_deduped{pptx_path.suffix}")
    with zipfile.ZipFile(rebuilt_path, "w", zipfile.ZIP_DEFLATED) as rebuilt_zip:
        for name in ordered_names:
            rebuilt_zip.writestr(name, latest_entries[name])
    rebuilt_path.replace(pptx_path)


def prepare_presentation() -> tuple[Presentation, object]:
    prs = Presentation(str(DEFAULT_TEMPLATE))
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    closing_index = manifest.slide_pools.ending if manifest.slide_pools and manifest.slide_pools.ending is not None else len(prs.slides) - 1
    closing_slide_id = prs.slides[closing_index].slide_id
    keep_indices = {
        manifest.slide_roles.welcome,
        manifest.slide_roles.theme,
        manifest.slide_roles.directory,
        manifest.slide_roles.body_template,
        closing_index,
    }
    for slide_index in range(len(prs.slides) - 1, -1, -1):
        if slide_index not in keep_indices:
            remove_slide(prs, slide_index)
    body_source = prs.slides[3]
    for _ in range(4):
        prs.slides.add_slide(body_source.slide_layout)
    ensure_last_slide(prs, closing_slide_id)
    return prs, manifest


def apply_cover(slide) -> None:
    add_textbox(
        slide,
        900000,
        1500000,
        8500000,
        520000,
        "装备制造行业数字化解决方案",
        size_pt=30,
        bold=True,
        color=ACCENT,
    )
    add_textbox(
        slide,
        900000,
        2120000,
        8200000,
        380000,
        "Digital Transformation Solution for Equipment Manufacturing",
        size_pt=14,
        color=INK,
    )
    add_textbox(
        slide,
        900000,
        2620000,
        8200000,
        800000,
        "通过数据驱动，打通从设计到售后的全链路，支撑研发协同化、生产精益化、管控实时化、服务智能化。",
        size_pt=15,
        color=MUTED,
    )
    add_textbox(
        slide,
        900000,
        5620000,
        4200000,
        220000,
        "SiE赛意 | 行业解决方案示意稿",
        size_pt=12,
        color=INK,
    )
    add_textbox(
        slide,
        9500000,
        5620000,
        1700000,
        220000,
        datetime.date.today().strftime("%Y/%m/%d"),
        size_pt=12,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    add_panel(slide, 900000, 3300000, 2500000, 90000, fill=ACCENT, line=ACCENT)


def apply_industry_judgement(slide) -> None:
    texts = sorted(pick_text_shapes(slide), key=lambda shape: (shape.top, shape.left))
    if len(texts) >= 3:
        title_shape = max(texts, key=lambda shape: shape.width)
        footer_shapes = sorted([shape for shape in texts if shape is not title_shape], key=lambda shape: shape.top)
        write_text(title_shape, "行业判断与转型命题", size_pt=28, bold=True, color=ACCENT, preserve_runs=True)
        write_text(footer_shapes[0], "INDUSTRY VIEW", size_pt=14, color=INK, preserve_runs=True)
        write_text(footer_shapes[1], "SiE赛意", size_pt=12, color=MUTED, preserve_runs=True)

    add_textbox(
        slide,
        600000,
        2550000,
        4800000,
        520000,
        "装备制造业的转型重点，不是局部提效，而是以数据驱动打通设计到售后的全链路协同。",
        size_pt=17,
        bold=True,
        color=INK,
    )
    card_specs = [
        ("行业定位", "多品种、小批量、非标定制是装备制造的基本业务特征。"),
        ("经营矛盾", "交期、质量、功能、个性化与智能化要求叠加，复杂度持续上升。"),
        ("转型命题", "通过数据驱动实现研发协同化、生产精益化、管控实时化、服务智能化。"),
    ]
    left = 600000
    top = 3470000
    width = 1450000
    gap = 180000
    for index, (title, detail) in enumerate(card_specs):
        box_left = left + index * (width + gap)
        add_panel(slide, box_left, top, width, 1350000, fill=(247, 249, 252), line=(214, 221, 229))
        add_textbox(slide, box_left + 120000, top + 120000, width - 240000, 180000, title, size_pt=15, bold=True, color=ACCENT)
        add_textbox(slide, box_left + 120000, top + 390000, width - 240000, 720000, detail, size_pt=11, color=MUTED)


def apply_agenda(slide, chapter_lines: list[str]) -> None:
    remove_all_shapes(slide)
    add_textbox(
        slide,
        630000,
        610000,
        1600000,
        500000,
        "目录",
        size_pt=28,
        bold=True,
        color=INK,
    )
    add_textbox(
        slide,
        720000,
        1280000,
        2000000,
        260000,
        "Agenda",
        size_pt=13,
        color=ACCENT,
    )
    add_panel(slide, 980000, 1850000, 9800000, 3200000, fill=(255, 255, 255), line=(220, 225, 231))
    add_textbox(
        slide,
        1100000,
        2050000,
        4000000,
        300000,
        "建议汇报结构",
        size_pt=20,
        bold=True,
        color=ACCENT,
    )
    add_textbox(
        slide,
        1100000,
        2420000,
        6500000,
        260000,
        "围绕“为什么转、转什么、怎么落地”三层问题组织整套方案。",
        size_pt=11,
        color=MUTED,
    )
    card_width = 2900000
    card_height = 760000
    gap_x = 220000
    gap_y = 180000
    origin_left = 1180000
    origin_top = 2850000
    for idx, line in enumerate(chapter_lines):
        row, col = divmod(idx, 3)
        left = origin_left + col * (card_width + gap_x)
        top = origin_top + row * (card_height + gap_y)
        add_panel(slide, left, top, card_width, card_height, fill=(247, 249, 252), line=(214, 221, 229))
        add_textbox(
            slide,
            left + 140000,
            top + 120000,
            500000,
            180000,
            f"{idx + 1:02d}",
            size_pt=15,
            bold=True,
            color=ACCENT,
        )
        add_textbox(
            slide,
            left + 140000,
            top + 320000,
            card_width - 280000,
            220000,
            line,
            size_pt=16,
            bold=True,
            color=INK,
        )
        add_textbox(
            slide,
            left + 140000,
            top + 560000,
            card_width - 280000,
            150000,
            "\u63d0\u6848\u7ae0\u8282",
            size_pt=9,
            color=LIGHT,
        )


def apply_closing(slide) -> None:
    add_textbox(
        slide,
        950000,
        1650000,
        9200000,
        520000,
        "结论",
        size_pt=28,
        bold=True,
        color=ACCENT,
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        slide,
        950000,
        2400000,
        9200000,
        920000,
        "装备制造业数字化转型的本质，不是再堆系统，而是以数据驱动重构从设计到售后的经营协同链路。",
        size_pt=22,
        bold=True,
        color=INK,
    )
    add_textbox(
        slide,
        950000,
        3600000,
        9200000,
        520000,
        "关键词：研发协同化 / 生产精益化 / 管控实时化 / 服务智能化",
        size_pt=15,
        color=MUTED,
    )
    add_textbox(
        slide,
        950000,
        4380000,
        9200000,
        280000,
        "SiE 方案建议先以核心业务数字化为抓手，再逐步延伸到数字工厂、服务化与柔性智能制造。",
        size_pt=12,
        color=LIGHT,
    )


def build_deck(output_path: Path) -> Path:
    prs, manifest = prepare_presentation()
    body_pages = build_body_pages()

    apply_cover(prs.slides[0])
    apply_industry_judgement(prs.slides[1])
    apply_agenda(prs.slides[2], AGENDA_LINES)

    for index, page in enumerate(body_pages, start=3):
        fill_body_slide(prs.slides[index], page, manifest)

    apply_closing(prs.slides[-1])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    dedupe_package_entries(output_path)
    return output_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a full SIE-style deck for equipment manufacturing digital transformation.")
    parser.add_argument("--output", default=str(OUTPUT_PPT), help="Output PPTX path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = build_deck(Path(args.output).resolve())
    print(output)


if __name__ == "__main__":
    main()
