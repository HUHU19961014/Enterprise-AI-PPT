from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

PROJECT_ROOT = Path(__file__).resolve().parents[2]
project_root_str = str(PROJECT_ROOT)
if project_root_str not in sys.path:
    sys.path.insert(0, project_root_str)

from tools.sie_autoppt.slide_ops import set_slide_metadata_names


SAMPLES_INPUT_DIR = PROJECT_ROOT / "samples" / "input"
REFERENCE_BODY_PATH = SAMPLES_INPUT_DIR / "reference_body_style.pptx"


UAT_PLAN_SAMPLE_HTML = """\
<div class="title">制造企业数字化升级蓝图</div>
<div class="subtitle">围绕现状、场景与治理动作梳理项目重点</div>
<div class="scope-title">关键业务场景</div>
<div class="scope-subtitle">聚焦供应链协同与交付风险</div>
<div class="focus-title">实施治理重点</div>
<div class="focus-subtitle">需要跨部门协同推进</div>
<div class="phase-time">Q1</div><div class="phase-name">规划</div><div class="phase-func">完成现状诊断与蓝图设计</div>
<div class="phase-time">Q2</div><div class="phase-name">建设</div><div class="phase-func">打通主数据与核心流程</div>
<div class="phase-time">Q3</div><div class="phase-name">试点</div><div class="phase-func">验证跨部门协同链路</div>
<div class="scenario">采购到库存链路口径不一致</div>
<div class="scenario">生产计划与执行之间存在断点</div>
<div class="scenario">质量追溯依赖人工汇总</div>
<div class="note">需要建立项目治理机制与节奏</div>
<div class="note">主数据口径必须同步收敛</div>
<div class="footer">同步风险与验收标准，避免试点和推广脱节</div>
"""


ARCHITECTURE_PROGRAM_SAMPLE_HTML = """\
<div class="title">ERP 架构蓝图规划</div>
<div class="subtitle">围绕平台架构、流程协同与治理机制制定分阶段方案</div>
<div class="scope-title">流程协同路径</div>
<div class="scope-subtitle">打通采购、生产、仓储与财务流程</div>
<div class="focus-title">治理要求与实施机制</div>
<div class="focus-subtitle">明确跨部门责任与项目验收标准</div>
<div class="phase-time">Q1</div><div class="phase-name">架构设计</div><div class="phase-func">梳理现状系统与目标蓝图</div>
<div class="phase-time">Q2</div><div class="phase-name">集成建设</div><div class="phase-func">完成主数据和核心接口建设</div>
<div class="phase-time">Q3</div><div class="phase-name">业务试点</div><div class="phase-func">验证关键流程闭环</div>
<div class="phase-time">Q4</div><div class="phase-name">推广治理</div><div class="phase-func">复制到更多工厂与业务单元</div>
<div class="scenario">采购到生产的计划协同需要统一流程</div>
<div class="scenario">库存与财务结算链路需要标准化</div>
<div class="scenario">跨系统接口需要形成稳定传输路径</div>
<div class="scenario">异常处理要纳入统一运营机制</div>
<div class="note">建立项目例会与风险升级机制</div>
<div class="note">明确架构评审、上线评审和验收口径</div>
<div class="footer">治理节奏必须和流程建设同步推进</div>
"""


PCB_ERP_GENERAL_SOLUTION_HTML = """\
<div class="title">PCB ERP 总体架构蓝图</div>
<div class="subtitle">面向集团化制造场景构建统一平台能力</div>
<div class="scope-title">关键流程推进路径</div>
<div class="scope-subtitle">以采购、计划、生产、财务四条主链路为核心</div>
<div class="focus-title">治理要求与实施重点</div>
<div class="focus-subtitle">建立跨部门协同机制与验收闭环</div>
<div class="phase-time">Q1</div><div class="phase-name">业务域设计</div><div class="phase-func">梳理销售、采购、生产与财务域边界</div>
<div class="phase-time">Q2</div><div class="phase-name">数据域建设</div><div class="phase-func">统一主数据模型与编码规则</div>
<div class="phase-time">Q3</div><div class="phase-name">集成域打通</div><div class="phase-func">建设 ERP 与 MES/WMS 接口链路</div>
<div class="phase-time">Q4</div><div class="phase-name">运营域闭环</div><div class="phase-func">形成监控预警与持续优化机制</div>
<div class="scenario">采购申请到下单流程需要统一审批路径</div>
<div class="scenario">生产计划到执行反馈需要实时协同机制</div>
<div class="scenario">库存、成本与财务数据需要自动对齐</div>
<div class="scenario">质量异常到纠正动作需要形成闭环</div>
<div class="note">建立项目 steering committee 与周节奏例会</div>
<div class="note">主数据治理必须纳入上线前强校验</div>
<div class="note">关键接口需要设定明确 owner 与 SLA</div>
<div class="note">试点验收要绑定业务指标与推广条件</div>
<div class="footer">所有治理动作需要贯穿方案、试点和推广阶段</div>
"""


REFERENCE_SLIDE_TEXT = {
    5: ("comparison_upgrade_reference", "价值：业务效能跃升\n传统人工追函模式"),
    6: ("pain_cards_reference", "痛点：标准不清晰\n链路难贯通\n组织协同慢"),
    16: ("capability_ring_reference", "赛意追溯产品亮点\nAI智能识别"),
    20: ("five_phase_path_reference", "追溯管理-外部追溯推进路径\n阶段二 数据应用"),
}


def _write_html_fixture(path: Path, content: str) -> None:
    path.write_text(content, encoding="utf-8")


def _build_reference_body_ppt(path: Path) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    while len(prs.slides) < 20:
        prs.slides.add_slide(blank_layout)

    for index, slide in enumerate(prs.slides, start=1):
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8.5), Inches(1.8))
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        marker_text = REFERENCE_SLIDE_TEXT.get(index, ("", f"Reference body slide {index}"))[1]
        for line_index, line in enumerate(marker_text.splitlines()):
            paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
            paragraph.text = line

    prs.save(path)
    set_slide_metadata_names(path, {slide_no: name for slide_no, (name, _) in REFERENCE_SLIDE_TEXT.items()})


def main() -> None:
    SAMPLES_INPUT_DIR.mkdir(parents=True, exist_ok=True)
    _write_html_fixture(SAMPLES_INPUT_DIR / "uat_plan_sample.html", UAT_PLAN_SAMPLE_HTML)
    _write_html_fixture(SAMPLES_INPUT_DIR / "architecture_program_sample.html", ARCHITECTURE_PROGRAM_SAMPLE_HTML)
    _write_html_fixture(SAMPLES_INPUT_DIR / "pcb_erp_general_solution.html", PCB_ERP_GENERAL_SOLUTION_HTML)
    _build_reference_body_ppt(REFERENCE_BODY_PATH)


if __name__ == "__main__":
    main()
