from __future__ import annotations

import argparse
import datetime
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

try:
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest
except ModuleNotFoundError:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_FLOW = PROJECT_ROOT / "output" / "sie_battery_passport_update_rule_flow.pptx"
OUTPUT_DECISION = PROJECT_ROOT / "output" / "sie_battery_passport_update_rule_decision.pptx"

ACCENT = (173, 5, 61)
INK = (35, 41, 46)
MUTED = (96, 108, 121)
LIGHT_TEXT = (146, 154, 161)
LINE = (224, 228, 232)
WHITE = (255, 255, 255)
PANEL = (248, 249, 251)
GREEN = (9, 177, 85)
GREEN_SOFT = (236, 248, 241)
TEAL = (72, 169, 152)
TEAL_SOFT = (235, 247, 245)
LIME = (107, 183, 72)
LIME_SOFT = (240, 249, 236)
GRAY_BOX = (146, 146, 146)
GRAY_SOFT = (244, 244, 244)
ACCENT_SOFT = (252, 241, 245)

TITLE = "Pack 护照更新 / 注销规则"
SUBTITLE = "围绕 Pack 故障后的维修、回收与再利用路径，统一定义护照更新与注销动作。"
LEAD_LABEL = "规则说明"
LEAD_LINE = "Pack 状态变化后，应先完成电池检测，再按修复、回收、用途变化和报废路径执行护照更新或注销。"
SUPPORT = "更新或注销并不取决于是否维修，而取决于状态变化后能否完成参数验证、用途判定与合规闭环。"
FOOTNOTE = "要求：1. 检测电池；2. 更新护照参数（容量、自放电率、功率、往返效能、内阻）。若海外现场不具备检测条件，合规建议直接报废。"
FAITHFUL_FORBIDDEN_TEXTS = ("核心判断", "结论摘要", "第一梯队", "优先推进", "流程版", "决策版")


@dataclass(frozen=True)
class DefinitionCard:
    key: str
    title: str
    detail: str
    accent: tuple[int, int, int]
    soft: tuple[int, int, int]


DEFINITIONS = (
    DefinitionCard(
        key="re-used",
        title="同用途复用",
        detail="简单检测或维修后继续按原用途使用，常见于备件替换或返修复用。",
        accent=TEAL,
        soft=TEAL_SOFT,
    ),
    DefinitionCard(
        key="remanufactured",
        title="原用途再制造",
        detail="经拆解、换芯或模块更换后恢复性能；容量恢复至原额定的 >=90%，单芯健康差异 <=3%。",
        accent=ACCENT,
        soft=ACCENT_SOFT,
    ),
    DefinitionCard(
        key="repurposed",
        title="变更用途再利用",
        detail="经处理后转作不同用途，例如 EV 电池回收后转用于储能。",
        accent=LIME,
        soft=LIME_SOFT,
    ),
)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def set_run_style(run, *, size: float, color: tuple[int, int, int], bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: float,
    color: tuple[int, int, int] = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=font_size, color=color, bold=bold)
    return shape


def add_bullets(
    slide,
    left: int,
    top: int,
    width: int,
    bullets: list[str] | tuple[str, ...],
    *,
    bullet_color: tuple[int, int, int],
    text_color: tuple[int, int, int] = MUTED,
    font_size: float = 10.8,
    row_gap: int = 250000,
    bullet_size: int = 68000,
):
    current_y = top
    for bullet in bullets:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left, current_y + 58000, bullet_size, bullet_size, fill=bullet_color)
        add_text(slide, left + 110000, current_y, width - 110000, 220000, bullet, font_size=font_size, color=text_color)
        current_y += row_gap


def add_shape(
    slide,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width: float = 1,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1: int, y1: int, x2: int, y2: int, *, color: tuple[int, int, int], width_pt: float = 1.1):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width_pt)
    return line


def add_arrow(slide, x1: int, y1: int, x2: int, y2: int, *, color: tuple[int, int, int], width_pt: float = 1.2) -> None:
    add_line(slide, x1, y1, x2, y2, color=color, width_pt=width_pt)
    if abs(y2 - y1) <= 30000:
        if x2 >= x1:
            tri = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x2 - 70000, y2 - 42000, 84000, 84000, fill=color)
            tri.rotation = 90
        else:
            tri = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x2 - 14000, y2 - 42000, 84000, 84000, fill=color)
            tri.rotation = 270
    elif abs(x2 - x1) <= 30000 and y2 >= y1:
        tri = add_shape(slide, MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x2 - 42000, y2 - 70000, 84000, 84000, fill=color)
        tri.rotation = 180


def keep_only_body_template(prs: Presentation) -> None:
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    keep_idx = manifest.slide_roles.body_template
    for index in range(len(prs.slides) - 1, -1, -1):
        if index != keep_idx:
            remove_slide(prs, index)


def render_title(slide) -> None:
    title_shape = slide.shapes[0]
    tf = title_shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    for text, color, bold in [
        ("Pack ", INK, True),
        ("护照更新 / 注销", ACCENT, True),
        ("规则", INK, True),
    ]:
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size=24.5, color=color, bold=bold)
    add_text(slide, 720000, 540000, 9000000, 190000, SUBTITLE, font_size=12.0, color=MUTED)


def render_lead(slide, slide_width: int) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 1090000, 70000, 270000, fill=ACCENT)
    add_text(slide, 850000, 1085000, 1800000, 180000, LEAD_LABEL, font_size=13.0, color=ACCENT, bold=True)
    add_text(slide, 720000, 1345000, slide_width - 1440000, 250000, LEAD_LINE, font_size=17.2, color=INK, bold=True)
    add_text(slide, 720000, 1605000, slide_width - 1440000, 150000, SUPPORT, font_size=11.2, color=MUTED)


def render_definition_ribbon(slide) -> None:
    left = 720000
    top = 1920000
    width = 3440000
    gap = 160000
    height = 630000
    for idx, card in enumerate(DEFINITIONS):
        x = left + idx * (width + gap)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, width, height, fill=WHITE, line=LINE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, top, width, 50000, fill=card.accent)
        add_text(slide, x + 170000, top + 115000, 900000, 120000, card.key, font_size=9.2, color=card.accent, bold=True)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x + 170000, top + 300000, 1080000, 200000, fill=card.soft)
        add_text(slide, x + 240000, top + 338000, 950000, 120000, card.title, font_size=11.8, color=card.accent, bold=True)
        add_text(slide, x + 1370000, top + 120000, 1840000, 340000, card.detail, font_size=9.6, color=MUTED)


def status_box(slide, left: int, top: int, width: int, height: int, text: str, *, fill: tuple[int, int, int], color: tuple[int, int, int] = WHITE, size: float = 14.4) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=fill)
    add_text(slide, left, top + 40000, width, height - 80000, text, font_size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)


def render_flow_canvas(slide, *, left: int, top: int, width: int, height: int, compact: bool = False) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=PANEL, line=LINE)
    add_text(slide, left + 180000, top + 120000, 1200000, 120000, "典型路径", font_size=12.0, color=INK, bold=True)
    add_text(slide, left + width - 1700000, top + 120000, 1500000, 120000, "从故障处置到护照动作", font_size=10.2, color=LIGHT_TEXT, align=PP_ALIGN.RIGHT)

    if compact:
        y_bottom = top + 1180000
        y_mid = top + 760000
        y_top = top + 340000
        status_box(slide, left + 160000, y_bottom, 980000, 360000, "Pack\n正常", fill=GREEN, size=13.0)
        status_box(slide, left + 1320000, y_bottom, 980000, 360000, "Pack\n故障", fill=GREEN, size=13.0)
        status_box(slide, left + 2480000, y_bottom, 1080000, 360000, "Pack\n恢复正常", fill=GREEN, size=13.0)
        status_box(slide, left + 4000000, y_mid, 1180000, 360000, "Pack\n作为备件使用", fill=TEAL, size=12.6)
        status_box(slide, left + 5590000, y_top, 1120000, 360000, "Pack\n回收商回收", fill=GRAY_SOFT, color=INK, size=12.8)
        status_box(slide, left + 5460000, y_bottom, 1200000, 360000, "Pack\n不更换用途", fill=LIME, size=12.6)
        status_box(slide, left + 6880000, y_mid, 1200000, 360000, "Pack\n更换用途", fill=LIME, size=12.6)
        status_box(slide, left + 7360000, y_top, 980000, 360000, "Pack\n报废", fill=GRAY_BOX, size=13.0)

        add_arrow(slide, left + 1140000, y_bottom + 180000, left + 1320000, y_bottom + 180000, color=MUTED, width_pt=1.3)
        add_arrow(slide, left + 2300000, y_bottom + 180000, left + 2480000, y_bottom + 180000, color=MUTED, width_pt=1.3)
        add_text(slide, left + 2380000, y_bottom - 120000, 460000, 100000, "现场维修", font_size=10.2, color=INK)

        add_line(slide, left + 2300000, y_bottom + 180000, left + 2300000, y_mid + 180000, color=MUTED, width_pt=1.2)
        add_arrow(slide, left + 2300000, y_mid + 180000, left + 4000000, y_mid + 180000, color=MUTED, width_pt=1.3)
        add_text(slide, left + 3040000, y_mid - 100000, 520000, 100000, "返厂维修", font_size=10.2, color=INK)

        add_line(slide, left + 2300000, y_bottom + 180000, left + 2300000, y_top + 180000, color=MUTED, width_pt=1.2)
        add_arrow(slide, left + 2300000, y_top + 180000, left + 5590000, y_top + 180000, color=MUTED, width_pt=1.3)
        add_text(slide, left + 3680000, y_top - 100000, 540000, 100000, "无法维修", font_size=10.2, color=INK)

        add_line(slide, left + 6150000, y_top + 360000, left + 6150000, y_bottom + 180000, color=MUTED, width_pt=1.2)
        add_arrow(slide, left + 6710000, y_top + 180000, left + 7360000, y_top + 180000, color=MUTED, width_pt=1.3)
        add_arrow(slide, left + 6710000, y_mid + 180000, left + 6880000, y_mid + 180000, color=MUTED, width_pt=1.3)
        add_line(slide, left + 6710000, y_mid + 180000, left + 6710000, y_bottom + 180000, color=MUTED, width_pt=1.2)
        add_arrow(slide, left + 6710000, y_bottom + 180000, left + 5460000, y_bottom + 180000, color=MUTED, width_pt=1.3)
        add_text(slide, left + 6460000, y_mid - 100000, 700000, 100000, "重新利用准备", font_size=10.2, color=INK)
        return

    y_bottom = top + 930000
    y_mid = top + 520000
    y_top = top + 150000
    status_box(slide, left + 160000, y_bottom, 1140000, 390000, "Pack\n正常", fill=GREEN)
    status_box(slide, left + 1540000, y_bottom, 1140000, 390000, "Pack\n故障", fill=GREEN)
    status_box(slide, left + 3100000, y_bottom, 1300000, 390000, "Pack\n恢复正常", fill=GREEN)
    status_box(slide, left + 4650000, y_mid, 1400000, 390000, "Pack\n作为备件使用", fill=TEAL, size=13.4)
    status_box(slide, left + 6420000, y_top, 1320000, 390000, "Pack\n回收商回收", fill=GRAY_SOFT, color=INK, size=13.6)
    status_box(slide, left + 7800000, y_bottom, 1320000, 390000, "Pack\n不更换用途", fill=LIME, size=13.4)
    status_box(slide, left + 9330000, y_mid, 1320000, 390000, "Pack\n更换用途", fill=LIME, size=13.4)
    status_box(slide, left + 10230000, y_top, 980000, 390000, "Pack\n报废", fill=GRAY_BOX)

    add_arrow(slide, left + 1300000, y_bottom + 195000, left + 1540000, y_bottom + 195000, color=MUTED, width_pt=1.4)
    add_arrow(slide, left + 2680000, y_bottom + 195000, left + 3100000, y_bottom + 195000, color=MUTED, width_pt=1.4)
    add_text(slide, left + 2850000, y_bottom - 130000, 520000, 110000, "现场维修", font_size=10.4, color=INK)

    add_line(slide, left + 2680000, y_bottom + 195000, left + 2680000, y_mid + 195000, color=MUTED, width_pt=1.3)
    add_arrow(slide, left + 2680000, y_mid + 195000, left + 4650000, y_mid + 195000, color=MUTED, width_pt=1.4)
    add_text(slide, left + 3480000, y_mid - 100000, 620000, 110000, "返厂维修", font_size=10.4, color=INK)

    add_line(slide, left + 2680000, y_bottom + 195000, left + 2680000, y_top + 195000, color=MUTED, width_pt=1.3)
    add_arrow(slide, left + 2680000, y_top + 195000, left + 6420000, y_top + 195000, color=MUTED, width_pt=1.4)
    add_text(slide, left + 4300000, y_top - 100000, 620000, 110000, "无法维修", font_size=10.4, color=INK)

    add_line(slide, left + 7080000, y_top + 390000, left + 7080000, y_bottom + 195000, color=MUTED, width_pt=1.3)
    add_arrow(slide, left + 7740000, y_top + 195000, left + 10230000, y_top + 195000, color=MUTED, width_pt=1.4)
    add_arrow(slide, left + 7740000, y_mid + 195000, left + 9330000, y_mid + 195000, color=MUTED, width_pt=1.4)
    add_line(slide, left + 7740000, y_mid + 195000, left + 7740000, y_bottom + 195000, color=MUTED, width_pt=1.3)
    add_arrow(slide, left + 7740000, y_bottom + 195000, left + 7800000, y_bottom + 195000, color=MUTED, width_pt=1.4)
    add_text(slide, left + 8400000, y_mid - 100000, 760000, 110000, "重新利用准备", font_size=10.4, color=INK)


def render_status_band(slide, *, left: int, top: int, compact: bool) -> None:
    widths = [1600000, 1450000, 1700000, 1450000, 1120000] if compact else [1750000, 1500000, 1750000, 1500000, 1100000]
    labels = [
        ("Original", GREEN),
        ("reused", TEAL),
        ("remanufactured", LIME),
        ("repurposed", (120, 188, 83)),
        ("wasted", GRAY_BOX),
    ]
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, 1080000, 250000, fill=GREEN)
    add_text(slide, left + 70000, top + 35000, 900000, 160000, "电池护照状态", font_size=10.4, color=INK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_VERTICAL_ANCHOR.MIDDLE)
    current_x = left + 1180000
    for width, (label, color) in zip(widths, labels):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.CHEVRON, current_x, top, width, 250000, fill=color)
        add_text(slide, current_x + 120000, top + 50000, width - 240000, 120000, label, font_size=11.0 if len(label) > 9 else 11.7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        current_x += width - 100000


def render_footer(slide, slide_width: int, *, top: int) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, top, slide_width - 1440000, 1, fill=LINE)
    add_text(slide, 720000, top + 120000, slide_width - 1440000, 180000, FOOTNOTE, font_size=10.1, color=MUTED)


def render_flow_variant(slide, slide_width: int) -> None:
    render_lead(slide, slide_width)
    render_definition_ribbon(slide)
    render_flow_canvas(slide, left=720000, top=2680000, width=slide_width - 1440000, height=2220000, compact=False)
    render_status_band(slide, left=720000, top=5050000, compact=False)
    render_footer(slide, slide_width, top=5460000)


def render_rule_card(slide, left: int, top: int, width: int, height: int, *, title: str, accent: tuple[int, int, int], summary: str, bullets: tuple[str, ...]) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=WHITE, line=LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, 50000, fill=accent)
    add_text(slide, left + 150000, top + 110000, width - 300000, 120000, title, font_size=12.4, color=INK, bold=True)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 150000, top + 280000, width - 300000, 210000, fill=ACCENT_SOFT if accent == ACCENT else (245, 247, 250))
    add_text(slide, left + 200000, top + 332000, width - 400000, 120000, summary, font_size=10.6, color=accent, bold=True)
    add_bullets(slide, left + 150000, top + 560000, width - 300000, bullets, bullet_color=accent, font_size=10.2, row_gap=240000)


def render_definition_stack(slide, left: int, top: int, width: int) -> None:
    add_text(slide, left, top, width, 120000, "状态定义", font_size=12.0, color=INK, bold=True)
    add_text(slide, left + width - 650000, top, 650000, 120000, "03", font_size=10.5, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.RIGHT)
    concise_details = (
        "按原用途继续使用",
        "恢复原用途性能后继续流转",
        "转为其他用途继续流转",
    )
    box_top = top + 170000
    box_height = 145000
    for idx, (card, detail) in enumerate(zip(DEFINITIONS, concise_details)):
        y = box_top + idx * 185000
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, y, width, box_height, fill=WHITE, line=LINE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 110000, y + 24000, 900000, 98000, fill=card.soft)
        add_text(slide, left + 165000, y + 42000, 820000, 70000, card.title, font_size=10.1, color=card.accent, bold=True)
        add_text(slide, left + 1090000, y + 35000, width - 1220000, 90000, detail, font_size=9.2, color=MUTED)


def render_decision_variant(slide, slide_width: int) -> None:
    render_lead(slide, slide_width)

    left_panel = 720000
    left_width = 7900000
    top = 1880000
    render_flow_canvas(slide, left=left_panel, top=top, width=left_width, height=2660000, compact=True)

    right_left = left_panel + left_width + 200000
    right_width = slide_width - right_left - 720000
    render_rule_card(
        slide,
        right_left,
        top,
        right_width,
        1050000,
        title="更新护照",
        accent=ACCENT,
        summary="满足检测与参数更新条件时执行",
        bullets=(
            "Pack 恢复正常后继续投入使用。",
            "作为备件复用，且用途不变。",
            "用途变更后继续流转，需同步更新用途与参数。",
        ),
    )
    render_rule_card(
        slide,
        right_left,
        top + 1180000,
        right_width,
        920000,
        title="注销护照",
        accent=GRAY_BOX,
        summary="缺乏检测条件或进入报废路径时执行",
        bullets=(
            "无法完成关键参数检测与验证。",
            "进入报废路径，不再继续流转。",
        ),
    )
    render_definition_stack(slide, right_left, top + 2240000, right_width)

    render_status_band(slide, left=720000, top=4760000, compact=True)
    render_footer(slide, slide_width, top=5210000)


def iter_shape_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text)
    return texts


def self_check_layout(prs: Presentation, *, narrative_mode: str = "faithful") -> list[str]:
    slide = prs.slides[0]
    issues: list[str] = []
    for idx, shape in enumerate(slide.shapes, 1):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            issues.append(f"shape {idx} is picture-based; expected editable output")
        if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
            issues.append(f"shape {idx} exceeds slide bounds")
    if narrative_mode == "faithful":
        all_text = "\n".join(iter_shape_texts(slide))
        for forbidden in FAITHFUL_FORBIDDEN_TEXTS:
            if forbidden in all_text:
                issues.append(f"found synthesized faithful-mode text: {forbidden}")
    return issues


def save_with_fallback(prs: Presentation, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        timestamp = datetime.datetime.now().strftime("%H%M%S_%f")[:9]
        fallback = output_path.with_stem(f"{output_path.stem}_{timestamp}")
        prs.save(str(fallback))
        return fallback


def build_slide(output_path: Path, *, layout: str, narrative_mode: str = "faithful") -> Path:
    prs = Presentation(str(DEFAULT_TEMPLATE))
    keep_only_body_template(prs)
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)

    render_title(slide)
    if layout == "flow":
        render_flow_variant(slide, slide_width)
    elif layout == "decision":
        render_decision_variant(slide, slide_width)
    else:
        raise ValueError(f"unsupported layout: {layout}")

    issues = self_check_layout(prs, narrative_mode=narrative_mode)
    if issues:
        raise ValueError("layout self-check failed: " + " | ".join(issues))
    return save_with_fallback(prs, output_path)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build SIE slides for battery passport update/cancel rules.")
    parser.add_argument("--layout", choices=["flow", "decision", "both"], default="both", help="Layout variant to render.")
    parser.add_argument("--narrative-mode", choices=["faithful", "consulting"], default="faithful", help="Content framing mode.")
    parser.add_argument("--output", help="Optional output path. Only valid when --layout is flow or decision.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    outputs: list[Path] = []
    if args.layout == "flow":
        output = Path(args.output).resolve() if args.output else OUTPUT_FLOW.resolve()
        outputs.append(build_slide(output, layout="flow", narrative_mode=args.narrative_mode))
    elif args.layout == "decision":
        output = Path(args.output).resolve() if args.output else OUTPUT_DECISION.resolve()
        outputs.append(build_slide(output, layout="decision", narrative_mode=args.narrative_mode))
    else:
        outputs.append(build_slide(OUTPUT_FLOW.resolve(), layout="flow", narrative_mode=args.narrative_mode))
        outputs.append(build_slide(OUTPUT_DECISION.resolve(), layout="decision", narrative_mode=args.narrative_mode))
    for output in outputs:
        print(output)


if __name__ == "__main__":
    main()
