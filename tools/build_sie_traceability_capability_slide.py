from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

try:
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest
except ModuleNotFoundError:
    import sys

    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from sie_autoppt.config import DEFAULT_TEMPLATE, FONT_NAME
    from sie_autoppt.slide_ops import remove_slide
    from sie_autoppt.template_manifest import load_template_manifest


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PPT = PROJECT_ROOT / "output" / "sie_supply_chain_traceability_capability.pptx"

ACCENT = (173, 5, 61)
INK = (35, 41, 46)
MUTED = (98, 109, 121)
LIGHT_TEXT = (146, 154, 161)
LINE = (224, 228, 232)
WHITE = (255, 255, 255)


@dataclass(frozen=True)
class CardSpec:
    index: str
    key: str
    title: str
    english: str
    accent: tuple[int, int, int]
    strip: tuple[int, int, int]
    summary: str
    bullets: tuple[str, str, str]
    closing: str


TITLE_SEGMENTS = [
    ("什么是", INK, True),
    ("“可用于对外证明”", ACCENT, True),
    ("的供应链追溯能力", INK, True),
]
SUBTITLE = "企业真正需要建设的，不是一套零散功能，而是一组同时面向客户、监管与海外市场的核心能力。"
HEADLINE = "对外审查并不止看“能不能查到”，而是进一步要求企业证明链路真实、证据可信、风险可控。"
SUPPORT = "换句话说，供应链追溯能力必须从“能查”升级为“能证”，最终落到“可控”。"
FOOTNOTE = "建议作为统一能力框架：先回答链路是否打通，再说明证据是否可信，最后证明风险是否可控。"

CARDS = (
    CardSpec(
        index="01",
        key="链",
        title="链路可追",
        english="Traceability",
        accent=(118, 86, 167),
        strip=(248, 243, 250),
        summary="先回答链路是否完整。",
        bullets=(
            "贯通矿源、原料、生产到交付的关键链路。",
            "支持订单、SN、批次级追溯与路径还原。",
            "形成跨供应商、跨工厂、跨系统的一体化视图。",
        ),
        closing="一句话：知道东西从哪来、怎么流转、最终到哪去。",
    ),
    CardSpec(
        index="02",
        key="证",
        title="证据可证",
        english="Proof",
        accent=ACCENT,
        strip=(252, 241, 245),
        summary="再回答证据是否可信。",
        bullets=(
            "业务数据、业务单据、合规文件一一对应。",
            "保证证明材料一致、可校验、可审计。",
            "支撑客户尽调、海关检查与审计核查。",
        ),
        closing="一句话：不仅有数据，更有能让外部相信的证据链。",
    ),
    CardSpec(
        index="03",
        key="控",
        title="风险可控",
        english="Control",
        accent=(53, 116, 136),
        strip=(241, 247, 249),
        summary="最终回答风险能否稳住。",
        bullets=(
            "快速锁定风险物料、风险批次与风险供应商。",
            "在抽查升级时快速响应、说明并出具材料。",
            "支撑风险隔离、经营决策与海外准入稳定交付。",
        ),
        closing="一句话：出了问题能说得清、找得准、控得住。",
    ),
)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def set_run_style(run, *, size: float, color: tuple[int, int, int], bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    font_size: float,
    color: tuple[int, int, int] = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = valign
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=font_size, color=color, bold=bold)
    return shape


def add_shape(
    slide,
    shape_type: MSO_AUTO_SHAPE_TYPE,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width: float = 1,
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    return shape


def keep_only_body_template(prs: Presentation) -> None:
    manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE)
    keep_idx = manifest.slide_roles.body_template
    for index in range(len(prs.slides) - 1, -1, -1):
        if index != keep_idx:
            remove_slide(prs, index)


def render_title(slide) -> None:
    title_shape = slide.shapes[0]
    tf = title_shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    for text, color, bold in TITLE_SEGMENTS:
        run = paragraph.add_run()
        run.text = text
        set_run_style(run, size=25, color=color, bold=bold)

    add_text(slide, 720000, 545000, 8400000, 190000, SUBTITLE, font_size=12.2, color=MUTED)


def render_judgement(slide, slide_width: int) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 1150000, 70000, 280000, fill=ACCENT)
    add_text(slide, 855000, 1140000, 1800000, 200000, "核心判断", font_size=13.2, color=ACCENT, bold=True)
    add_text(slide, 720000, 1420000, slide_width - 1440000, 360000, HEADLINE, font_size=18, color=INK, bold=True)
    add_text(slide, 720000, 1825000, slide_width - 1440000, 180000, SUPPORT, font_size=11.3, color=MUTED)


def card_geometries(slide_width: int) -> list[tuple[int, int, int, int]]:
    left = 720000
    gap = 240000
    usable_width = slide_width - left * 2
    width = int((usable_width - gap * 2) / 3)
    top = 2170000
    height = 2700000
    return [(left + idx * (width + gap), top, width, height) for idx in range(3)]


def render_card(slide, geom: tuple[int, int, int, int], spec: CardSpec) -> None:
    left, top, width, height = geom
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height, fill=WHITE, line=LINE, line_width=1)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, 65000, fill=spec.accent)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + 180000, top + 170000, 360000, 360000, fill=spec.accent)
    add_text(slide, left + 180000, top + 240000, 360000, 120000, spec.key, font_size=19, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, left + 640000, top + 150000, 1200000, 170000, spec.title, font_size=17.5, color=INK, bold=True)
    add_text(slide, left + 640000, top + 360000, 1000000, 110000, spec.english, font_size=8.5, color=spec.accent, bold=True)
    add_text(slide, left + width - 450000, top + 185000, 280000, 120000, spec.index, font_size=10.5, color=LIGHT_TEXT, bold=True, align=PP_ALIGN.RIGHT)

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left + 180000, top + 640000, width - 360000, 320000, fill=spec.strip)
    add_text(slide, left + 235000, top + 708000, width - 470000, 180000, spec.summary, font_size=12.8, color=spec.accent, bold=True)

    bullet_y = top + 1080000
    for bullet in spec.bullets:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, left + 185000, bullet_y + 65000, 70000, 70000, fill=spec.accent)
        add_text(slide, left + 315000, bullet_y, width - 470000, 250000, bullet, font_size=10.6, color=MUTED)
        bullet_y += 430000

    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, left + 180000, top + height - 420000, width - 360000, 1, fill=LINE)
    add_text(slide, left + 180000, top + height - 325000, width - 360000, 180000, spec.closing, font_size=10.1, color=INK, bold=True)


def render_footer(slide, slide_width: int) -> None:
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 720000, 5270000, slide_width - 1440000, 1, fill=LINE)
    add_text(slide, 720000, 5390000, slide_width - 1440000, 160000, FOOTNOTE, font_size=10.5, color=MUTED)


def build_slide(output_path: Path) -> Path:
    prs = Presentation(str(DEFAULT_TEMPLATE))
    keep_only_body_template(prs)
    slide = prs.slides[0]
    slide_width = int(prs.slide_width)

    render_title(slide)
    render_judgement(slide, slide_width)
    for geom, spec in zip(card_geometries(slide_width), CARDS):
        render_card(slide, geom, spec)
    render_footer(slide, slide_width)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        timestamp = datetime.datetime.now().strftime("%H%M%S_%f")[:9]
        fallback = output_path.with_stem(f"{output_path.stem}_{timestamp}")
        prs.save(str(fallback))
        return fallback


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build a consulting-style SIE body slide for externally provable supply-chain traceability capability.")
    parser.add_argument("--output", default=str(OUTPUT_PPT), help="Output PPTX path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = build_slide(Path(args.output).resolve())
    print(output)


if __name__ == "__main__":
    main()
