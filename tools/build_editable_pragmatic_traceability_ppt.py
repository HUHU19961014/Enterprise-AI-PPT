from __future__ import annotations

import argparse
import math
from dataclasses import dataclass
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_HTML = PROJECT_ROOT / "input" / "pragmatic_traceability_planning.html"
DEFAULT_OUTPUT = PROJECT_ROOT / "output" / "pragmatic_traceability_planning_editable.pptx"
DEFAULT_SLIDE_WIDTH = Inches(13.333333)
DEFAULT_SLIDE_HEIGHT = Inches(7.5)

SLIDE_W = 1200
SLIDE_H = 675

SIE_WINE = (138, 21, 56)
TEXT_MAIN = (51, 51, 51)
TEXT_MUTED = (102, 102, 102)
BG_GRAY = (245, 247, 250)
DARK_GRAY = (74, 74, 74)
BORDER = (224, 224, 224)
LIGHT_GRAY = (208, 208, 208)
WHITE = (255, 255, 255)
SHADOW = (120, 120, 120)


@dataclass(frozen=True)
class RichTextPart:
    text: str
    highlight: bool = False


@dataclass(frozen=True)
class BulletItem:
    prefix: str
    detail: str


@dataclass(frozen=True)
class TextSection:
    title: str
    accent_color: tuple[int, int, int]
    items: list[BulletItem]


@dataclass(frozen=True)
class PointSpec:
    name: str
    x_ratio: float
    y_ratio: float
    size_px: float
    fill: tuple[int, int, int]
    line: tuple[int, int, int] | None
    label_dx_px: float
    label_dy_px: float
    label_width_px: float
    label_color: tuple[int, int, int]
    label_size: float


@dataclass(frozen=True)
class SlideSpec:
    title_parts: list[RichTextPart]
    system_lines: list[str]
    manual_lines: list[str]
    x_axis: str
    y_axis: str
    chasm: str
    arrow: str
    sections: list[TextSection]
    value_title: str
    value_items: list[BulletItem]
    points: list[PointSpec]


def px_x(value: float, slide_width: Emu) -> Emu:
    return int(value / SLIDE_W * slide_width)


def px_y(value: float, slide_height: Emu) -> Emu:
    return int(value / SLIDE_H * slide_height)


def rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def normalize_text(text: str) -> str:
    replacements = {
        "锛?": "：",
        "Ŗē": "：",
        "鈥?": "•",
        "鉁?": "✓",
    }
    cleaned = text
    for old, new in replacements.items():
        cleaned = cleaned.replace(old, new)
    return " ".join(cleaned.replace("\xa0", " ").split())


def weighted_text_length(text: str) -> float:
    total = 0.0
    for char in text:
        if char.isspace():
            total += 0.35
        elif ord(char) < 128:
            total += 0.6
        else:
            total += 1.0
    return total


def estimate_lines(text: str, width_px: float, font_size: float) -> int:
    if not text:
        return 1
    chars_per_line = max(1.0, width_px / max(font_size * 1.15, 1))
    return max(1, math.ceil(weighted_text_length(text) / chars_per_line))


def estimate_text_height_px(text: str, width_px: float, font_size: float, line_height: float = 1.45, min_lines: int = 1) -> float:
    return max(min_lines, estimate_lines(text, width_px, font_size)) * font_size * line_height


def apply_shadow(shape, *, transparency: float = 0.88) -> None:
    shape.fill.transparency = transparency
    shape.line.fill.background()


def add_shadow_rect(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    dx: Emu,
    dy: Emu,
    rounded: bool = False,
    transparency: float = 0.88,
):
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left + dx,
        top + dy,
        width,
        height,
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb(SHADOW)
    apply_shadow(shadow, transparency=transparency)
    return shadow


def add_shadow_circle(
    slide,
    center_x: Emu,
    center_y: Emu,
    diameter: Emu,
    *,
    dx: Emu,
    dy: Emu,
    transparency: float = 0.82,
):
    shadow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        center_x - diameter // 2 + dx,
        center_y - diameter // 2 + dy,
        diameter,
        diameter,
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = rgb(SHADOW)
    apply_shadow(shadow, transparency=transparency)
    return shadow


def add_textbox(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    font_name: str = "Microsoft YaHei",
    font_size: float = 14,
    color: tuple[int, int, int] = TEXT_MAIN,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_VERTICAL_ANCHOR = MSO_VERTICAL_ANCHOR.TOP,
    fill: tuple[int, int, int] | None = None,
    line: tuple[int, int, int] | None = None,
    line_width_pt: float = 1,
    shape_type: MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE.RECTANGLE,
):
    if fill is not None or line is not None:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Pt(line_width_pt)
        text_frame = shape.text_frame
    else:
        shape = slide.shapes.add_textbox(left, top, width, height)
        text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = vertical_anchor
    text_frame.margin_left = Inches(0.02)
    text_frame.margin_right = Inches(0.02)
    text_frame.margin_top = Inches(0.01)
    text_frame.margin_bottom = Inches(0.01)
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape, paragraph, run


def add_rect(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    *,
    fill: tuple[int, int, int],
    line: tuple[int, int, int] | None = None,
    line_width_pt: float = 1,
    rounded: bool = False,
    shadow: bool = False,
    shadow_dx_px: float = 4,
    shadow_dy_px: float = 4,
    shadow_transparency: float = 0.9,
):
    if shadow:
        add_shadow_rect(
            slide,
            left,
            top,
            width,
            height,
            dx=px_x(shadow_dx_px, DEFAULT_SLIDE_WIDTH),
            dy=px_y(shadow_dy_px, DEFAULT_SLIDE_HEIGHT),
            rounded=rounded,
            transparency=shadow_transparency,
        )
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top,
        width,
        height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width_pt)
    return shape


def add_circle(
    slide,
    center_x: Emu,
    center_y: Emu,
    diameter: Emu,
    *,
    fill: tuple[int, int, int],
    line: tuple[int, int, int] | None = None,
    line_width_pt: float = 1,
    shadow: bool = False,
):
    if shadow:
        add_shadow_circle(
            slide,
            center_x,
            center_y,
            diameter,
            dx=px_x(2, DEFAULT_SLIDE_WIDTH),
            dy=px_y(3, DEFAULT_SLIDE_HEIGHT),
        )
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        center_x - diameter // 2,
        center_y - diameter // 2,
        diameter,
        diameter,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width_pt)
    return shape


def add_line(
    slide,
    x1: Emu,
    y1: Emu,
    x2: Emu,
    y2: Emu,
    *,
    color: tuple[int, int, int],
    width_pt: float = 1,
    dash: MSO_LINE_DASH_STYLE | None = None,
):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width_pt)
    if dash is not None:
        line.line.dash_style = dash
    return line


def extract_text_with_highlight(node: Tag) -> list[RichTextPart]:
    parts: list[RichTextPart] = []
    for child in node.children:
        if isinstance(child, NavigableString):
            text = normalize_text(str(child))
            if text:
                parts.append(RichTextPart(text))
        elif isinstance(child, Tag):
            parts.append(
                RichTextPart(
                    normalize_text(child.get_text(strip=False)),
                    "highlight" in set(child.get("class", [])),
                )
            )
    return [part for part in parts if part.text]


def extract_li_parts(li_node: Tag) -> BulletItem:
    li_copy = BeautifulSoup(str(li_node), "html.parser").find("li")
    strong = li_copy.find("strong")
    if strong is None:
        return BulletItem("", normalize_text(li_copy.get_text(" ", strip=True)))
    prefix = normalize_text(strong.get_text(" ", strip=True)).rstrip("：:")
    strong.extract()
    detail = normalize_text(li_copy.get_text(" ", strip=True)).lstrip("：: ")
    return BulletItem(prefix, detail)


def extract_label_lines(node: Tag) -> list[str]:
    node_copy = BeautifulSoup(str(node), "html.parser")
    for br in node_copy.find_all("br"):
        br.replace_with("\n")
    return [normalize_text(line) for line in node_copy.get_text("\n", strip=True).splitlines() if normalize_text(line)]


def build_point_specs(soup: BeautifulSoup) -> list[PointSpec]:
    names = {
        "yingfa": normalize_text(soup.select_one(".pt-yingfa .point-label").get_text(" ", strip=True)),
        "trina": normalize_text(soup.select_one(".pt-top.trina .point-label").get_text(" ", strip=True)),
        "ja": normalize_text(soup.select_one(".pt-top.ja .point-label").get_text(" ", strip=True)),
        "jinko": normalize_text(soup.select_one(".pt-top.jinko .point-label").get_text(" ", strip=True)),
        "yidao": normalize_text(soup.select_one(".pt-gray.yidao .point-label").get_text(" ", strip=True)),
        "hengdian": normalize_text(soup.select_one(".pt-gray.hengdian .point-label").get_text(" ", strip=True)),
    }
    return [
        PointSpec(names["yingfa"], 0.35, 0.60, 20, DARK_GRAY, WHITE, -22, -30, 90, TEXT_MAIN, 15),
        PointSpec(names["trina"], 0.85, 0.20, 14, SIE_WINE, None, -76, -12, 72, SIE_WINE, 13),
        PointSpec(names["ja"], 0.80, 0.25, 14, SIE_WINE, None, 18, -10, 72, SIE_WINE, 13),
        PointSpec(names["jinko"], 0.75, 0.30, 14, SIE_WINE, None, -8, -34, 72, SIE_WINE, 13),
        PointSpec(names["yidao"], 0.45, 0.55, 10, LIGHT_GRAY, None, 16, -2, 74, (153, 153, 153), 12),
        PointSpec(names["hengdian"], 0.48, 0.60, 10, LIGHT_GRAY, None, 16, -2, 74, (153, 153, 153), 12),
    ]


def parse_html(html_path: Path) -> SlideSpec:
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    title = soup.select_one(".header h1")
    sections: list[TextSection] = []
    for block in soup.select(".text-block"):
        title_node = block.select_one(".text-block-title")
        accent = DARK_GRAY if "gray" in title_node.get("class", []) else SIE_WINE
        items = [extract_li_parts(li) for li in block.select("li")]
        sections.append(TextSection(normalize_text(title_node.get_text(" ", strip=True)), accent, items))

    return SlideSpec(
        title_parts=extract_text_with_highlight(title),
        system_lines=extract_label_lines(soup.select_one(".q-tr .q-label")),
        manual_lines=extract_label_lines(soup.select_one(".q-bl .q-label")),
        x_axis=normalize_text(soup.select_one(".x-axis-label").get_text(" ", strip=True)),
        y_axis=normalize_text(soup.select_one(".y-axis-label").get_text(" ", strip=True)),
        chasm=normalize_text(soup.select_one(".chasm-label").get_text(" ", strip=True)),
        arrow=normalize_text(soup.select_one(".arrow-label").get_text(" ", strip=True)),
        sections=sections,
        value_title=normalize_text(soup.select_one(".value-box-title").get_text(" ", strip=True)),
        value_items=[extract_li_parts(li) for li in soup.select(".value-box li")],
        points=build_point_specs(soup),
    )


def add_title(slide, slide_width: Emu, slide_height: Emu, spec: SlideSpec) -> None:
    _, paragraph, _ = add_textbox(
        slide,
        px_x(40, slide_width),
        px_y(24, slide_height),
        px_x(1120, slide_width),
        px_y(42, slide_height),
        font_size=24,
        color=TEXT_MAIN,
        bold=True,
    )
    paragraph.clear()
    paragraph.alignment = PP_ALIGN.LEFT
    for part in spec.title_parts:
        run = paragraph.add_run()
        run.text = part.text
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = rgb(SIE_WINE if part.highlight else TEXT_MAIN)

    divider = add_rect(
        slide,
        px_x(40, slide_width),
        px_y(84, slide_height),
        px_x(1080, slide_width),
        px_y(2, slide_height),
        fill=SIE_WINE,
    )
    divider.line.fill.background()


def add_axis_labels(slide, matrix_box: dict, spec: SlideSpec, slide_width: Emu, slide_height: Emu) -> None:
    y_shape, _, y_run = add_textbox(
        slide,
        px_x(-62, slide_width),
        matrix_box["top"] + px_y(220, slide_height),
        px_x(220, slide_width),
        px_y(24, slide_height),
        font_size=13,
        color=TEXT_MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE,
    )
    y_run.text = spec.y_axis
    y_shape.rotation = 270

    _, _, x_run = add_textbox(
        slide,
        matrix_box["left"] + px_x(126, slide_width),
        matrix_box["top"] + matrix_box["height"] + px_y(12, slide_height),
        px_x(320, slide_width),
        px_y(24, slide_height),
        font_size=14,
        color=TEXT_MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    x_run.text = spec.x_axis


def add_matrix(slide, spec: SlideSpec, slide_width: Emu, slide_height: Emu) -> None:
    matrix_box = {
        "left": px_x(75, slide_width),
        "top": px_y(128, slide_height),
        "width": px_x(545, slide_width),
        "height": px_y(500, slide_height),
    }

    add_rect(
        slide,
        matrix_box["left"],
        matrix_box["top"],
        matrix_box["width"],
        matrix_box["height"],
        fill=BG_GRAY,
        shadow=True,
        shadow_dx_px=3,
        shadow_dy_px=3,
        shadow_transparency=0.95,
    )
    add_rect(
        slide,
        matrix_box["left"] + matrix_box["width"] // 2,
        matrix_box["top"],
        matrix_box["width"] // 2,
        matrix_box["height"] // 2,
        fill=(250, 244, 246),
    )

    add_axis_labels(slide, matrix_box, spec, slide_width, slide_height)

    add_line(slide, matrix_box["left"], matrix_box["top"], matrix_box["left"], matrix_box["top"] + matrix_box["height"], color=DARK_GRAY, width_pt=2)
    add_line(
        slide,
        matrix_box["left"],
        matrix_box["top"] + matrix_box["height"],
        matrix_box["left"] + matrix_box["width"],
        matrix_box["top"] + matrix_box["height"],
        color=DARK_GRAY,
        width_pt=2,
    )
    add_line(
        slide,
        matrix_box["left"] + matrix_box["width"] // 2,
        matrix_box["top"],
        matrix_box["left"] + matrix_box["width"] // 2,
        matrix_box["top"] + matrix_box["height"],
        color=(204, 204, 204),
        width_pt=1,
        dash=MSO_LINE_DASH_STYLE.DASH,
    )
    add_line(
        slide,
        matrix_box["left"],
        matrix_box["top"] + matrix_box["height"] // 2,
        matrix_box["left"] + matrix_box["width"],
        matrix_box["top"] + matrix_box["height"] // 2,
        color=(204, 204, 204),
        width_pt=1,
        dash=MSO_LINE_DASH_STYLE.DASH,
    )

    qtr_shape, qtr_paragraph, _ = add_textbox(
        slide,
        matrix_box["left"] + px_x(390, slide_width),
        matrix_box["top"] + px_y(12, slide_height),
        px_x(146, slide_width),
        px_y(60, slide_height),
        font_size=15,
        color=SIE_WINE,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )
    qtr_paragraph.clear()
    run1 = qtr_paragraph.add_run()
    run1.text = spec.system_lines[0]
    run1.font.name = "Microsoft YaHei"
    run1.font.size = Pt(15)
    run1.font.bold = True
    run1.font.color.rgb = rgb(SIE_WINE)
    sub = qtr_shape.text_frame.add_paragraph()
    sub.alignment = PP_ALIGN.RIGHT
    sub_run = sub.add_run()
    sub_run.text = spec.system_lines[1]
    sub_run.font.name = "Microsoft YaHei"
    sub_run.font.size = Pt(11.5)
    sub_run.font.color.rgb = rgb(TEXT_MUTED)

    qbl_shape, qbl_paragraph, _ = add_textbox(
        slide,
        matrix_box["left"] + px_x(16, slide_width),
        matrix_box["top"] + px_y(432, slide_height),
        px_x(192, slide_width),
        px_y(60, slide_height),
        font_size=15,
        color=TEXT_MUTED,
        bold=True,
    )
    qbl_paragraph.clear()
    q_run = qbl_paragraph.add_run()
    q_run.text = spec.manual_lines[0]
    q_run.font.name = "Microsoft YaHei"
    q_run.font.size = Pt(15)
    q_run.font.bold = True
    q_run.font.color.rgb = rgb(TEXT_MUTED)
    sub2 = qbl_shape.text_frame.add_paragraph()
    sub2_run = sub2.add_run()
    sub2_run.text = spec.manual_lines[1]
    sub2_run.font.name = "Microsoft YaHei"
    sub2_run.font.size = Pt(11.5)
    sub2_run.font.color.rgb = rgb(TEXT_MUTED)

    add_line(
        slide,
        matrix_box["left"] + px_x(110, slide_width),
        matrix_box["top"] + px_y(75, slide_height),
        matrix_box["left"] + px_x(535, slide_width),
        matrix_box["top"] + px_y(530, slide_height),
        color=(187, 187, 187),
        width_pt=2,
        dash=MSO_LINE_DASH_STYLE.DASH,
    )

    chasm_label, _, chasm_run = add_textbox(
        slide,
        matrix_box["left"] + px_x(350, slide_width),
        matrix_box["top"] + px_y(247, slide_height),
        px_x(118, slide_width),
        px_y(24, slide_height),
        font_size=14,
        color=TEXT_MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
        fill=WHITE,
    )
    chasm_label.rotation = 315
    chasm_run.text = spec.chasm

    arrow_bar = add_rect(
        slide,
        matrix_box["left"] + px_x(198, slide_width),
        matrix_box["top"] + px_y(262, slide_height),
        px_x(156, slide_width),
        px_y(4, slide_height),
        fill=SIE_WINE,
    )
    arrow_bar.rotation = 325

    arrow_head = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        matrix_box["left"] + px_x(330, slide_width),
        matrix_box["top"] + px_y(214, slide_height),
        px_x(20, slide_width),
        px_y(20, slide_height),
    )
    arrow_head.fill.solid()
    arrow_head.fill.fore_color.rgb = rgb(SIE_WINE)
    arrow_head.line.fill.background()
    arrow_head.rotation = 55

    arrow_label, _, arrow_run = add_textbox(
        slide,
        matrix_box["left"] + px_x(186, slide_width),
        matrix_box["top"] + px_y(208, slide_height),
        px_x(126, slide_width),
        px_y(26, slide_height),
        font_size=12.5,
        color=SIE_WINE,
        bold=True,
        align=PP_ALIGN.CENTER,
        fill=WHITE,
        line=SIE_WINE,
        shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
    )
    arrow_label.rotation = 325
    arrow_run.text = spec.arrow

    for point in spec.points:
        center_x = matrix_box["left"] + int(matrix_box["width"] * point.x_ratio)
        center_y = matrix_box["top"] + int(matrix_box["height"] * point.y_ratio)
        add_circle(
            slide,
            center_x,
            center_y,
            px_x(point.size_px, slide_width),
            fill=point.fill,
            line=point.line,
            line_width_pt=2 if point.line else 1,
            shadow=point.name in {"英发睿能", "天合光能", "晶澳科技", "晶科能源"},
        )
        _, _, label_run = add_textbox(
            slide,
            center_x + px_x(point.label_dx_px, slide_width),
            center_y + px_y(point.label_dy_px, slide_height),
            px_x(point.label_width_px, slide_width),
            px_y(24, slide_height),
            font_size=point.label_size,
            color=point.label_color,
            bold=True,
        )
        label_run.text = point.name


def add_bullet_item(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    item: BulletItem,
    *,
    slide_width: Emu,
    slide_height: Emu,
    color: tuple[int, int, int],
    dot_color: tuple[int, int, int],
    font_size: float = 13.2,
    detail_size: float = 12.2,
) -> float:
    _, _, dot_run = add_textbox(
        slide,
        left,
        top,
        px_x(14, slide_width),
        px_y(18, slide_height),
        font_size=13,
        color=dot_color,
        bold=True,
    )
    dot_run.text = "•"

    full_text = f"{item.prefix}：{item.detail}" if item.prefix else item.detail
    width_px = width / slide_width * SLIDE_W - 14
    text_height_px = estimate_text_height_px(full_text, width_px, detail_size, line_height=1.5, min_lines=2)
    box_height_px = max(34, text_height_px + 8)

    _, paragraph, _ = add_textbox(
        slide,
        left + px_x(14, slide_width),
        top - px_y(1, slide_height),
        width - px_x(14, slide_width),
        px_y(box_height_px, slide_height),
        font_size=font_size,
        color=color,
        vertical_anchor=MSO_VERTICAL_ANCHOR.TOP,
    )
    paragraph.clear()
    if item.prefix:
        run1 = paragraph.add_run()
        run1.text = f"{item.prefix}："
        run1.font.name = "Microsoft YaHei"
        run1.font.size = Pt(font_size)
        run1.font.bold = True
        run1.font.color.rgb = rgb(TEXT_MAIN)
    run2 = paragraph.add_run()
    run2.text = item.detail
    run2.font.name = "Microsoft YaHei"
    run2.font.size = Pt(detail_size)
    run2.font.color.rgb = rgb(color)
    return box_height_px


def add_right_panel(slide, spec: SlideSpec, slide_width: Emu, slide_height: Emu) -> None:
    panel_left = px_x(660, slide_width)
    panel_width = px_x(470, slide_width)
    content_width = px_x(430, slide_width)
    current_y_px = 120

    for section in spec.sections:
        bar = add_rect(
            slide,
            panel_left,
            px_y(current_y_px, slide_height),
            px_x(4, slide_width),
            px_y(16, slide_height),
            fill=section.accent_color,
            shadow=True,
            shadow_dx_px=1,
            shadow_dy_px=2,
            shadow_transparency=0.9,
        )
        bar.line.fill.background()

        _, _, title_run = add_textbox(
            slide,
            panel_left + px_x(10, slide_width),
            px_y(current_y_px - 2, slide_height),
            panel_width,
            px_y(28, slide_height),
            font_size=15.5,
            color=TEXT_MAIN,
            bold=True,
        )
        title_run.text = section.title

        divider = add_rect(
            slide,
            panel_left + px_x(10, slide_width),
            px_y(current_y_px + 28, slide_height),
            content_width,
            px_y(1, slide_height),
            fill=BORDER,
        )
        divider.line.fill.background()

        bullet_y_px = current_y_px + 52
        for item in section.items:
            used_height = add_bullet_item(
                slide,
                panel_left,
                px_y(bullet_y_px, slide_height),
                content_width,
                item,
                slide_width=slide_width,
                slide_height=slide_height,
                color=TEXT_MUTED,
                dot_color=TEXT_MAIN,
            )
            bullet_y_px += used_height + 14
        current_y_px = bullet_y_px + 16

    value_left = panel_left
    value_top = px_y(455, slide_height)
    value_width = px_x(440, slide_width)
    value_height = px_y(202, slide_height)
    add_rect(
        slide,
        value_left,
        value_top,
        value_width,
        value_height,
        fill=SIE_WINE,
        rounded=True,
    )

    _, _, title_run = add_textbox(
        slide,
        value_left + px_x(16, slide_width),
        value_top + px_y(12, slide_height),
        px_x(300, slide_width),
        px_y(26, slide_height),
        font_size=15.5,
        color=WHITE,
        bold=True,
    )
    title_run.text = spec.value_title

    divider = add_rect(
        slide,
        value_left + px_x(16, slide_width),
        value_top + px_y(44, slide_height),
        value_width - px_x(32, slide_width),
        px_y(1, slide_height),
        fill=WHITE,
    )
    divider.fill.transparency = 0.8
    divider.line.fill.background()

    current_y_px = 500
    for item in spec.value_items:
        _, _, tick_run = add_textbox(
            slide,
            value_left + px_x(16, slide_width),
            px_y(current_y_px, slide_height),
            px_x(14, slide_width),
            px_y(18, slide_height),
            font_size=13,
            color=WHITE,
            bold=True,
        )
        tick_run.text = "✓"

        full_text = f"{item.prefix}：{item.detail}" if item.prefix else item.detail
        width_px = (value_width - px_x(48, slide_width)) / slide_width * SLIDE_W
        box_height_px = max(34, estimate_text_height_px(full_text, width_px, 12.4, line_height=1.54, min_lines=2) + 6)
        _, paragraph, _ = add_textbox(
            slide,
            value_left + px_x(32, slide_width),
            px_y(current_y_px - 1, slide_height),
            value_width - px_x(48, slide_width),
            px_y(box_height_px, slide_height),
            font_size=13.2,
            color=WHITE,
        )
        paragraph.clear()
        run1 = paragraph.add_run()
        run1.text = f"{item.prefix}："
        run1.font.name = "Microsoft YaHei"
        run1.font.size = Pt(13.2)
        run1.font.bold = True
        run1.font.color.rgb = rgb(WHITE)
        run2 = paragraph.add_run()
        run2.text = item.detail
        run2.font.name = "Microsoft YaHei"
        run2.font.size = Pt(12.4)
        run2.font.color.rgb = rgb(WHITE)
        current_y_px += box_height_px + 10


def build_ppt(html_path: Path, output_path: Path) -> Path:
    spec = parse_html(html_path)
    prs = Presentation()
    prs.slide_width = DEFAULT_SLIDE_WIDTH
    prs.slide_height = DEFAULT_SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, fill=WHITE)
    add_title(slide, prs.slide_width, prs.slide_height, spec)
    add_matrix(slide, spec, prs.slide_width, prs.slide_height)
    add_right_panel(slide, spec, prs.slide_width, prs.slide_height)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build an editable single-page PPT from a pragmatic traceability HTML layout.")
    parser.add_argument("--html", default=str(DEFAULT_HTML), help="Source HTML path.")
    parser.add_argument("--output", default=str(DEFAULT_OUTPUT), help="Target PPTX path.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    output = build_ppt(Path(args.html).resolve(), Path(args.output).resolve())
    print(str(output))


if __name__ == "__main__":
    main()
