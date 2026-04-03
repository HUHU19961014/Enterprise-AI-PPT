import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .config import COLOR_ACTIVE, FONT_NAME, IDX_BODY_TEMPLATE


REFERENCE_STYLE_LIBRARY = {
    "comparison_upgrade": {
        "name": "左右对比升级页",
        "source_slide": 5,
        "scenes": ["价值对比", "旧方式 vs 新方式", "能力边界与工程化解法"],
    },
    "capability_ring": {
        "name": "能力亮点环形页",
        "source_slide": 16,
        "scenes": ["能力亮点", "实践原则", "方案卖点"],
    },
    "five_phase_path": {
        "name": "五阶段推进路径页",
        "source_slide": 20,
        "scenes": ["实施路径", "方法路线", "阶段推进"],
    },
    "pain_cards": {
        "name": "三栏痛点拆解页",
        "source_slide": 6,
        "scenes": ["痛点拆解", "挑战分析", "风险拆解"],
    },
}


def get_reference_slide_no(style_id: str | None) -> int | None:
    if not style_id:
        return None
    style = REFERENCE_STYLE_LIBRARY.get(style_id)
    if not style:
        return None
    return int(style["source_slide"])


def build_reference_import_plan(body_pages) -> list[tuple[int, int]]:
    plan = []
    for offset, page in enumerate(body_pages):
        source_slide = get_reference_slide_no(getattr(page, "reference_style_id", None))
        if source_slide is None:
            continue
        target_slide = IDX_BODY_TEMPLATE + 1 + offset * 2
        plan.append((target_slide, source_slide))
    return plan


def _replace_text(shape, text: str, color=None, size_pt: int | None = None, align=PP_ALIGN.LEFT, bold=None):
    if not getattr(shape, "has_text_frame", False):
        return
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            if color is not None:
                run.font.color.rgb = RGBColor(*color)
            if size_pt is not None:
                run.font.size = Pt(size_pt)
            if bold is not None:
                run.font.bold = bold


def _shape(slide, index: int):
    return slide.shapes[index - 1]


def _normalize(text: str, max_chars: int) -> str:
    compact = re.sub(r"\s+", " ", text).strip()
    if len(compact) <= max_chars:
        return compact
    return compact[: max_chars - 1].rstrip(" ,;，；。") + "…"


def _cards(payload, key: str, count: int) -> list[dict[str, str]]:
    cards = list(payload.get(key, []))
    while len(cards) < count:
        cards.append({"title": "", "detail": ""})
    return cards[:count]


def fill_reference_style_slide(slide, page) -> bool:
    style_id = getattr(page, "reference_style_id", None)
    if style_id == "comparison_upgrade":
        _fill_comparison_upgrade(slide, page)
        return True
    if style_id == "capability_ring":
        _fill_capability_ring(slide, page)
        return True
    if style_id == "five_phase_path":
        _fill_five_phase_path(slide, page)
        return True
    if style_id == "pain_cards":
        _fill_pain_cards(slide, page)
        return True
    return False


def populate_reference_body_pages(pptx_path: Path, body_pages) -> bool:
    has_reference_pages = any(getattr(page, "reference_style_id", None) for page in body_pages)
    if not has_reference_pages:
        return False
    prs = Presentation(str(pptx_path))
    for offset, page in enumerate(body_pages):
        if not getattr(page, "reference_style_id", None):
            continue
        slide_index = IDX_BODY_TEMPLATE + offset * 2
        if slide_index >= len(prs.slides):
            continue
        fill_reference_style_slide(prs.slides[slide_index], page)
    prs.save(str(pptx_path))
    return True


def _fill_comparison_upgrade(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    _replace_text(_shape(slide, 56), payload.get("headline", page.subtitle), color=(0, 0, 0), size_pt=18, align=PP_ALIGN.CENTER, bold=True)

    _replace_text(_shape(slide, 2), payload.get("left_label", "● 传统方式"), color=(91, 112, 135), size_pt=17, bold=True)
    _replace_text(_shape(slide, 15), payload.get("right_label", "● 工程化方式"), color=COLOR_ACTIVE, size_pt=17, bold=True)

    for slot, card in zip([4, 7, 10, 13], _cards(payload, "left_cards", 4)):
        _replace_text(_shape(slide, slot), card.get("title", ""), color=(60, 76, 96), size_pt=16, bold=True)
    for slot, card in zip([5, 8, 11, 14], _cards(payload, "left_cards", 4)):
        detail = _cards(payload, "left_cards", 4)[[5, 8, 11, 14].index(slot)].get("detail", "")
        _replace_text(_shape(slide, slot), detail, color=(91, 112, 135), size_pt=10)

    for slot, card in zip([17, 20, 23, 26], _cards(payload, "right_cards", 4)):
        _replace_text(_shape(slide, slot), card.get("title", ""), color=(255, 255, 255), size_pt=16, bold=True)
    for slot, card in zip([18, 21, 24, 27], _cards(payload, "right_cards", 4)):
        detail = _cards(payload, "right_cards", 4)[[18, 21, 24, 27].index(slot)].get("detail", "")
        _replace_text(_shape(slide, slot), detail, color=(255, 255, 255), size_pt=10)

    _replace_text(_shape(slide, 30), payload.get("center_kicker", "ENGINEERING COPILOT"), color=(255, 255, 255), size_pt=12, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 31), payload.get("center_title", "结构化内容交付中枢"), color=COLOR_ACTIVE, size_pt=24, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 32), payload.get("center_subtitle", ""), color=(91, 112, 135), size_pt=12, align=PP_ALIGN.CENTER)
    _replace_text(_shape(slide, 35), payload.get("center_top_label", "协同原则层"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 36), payload.get("center_section_title", "人机协同 (Human in the Loop)"), color=(0, 0, 0), size_pt=18, bold=True)
    _replace_text(_shape(slide, 38), payload.get("center_row1_left", "结构化输入"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 40), payload.get("center_row1_right", "模板化渲染"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 42), payload.get("center_row2_left", "视觉 QA"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 44), payload.get("center_row2_right", "人工微调"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 45), payload.get("center_divider", "♦ 保持骨架稳定，保留创作空间 ♦"), color=COLOR_ACTIVE, size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 48), payload.get("center_bottom_label", "交付结果层"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 49), payload.get("center_bottom_title", "高质量输出 (Deliverable)"), color=(0, 0, 0), size_pt=18, bold=True)
    _replace_text(_shape(slide, 51), payload.get("center_bottom_left", "稳定复用"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 53), payload.get("center_bottom_right", "风格一致"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 55), payload.get("center_bottom_footer", "效率与质量兼得"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 60), payload.get("bottom_left_caption", "能力边界"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 61), payload.get("bottom_right_caption", "工程化交付"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)


def _fill_capability_ring(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    items = list(payload.get("items", []))
    while len(items) < 7:
        items.append({"title": "", "detail": ""})
    for title_idx, detail_idx, item in zip([18, 20, 22, 24, 26, 28, 30], [19, 21, 23, 25, 27, 29, 31], items[:7]):
        _replace_text(_shape(slide, title_idx), _normalize(item.get("title", ""), 14), color=(255, 255, 255), size_pt=15, align=PP_ALIGN.CENTER, bold=True)
        _replace_text(_shape(slide, detail_idx), _normalize(item.get("detail", ""), 40), color=(60, 76, 96), size_pt=10, align=PP_ALIGN.CENTER)


def _fill_five_phase_path(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    _replace_text(_shape(slide, 60), payload.get("intro", page.subtitle), color=(0, 0, 0), size_pt=16, bold=True)

    stages = list(payload.get("stages", []))
    while len(stages) < 5:
        stages.append({"header": "", "tasks": []})
    for slot, stage in zip([6, 4, 2, 8, 10], stages[:5]):
        _replace_text(_shape(slide, slot), stage.get("header", ""), color=(255, 255, 255), size_pt=14, align=PP_ALIGN.CENTER, bold=True)

    task_slots = {
        0: [34, 35, 31, 32, 36],
        1: [12, 13, 16, 14, 15],
        2: [17, 18, 20],
        3: [19, 21, 22, 23, 25],
        4: [24, 26, 27, 33],
    }
    for stage_idx, slots in task_slots.items():
        tasks = list(stages[stage_idx].get("tasks", []))
        while len(tasks) < len(slots):
            tasks.append("")
        for slot, task in zip(slots, tasks):
            _replace_text(_shape(slide, slot), _normalize(task, 18), color=(0, 0, 0), size_pt=10, align=PP_ALIGN.CENTER)

    legend = list(payload.get("legend", ["AI", "Python", "Template", "Designer"]))
    while len(legend) < 4:
        legend.append("")
    for slot, label in zip([28, 41, 29, 30], legend[:4]):
        _replace_text(_shape(slide, slot), label, color=(0, 0, 0), size_pt=10, align=PP_ALIGN.CENTER)


def _fill_pain_cards(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    _replace_text(_shape(slide, 35), payload.get("lead", page.subtitle), color=(0, 0, 0), size_pt=18, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 34), payload.get("bottom_banner", "先识别问题边界，再设计自动化路径。"), color=(255, 255, 255), size_pt=16, align=PP_ALIGN.CENTER, bold=True)

    cards = _cards(payload, "cards", 3)
    for slot, card in zip([5, 15, 25], cards):
        _replace_text(_shape(slide, slot), card.get("title", ""), color=COLOR_ACTIVE, size_pt=17, bold=True)
    for slot, card in zip([6, 16, 26], cards):
        _replace_text(_shape(slide, slot), card.get("detail", ""), color=(91, 112, 135), size_pt=12)
    bullet_slots = [[9, 10, 11], [19, 20, 21], [29, 30, 31]]
    for slots, card in zip(bullet_slots, cards):
        points = list(card.get("points", []))
        while len(points) < 3:
            points.append("")
        for slot, point in zip(slots, points):
            _replace_text(_shape(slide, slot), point, color=(91, 112, 135), size_pt=11)
