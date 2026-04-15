import gc
import shutil
import warnings
from pathlib import Path

from pptx import Presentation

from .generation_runtime import refresh_legacy_directory_clones, render_with_legacy_clone, render_with_preallocated_pool
from .generation_support import (
    apply_reference_body_slides,
    build_output_path,
    build_page_render_traces,
    build_preflight_notes,
    refresh_preallocated_directory_assets,
    reference_import_unavailable_reason,
    validate_slide_pool_configuration,
    warn_if_reference_import_disabled,
)
from .pipeline import plan_deck_from_html, plan_deck_from_json
from .presentation_ops import ensure_last_slide
from .body_renderers import apply_theme_title
from ..config import DEFAULT_MIN_TEMPLATE_SLIDES, DEFAULT_OUTPUT_DIR
from ..models import DeckPlan, DeckRenderTrace, GenerationArtifacts
from .reference_styles import populate_reference_body_pages
from ..template_manifest import load_template_manifest


def _generate_ppt_from_plan(
    deck_plan: DeckPlan,
    input_kind: str,
    template_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    manifest = load_template_manifest(template_path=template_path)
    deck = deck_plan.deck
    body_pages = deck.body_pages

    final_output_dir = output_dir or DEFAULT_OUTPUT_DIR
    out = build_output_path(final_output_dir, output_prefix)
    shutil.copy2(template_path, out)

    prs = Presentation(str(out))
    if len(prs.slides) < DEFAULT_MIN_TEMPLATE_SLIDES:
        raise ValueError(
            "Template slide count is below the minimum required for SIE rendering: "
            f"expected at least {DEFAULT_MIN_TEMPLATE_SLIDES}, got {len(prs.slides)}."
        )

    preflight_notes = build_preflight_notes(body_pages, manifest, len(prs.slides), reference_body_path)
    apply_theme_title(prs, deck.cover_title, manifest)
    if manifest.slide_pools and manifest.slide_pools.ending is not None and manifest.slide_pools.ending < len(prs.slides):
        thanks_slide_id = int(prs.slides._sldIdLst[manifest.slide_pools.ending].id)
    else:
        thanks_slide_id = int(prs.slides._sldIdLst[len(prs.slides) - 1].id)

    used_preallocated_pool = bool(manifest.slide_pools)
    if used_preallocated_pool:
        render_with_preallocated_pool(prs, body_pages, active_start, manifest)
    else:
        warnings.warn(
            "Template does not provide a preallocated slide pool; using deprecated legacy runtime cloning for this deck.",
            stacklevel=2,
        )
        render_with_legacy_clone(prs, body_pages, active_start, manifest)

    ensure_last_slide(prs, thanks_slide_id)
    prs.save(str(out))
    prs = None
    gc.collect()

    reference_import_reason = reference_import_unavailable_reason(body_pages, reference_body_path)
    warn_if_reference_import_disabled(body_pages, reference_body_path)
    reference_import_applied = apply_reference_body_slides(out, reference_body_path, body_pages, manifest)
    if reference_import_applied:
        populate_reference_body_pages(out, body_pages, manifest=manifest)
        reference_import_reason = ""
    elif any(page.reference_style_id for page in body_pages) and not reference_import_reason:
        reference_import_reason = "reference style import failed during native package merge"

    if used_preallocated_pool and not refresh_preallocated_directory_assets(out, len(body_pages), manifest):
        warnings.warn(
            "Preallocated directory slide assets could not be fully refreshed after save; manual review is recommended.",
            stacklevel=2,
        )

    if not used_preallocated_pool:
        refresh_legacy_directory_clones(out, body_pages, active_start, len(body_pages), manifest)

    render_trace = DeckRenderTrace(
        input_kind=input_kind,
        body_render_mode="preallocated_pool" if used_preallocated_pool else "legacy_clone",
        reference_import_applied=reference_import_applied,
        reference_import_reason=reference_import_reason,
        preflight_notes=preflight_notes,
        page_traces=build_page_render_traces(body_pages, reference_import_applied, reference_import_reason),
    )
    return GenerationArtifacts(
        output_path=out,
        deck_plan=deck_plan,
        render_trace=render_trace,
    )


def generate_ppt_artifacts_from_html(
    template_path: Path,
    html_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    chapters: int | None,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    if not html_path.exists():
        raise FileNotFoundError(f"HTML not found: {html_path}")
    deck_plan = plan_deck_from_html(html_path, chapters)
    return _generate_ppt_from_plan(
        deck_plan=deck_plan,
        input_kind="html",
        template_path=template_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        active_start=active_start,
        output_dir=output_dir,
    )


def generate_ppt_artifacts_from_deck_spec(
    template_path: Path,
    deck_spec_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    if not deck_spec_path.exists():
        raise FileNotFoundError(f"Deck JSON not found: {deck_spec_path}")
    deck_plan = plan_deck_from_json(deck_spec_path)
    return _generate_ppt_from_plan(
        deck_plan=deck_plan,
        input_kind="deck_spec_json",
        template_path=template_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        active_start=active_start,
        output_dir=output_dir,
    )


def generate_ppt_artifacts_from_deck_plan(
    deck_plan: DeckPlan,
    input_kind: str,
    template_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    return _generate_ppt_from_plan(
        deck_plan=deck_plan,
        input_kind=input_kind,
        template_path=template_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        active_start=active_start,
        output_dir=output_dir,
    )


def generate_ppt(
    template_path: Path,
    html_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    chapters: int | None,
    active_start: int,
    output_dir: Path | None = None,
):
    artifacts = generate_ppt_artifacts_from_html(
        template_path=template_path,
        html_path=html_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        chapters=chapters,
        active_start=active_start,
        output_dir=output_dir,
    )
    return artifacts.output_path, artifacts.deck_plan.pattern_ids, artifacts.deck_plan.chapter_lines


__all__ = [
    "build_output_path",
    "generate_ppt",
    "generate_ppt_artifacts_from_deck_plan",
    "generate_ppt_artifacts_from_deck_spec",
    "generate_ppt_artifacts_from_html",
    "validate_slide_pool_configuration",
]
