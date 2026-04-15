import gc
import warnings
from pathlib import Path

from pptx import Presentation

from .generation_support import LEGACY_CLONE_DEPRECATION_MESSAGE
from .openxml_slide_ops import copy_slide_xml_assets, slide_assets_preserved
from .presentation_ops import clone_slide_after, remove_slide
from .body_renderers import fill_body_slide, fill_directory_slide
from ..models import BodyPageSpec
from ..planning.deck_planner import build_directory_window
from ..template_manifest import TemplateManifest


def directory_lines_for_page(body_pages: list[BodyPageSpec], active_index: int) -> tuple[list[str], int]:
    return build_directory_window(body_pages, active_index)


def render_with_preallocated_pool(
    prs: Presentation,
    body_pages: list[BodyPageSpec],
    active_start: int,
    manifest: TemplateManifest,
):
    if not manifest.slide_pools:
        raise ValueError("Template manifest does not define slide pools.")

    used_directory_indices = list(manifest.slide_pools.directory[: len(body_pages)])
    used_body_indices = list(manifest.slide_pools.body[: len(body_pages)])
    unused_pairs = list(
        zip(
            manifest.slide_pools.directory[len(body_pages):],
            manifest.slide_pools.body[len(body_pages):],
            strict=False,
        )
    )

    for directory_idx, body_idx in sorted(unused_pairs, reverse=True):
        for slide_index in sorted((directory_idx, body_idx), reverse=True):
            if slide_index < len(prs.slides):
                remove_slide(prs, slide_index)

    directory_slides = [prs.slides[index] for index in used_directory_indices]
    body_slides = [prs.slides[index] for index in used_body_indices]

    for chapter_idx, directory_slide in enumerate(directory_slides):
        window_lines, window_active_index = directory_lines_for_page(body_pages, active_start + chapter_idx)
        fill_directory_slide(directory_slide, window_lines, window_active_index, manifest)
    for page, body_slide in zip(body_pages, body_slides):
        fill_body_slide(body_slide, page, manifest)


def render_with_legacy_clone(
    prs: Presentation,
    body_pages: list[BodyPageSpec],
    active_start: int,
    manifest: TemplateManifest,
):
    warnings.warn(LEGACY_CLONE_DEPRECATION_MESSAGE, DeprecationWarning, stacklevel=2)
    directory_idx = manifest.slide_roles.directory
    body_template_idx = manifest.slide_roles.body_template
    directory_slides = [prs.slides[directory_idx]]
    body_slides = [prs.slides[body_template_idx]]
    insert_after = body_template_idx
    for _ in body_pages[1:]:
        new_directory = clone_slide_after(prs, directory_idx, insert_after, keep_rel_ids=True)
        directory_slides.append(new_directory)
        insert_after += 1

        new_body = clone_slide_after(prs, body_template_idx, insert_after, keep_rel_ids=False)
        body_slides.append(new_body)
        insert_after += 1

    for chapter_idx, directory_slide in enumerate(directory_slides):
        window_lines, window_active_index = directory_lines_for_page(body_pages, active_start + chapter_idx)
        fill_directory_slide(directory_slide, window_lines, window_active_index, manifest)
    for page, body_slide in zip(body_pages, body_slides):
        fill_body_slide(body_slide, page, manifest)


def refresh_legacy_directory_clones(
    pptx_path: Path,
    body_pages: list[BodyPageSpec],
    active_start: int,
    body_page_count: int,
    manifest: TemplateManifest,
):
    warnings.warn(LEGACY_CLONE_DEPRECATION_MESSAGE, DeprecationWarning, stacklevel=2)
    directory_idx = manifest.slide_roles.directory
    targets = [directory_idx + 1 + i * 2 for i in range(1, body_page_count)]
    if not targets:
        return True

    source_idx = directory_idx + 1
    for _ in range(3):
        if not copy_slide_xml_assets(pptx_path, source_idx=source_idx, target_indices=targets):
            continue

        prs_reloaded = Presentation(str(pptx_path))
        window_lines, window_active_index = directory_lines_for_page(body_pages, active_start)
        fill_directory_slide(prs_reloaded.slides[directory_idx], window_lines, window_active_index, manifest)
        for offset, directory_slide_no in enumerate(targets, start=1):
            slide_index = directory_slide_no - 1
            if slide_index < len(prs_reloaded.slides):
                window_lines, window_active_index = directory_lines_for_page(body_pages, active_start + offset)
                fill_directory_slide(prs_reloaded.slides[slide_index], window_lines, window_active_index, manifest)
        prs_reloaded.save(str(pptx_path))
        prs_reloaded = None
        gc.collect()

        if slide_assets_preserved(pptx_path, source_idx=source_idx, target_indices=targets):
            return True

    raise RuntimeError(
        "Legacy clone directory-slide asset refresh failed after 3 retries. "
        f"Source slide index: {source_idx}; target slide indices: {targets}. "
        "Please migrate this template to the preallocated slide-pool path."
    )


__all__ = [
    "directory_lines_for_page",
    "refresh_legacy_directory_clones",
    "render_with_legacy_clone",
    "render_with_preallocated_pool",
]
