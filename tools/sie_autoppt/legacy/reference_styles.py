import re
import unicodedata
import zipfile
from dataclasses import replace
from functools import lru_cache
from pathlib import Path
from typing import Any, cast
from xml.etree import ElementTree

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from ..config import COLOR_ACTIVE
from ..models import validate_body_page_payload
from ..template_manifest import TemplateManifest, load_template_manifest
from ..text_ops import write_text

PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


REFERENCE_STYLE_LIBRARY = {
    "comparison_upgrade": {
        "name": "左右对比升级页",
        "source_slide_name": "comparison_upgrade_reference",
        "source_slide": 5,
        "match_any_text": ["价值：业务效能跃升", "传统人工追溯模式"],
        "scenes": ["价值对比", "旧方式 vs 新方式", "能力边界与工程化解法"],
    },
    "capability_ring": {
        "name": "能力亮点环形页",
        "source_slide_name": "capability_ring_reference",
        "source_slide": 16,
        "match_any_text": ["赛意追溯产品亮点", "AI智能识别"],
        "scenes": ["能力亮点", "实践原则", "方案卖点"],
    },
    "five_phase_path": {
        "name": "五阶段推进路径页",
        "source_slide_name": "five_phase_path_reference",
        "source_slide": 20,
        "match_any_text": ["追溯管理-外部追溯推进路径", "阶段五 数据应用"],
        "scenes": ["实施路径", "方法路线", "阶段推进"],
    },
    "pain_cards": {
        "name": "三栏痛点拆解页",
        "source_slide_name": "pain_cards_reference",
        "source_slide": 6,
        "match_any_text": ["痛点：标准不清晰", "链路难贯通", "组织协同慢"],
        "scenes": ["痛点拆解", "挑战分析", "风险拆解"],
    },
}


def get_reference_slide_no(style_id: str | None) -> int | None:
    if not style_id:
        return None
    style = REFERENCE_STYLE_LIBRARY.get(style_id)
    if not style:
        return None
    return int(cast(int, style["source_slide"]))


def _normalize_reference_text(text: str) -> str:
    normalized = unicodedata.normalize("NFKC", text).lower()
    normalized = re.sub(r"\s+", "", normalized)
    return normalized


@lru_cache(maxsize=8)
def _slide_text_index(reference_body_path_str: str) -> tuple[str, ...]:
    prs = Presentation(reference_body_path_str)
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = " ".join(shape.text_frame.text.split())
                if text:
                    texts.append(text)
        slide_texts.append(_normalize_reference_text(" ".join(texts)))
    return tuple(slide_texts)


@lru_cache(maxsize=8)
def _slide_metadata_name_index(reference_body_path_str: str) -> tuple[str, ...]:
    prs = Presentation(reference_body_path_str)
    metadata_names = []
    with zipfile.ZipFile(reference_body_path_str) as package:
        for slide in prs.slides:
            part_name = str(slide.part.partname).lstrip("/")
            metadata_name = ""
            try:
                root = ElementTree.fromstring(package.read(part_name))
            except KeyError:
                metadata_names.append(metadata_name)
                continue
            c_sld = root.find(f"{{{PRESENTATION_NS}}}cSld")
            if c_sld is not None:
                metadata_name = _normalize_reference_text(c_sld.attrib.get("name", ""))
            metadata_names.append(metadata_name)
    return tuple(metadata_names)


def locate_reference_slide_no(style_id: str | None, reference_body_path: Path | None) -> int | None:
    if not style_id:
        return None
    style = REFERENCE_STYLE_LIBRARY.get(style_id)
    if not style:
        return None
    if reference_body_path and reference_body_path.exists():
        metadata_target = _normalize_reference_text(str(style.get("source_slide_name", "")))
        if metadata_target:
            metadata_names = _slide_metadata_name_index(str(reference_body_path.resolve()))
            for index, metadata_name in enumerate(metadata_names, start=1):
                if metadata_name == metadata_target:
                    return index
        slide_texts = _slide_text_index(str(reference_body_path.resolve()))
        match_any_text = cast(list[Any], style.get("match_any_text", []))
        for marker in match_any_text:
            normalized_marker = _normalize_reference_text(str(marker))
            for index, slide_text in enumerate(slide_texts, start=1):
                if normalized_marker and normalized_marker in slide_text:
                    return index
    return get_reference_slide_no(style_id)


def build_reference_import_plan(
    body_pages,
    reference_body_path: Path | None = None,
    manifest: TemplateManifest | None = None,
) -> list[tuple[int, int]]:
    manifest = manifest or load_template_manifest()
    plan = []
    for offset, page in enumerate(body_pages):
        source_slide = locate_reference_slide_no(getattr(page, "reference_style_id", None), reference_body_path)
        if source_slide is None:
            continue
        target_slide = manifest.slide_roles.body_template + 1 + offset * 2
        plan.append((target_slide, source_slide))
    return plan


def _replace_text(shape, text: str, color=None, size_pt: int | None = None, align=PP_ALIGN.LEFT, bold=None):
    write_text(
        shape,
        text,
        color=color,
        size_pt=size_pt,
        bold=bold,
        align=align,
        preserve_runs=False,
    )


def _shape(slide, index: int):
    return slide.shapes[index - 1]


def _has_shape_slots(slide, max_index: int) -> bool:
    return len(slide.shapes) >= max_index


def _normalize(text: str, max_chars: int) -> str:
    compact = re.sub(r"\s+", " ", text).strip()
    if len(compact) <= max_chars:
        return compact
    return compact[: max_chars - 1].rstrip(" ,;，；。") + "…"


def _cards(payload, key: str, count: int) -> list[dict[str, str]]:
    cards = list(payload.get(key, []))
    while len(cards) < count:
        cards.append({"title": "", "detail": ""})
    return cards[:count]


def fill_reference_style_slide(slide, page) -> bool:
    validated_payload = validate_body_page_payload(getattr(page, "pattern_id", ""), getattr(page, "payload", {}))
    if hasattr(validated_payload, "model_dump"):
        payload = validated_payload.model_dump(mode="json")
    else:
        payload = dict(validated_payload)
    page = replace(page, payload=payload)
    style_id = getattr(page, "reference_style_id", None)
    if style_id == "comparison_upgrade":
        if not _has_shape_slots(slide, 61):
            return False
        _fill_comparison_upgrade(slide, page)
        return True
    if style_id == "capability_ring":
        if not _has_shape_slots(slide, 31):
            return False
        _fill_capability_ring(slide, page)
        return True
    if style_id == "five_phase_path":
        if not _has_shape_slots(slide, 60):
            return False
        _fill_five_phase_path(slide, page)
        return True
    if style_id == "pain_cards":
        if not _has_shape_slots(slide, 35):
            return False
        _fill_pain_cards(slide, page)
        return True
    return False


def populate_reference_body_pages(pptx_path: Path, body_pages, manifest: TemplateManifest | None = None) -> bool:
    manifest = manifest or load_template_manifest()
    has_reference_pages = any(getattr(page, "reference_style_id", None) for page in body_pages)
    if not has_reference_pages:
        return False
    prs = Presentation(str(pptx_path))
    for offset, page in enumerate(body_pages):
        if not getattr(page, "reference_style_id", None):
            continue
        slide_index = manifest.slide_roles.body_template + offset * 2
        if slide_index >= len(prs.slides):
            continue
        fill_reference_style_slide(prs.slides[slide_index], page)
    prs.save(str(pptx_path))
    return True


def _fill_comparison_upgrade(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    _replace_text(_shape(slide, 56), payload.get("headline", page.subtitle), color=(0, 0, 0), size_pt=18, align=PP_ALIGN.CENTER, bold=True)

    _replace_text(_shape(slide, 2), payload.get("left_label", "● 传统方式"), color=(91, 112, 135), size_pt=17, bold=True)
    _replace_text(_shape(slide, 15), payload.get("right_label", "● 工程化方式"), color=COLOR_ACTIVE, size_pt=17, bold=True)

    for slot, card in zip([4, 7, 10, 13], _cards(payload, "left_cards", 4)):
        _replace_text(_shape(slide, slot), card.get("title", ""), color=(60, 76, 96), size_pt=16, bold=True)
    for slot, card in zip([5, 8, 11, 14], _cards(payload, "left_cards", 4)):
        detail = _cards(payload, "left_cards", 4)[[5, 8, 11, 14].index(slot)].get("detail", "")
        _replace_text(_shape(slide, slot), detail, color=(91, 112, 135), size_pt=10)

    for slot, card in zip([17, 20, 23, 26], _cards(payload, "right_cards", 4)):
        _replace_text(_shape(slide, slot), card.get("title", ""), color=(255, 255, 255), size_pt=16, bold=True)
    for slot, card in zip([18, 21, 24, 27], _cards(payload, "right_cards", 4)):
        detail = _cards(payload, "right_cards", 4)[[18, 21, 24, 27].index(slot)].get("detail", "")
        _replace_text(_shape(slide, slot), detail, color=(255, 255, 255), size_pt=10)

    _replace_text(_shape(slide, 30), payload.get("center_kicker", "ENGINEERING COPILOT"), color=(255, 255, 255), size_pt=12, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 31), payload.get("center_title", "结构化内容交付中枢"), color=COLOR_ACTIVE, size_pt=24, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 32), payload.get("center_subtitle", ""), color=(91, 112, 135), size_pt=12, align=PP_ALIGN.CENTER)
    _replace_text(_shape(slide, 35), payload.get("center_top_label", "协同原则层"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 36), payload.get("center_section_title", "人机协同 (Human in the Loop)"), color=(0, 0, 0), size_pt=18, bold=True)
    _replace_text(_shape(slide, 38), payload.get("center_row1_left", "结构化输入"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 40), payload.get("center_row1_right", "模板化渲染"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 42), payload.get("center_row2_left", "视觉 QA"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 44), payload.get("center_row2_right", "人工微调"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 45), payload.get("center_divider", "♦ 保持骨架稳定，保留创作空间 ♦"), color=COLOR_ACTIVE, size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 48), payload.get("center_bottom_label", "交付结果层"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 49), payload.get("center_bottom_title", "高质量输出 (Deliverable)"), color=(0, 0, 0), size_pt=18, bold=True)
    _replace_text(_shape(slide, 51), payload.get("center_bottom_left", "稳定复用"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 53), payload.get("center_bottom_right", "风格一致"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 55), payload.get("center_bottom_footer", "效率与质量兼得"), color=(60, 76, 96), size_pt=11, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 60), payload.get("bottom_left_caption", "能力边界"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 61), payload.get("bottom_right_caption", "工程化交付"), color=(255, 255, 255), size_pt=10, align=PP_ALIGN.CENTER, bold=True)


def _fill_capability_ring(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    items = list(payload.get("items", []))
    while len(items) < 7:
        items.append({"title": "", "detail": ""})
    for title_idx, detail_idx, item in zip([18, 20, 22, 24, 26, 28, 30], [19, 21, 23, 25, 27, 29, 31], items[:7]):
        _replace_text(_shape(slide, title_idx), _normalize(item.get("title", ""), 14), color=(255, 255, 255), size_pt=15, align=PP_ALIGN.CENTER, bold=True)
        _replace_text(_shape(slide, detail_idx), _normalize(item.get("detail", ""), 40), color=(60, 76, 96), size_pt=10, align=PP_ALIGN.CENTER)


def _fill_five_phase_path(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    _replace_text(_shape(slide, 60), payload.get("intro", page.subtitle), color=(0, 0, 0), size_pt=16, bold=True)

    stages = list(payload.get("stages", []))
    while len(stages) < 5:
        stages.append({"header": "", "tasks": []})
    for slot, stage in zip([6, 4, 2, 8, 10], stages[:5]):
        _replace_text(_shape(slide, slot), stage.get("header", ""), color=(255, 255, 255), size_pt=14, align=PP_ALIGN.CENTER, bold=True)

    task_slots = {
        0: [34, 35, 31, 32, 36],
        1: [12, 13, 16, 14, 15],
        2: [17, 18, 20],
        3: [19, 21, 22, 23, 25],
        4: [24, 26, 27, 33],
    }
    for stage_idx, slots in task_slots.items():
        tasks = list(stages[stage_idx].get("tasks", []))
        while len(tasks) < len(slots):
            tasks.append("")
        for slot, task in zip(slots, tasks):
            _replace_text(_shape(slide, slot), _normalize(task, 18), color=(0, 0, 0), size_pt=10, align=PP_ALIGN.CENTER)

    legend = list(payload.get("legend", ["AI", "Python", "Template", "Designer"]))
    while len(legend) < 4:
        legend.append("")
    for slot, label in zip([28, 41, 29, 30], legend[:4]):
        _replace_text(_shape(slide, slot), label, color=(0, 0, 0), size_pt=10, align=PP_ALIGN.CENTER)


def _fill_pain_cards(slide, page):
    payload = page.payload
    _replace_text(_shape(slide, 1), page.title, color=COLOR_ACTIVE, size_pt=30, bold=True)
    _replace_text(_shape(slide, 35), payload.get("lead", page.subtitle), color=(0, 0, 0), size_pt=18, align=PP_ALIGN.CENTER, bold=True)
    _replace_text(_shape(slide, 34), payload.get("bottom_banner", "先识别问题边界，再设计自动化路径。"), color=(255, 255, 255), size_pt=16, align=PP_ALIGN.CENTER, bold=True)

    cards = _cards(payload, "cards", 3)
    for slot, card in zip([5, 15, 25], cards):
        _replace_text(_shape(slide, slot), card.get("title", ""), color=COLOR_ACTIVE, size_pt=17, bold=True)
    for slot, card in zip([6, 16, 26], cards):
        _replace_text(_shape(slide, slot), card.get("detail", ""), color=(91, 112, 135), size_pt=12)
    bullet_slots = [[9, 10, 11], [19, 20, 21], [29, 30, 31]]
    for slots, card in zip(bullet_slots, cards):
        points = list(card.get("points", []))
        while len(points) < 3:
            points.append("")
        for slot, point in zip(slots, points):
            _replace_text(_shape(slide, slot), point, color=(91, 112, 135), size_pt=11)
