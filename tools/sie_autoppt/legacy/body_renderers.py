"""
Legacy V1 body slide renderers.

This module is still used by the classic generator path and scenario scripts.
New layout work should prefer the V2 semantic renderers under tools/sie_autoppt/v2/renderers/.
"""

import re
import textwrap
from dataclasses import replace

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from ..config import COLOR_ACTIVE, COLOR_INACTIVE
from ..models import BodyPageSpec, validate_body_page_payload
from ..style_guide import deep_merge_dict
from ..template_manifest import TemplateManifest
from ..text_ops import add_textbox, write_text


class _ManifestValue:
    def __init__(self, value, context: str):
        self._value = value
        self._context = context

    def __getitem__(self, key):
        if isinstance(self._value, dict):
            if key not in self._value:
                raise KeyError(f"[渲染器错误] manifest 缺少必要字段 '{key}'（位置：{self._context}）")
            return _wrap_manifest_value(self._value[key], f"{self._context}.{key}")
        if isinstance(self._value, (list, tuple)):
            try:
                value = self._value[key]
            except (IndexError, TypeError) as exc:
                raise KeyError(f"[渲染器错误] manifest 列表索引无效 '{key}'（位置：{self._context}）") from exc
            return _wrap_manifest_value(value, f"{self._context}[{key}]")
        raise TypeError(f"[渲染器错误] manifest 字段不支持索引读取（位置：{self._context}）")

    def get(self, key, default=None):
        if not isinstance(self._value, dict):
            raise TypeError(f"[渲染器错误] manifest 字段不支持 get() 调用（位置：{self._context}）")
        if key not in self._value:
            return default
        return _wrap_manifest_value(self._value[key], f"{self._context}.{key}")

    def __iter__(self):
        if isinstance(self._value, (list, tuple)):
            for index, item in enumerate(self._value):
                yield _wrap_manifest_value(item, f"{self._context}[{index}]")
            return
        raise TypeError(f"[渲染器错误] manifest 字段不可迭代（位置：{self._context}）")

    def __len__(self):
        return len(self._value)

    def __bool__(self):
        return bool(self._value)

    def __int__(self):
        try:
            return int(self._value)
        except (TypeError, ValueError) as exc:
            raise ValueError(
                f"[渲染器错误] manifest 字段无法转换为整数（位置：{self._context}，值：{self._value!r}）"
            ) from exc

    def __float__(self):
        try:
            return float(self._value)
        except (TypeError, ValueError) as exc:
            raise ValueError(
                f"[渲染器错误] manifest 字段无法转换为浮点数（位置：{self._context}，值：{self._value!r}）"
            ) from exc

    def __str__(self):
        return str(self._value)

    def __repr__(self):
        return repr(self._value)


def _wrap_manifest_value(value, context: str):
    if isinstance(value, _ManifestValue):
        return value
    return _ManifestValue(value, context)


def normalize_text_for_body_box(text: str, max_chars: int = 44) -> str:
    compact = re.sub(r"\s+", " ", text).strip()
    lines = textwrap.wrap(compact, width=max_chars)
    return "\n".join(lines[:4])


def split_title_detail(text: str) -> tuple[str, str]:
    normalized = re.sub(r"\s+", " ", text).strip()
    for sep in ("：", ":", " - ", " | "):
        if sep in normalized:
            title, detail = normalized.split(sep, 1)
            return title.strip(), detail.strip()
    return normalized, normalized


def choose_title_font_size(text: str, default: int = 24) -> int:
    n = len(text)
    if n > 28:
        return 16
    if n > 22:
        return 18
    if n > 16:
        return 20
    return default


def choose_font_size_by_length(text: str, base: int = 13) -> int:
    n = len(text)
    if n > 120:
        return max(10, base - 3)
    if n > 90:
        return max(11, base - 2)
    if n > 70:
        return max(12, base - 1)
    return base


def pick_text_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def _rgb(values):
    return tuple(int(value) for value in values)


def _accent_rgb(manifest: TemplateManifest) -> tuple[int, int, int]:
    return manifest.style_guide.accent_rgb or COLOR_ACTIVE


def _inactive_rgb(manifest: TemplateManifest) -> tuple[int, int, int]:
    return manifest.style_guide.inactive_rgb or COLOR_INACTIVE


def _layout(manifest: TemplateManifest, name: str) -> dict[str, object]:
    return _wrap_manifest_value(manifest.render_layout(name), f"render_layouts.{name}")


def _resolve_layout_name(manifest: TemplateManifest, page: BodyPageSpec, default_name: str) -> str:
    if page.layout_variant and page.layout_variant in manifest.render_layouts:
        return page.layout_variant
    return default_name


def _layout_for_page(manifest: TemplateManifest, page: BodyPageSpec, default_name: str) -> dict[str, object]:
    resolved_name = _resolve_layout_name(manifest, page, default_name)
    if resolved_name == default_name:
        return _layout(manifest, default_name)

    merged_layout = deep_merge_dict(manifest.render_layout(default_name), manifest.render_layout(resolved_name))
    return _wrap_manifest_value(merged_layout, f"render_layouts.{default_name}+{resolved_name}")


def _extract_system_tags(bullets: list[str]) -> list[str]:
    systems = []
    for bullet in bullets:
        for system in re.findall(r"\b[A-Z][A-Z0-9/-]{1,}\b", bullet):
            if system not in systems:
                systems.append(system)
    return systems[:5]


def fill_directory_slide(slide, chapter_lines, active_chapter_index: int, manifest: TemplateManifest):
    texts = sorted(pick_text_shapes(slide), key=lambda shape: (shape.top, shape.left))
    title_boxes = [shape for shape in texts if manifest.selectors.directory_title.matches(shape)]
    title_boxes = sorted(title_boxes, key=lambda shape: (shape.top, shape.left))
    safe_active_index = max(0, min(active_chapter_index, len(chapter_lines) - 1))
    accent_rgb = _accent_rgb(manifest)
    inactive_rgb = _inactive_rgb(manifest)
    for i, shape in enumerate(title_boxes[: len(chapter_lines)]):
        color = accent_rgb if i == safe_active_index else inactive_rgb
        write_text(
            shape,
            chapter_lines[i],
            color=color,
            size_pt=int(manifest.fonts.directory_title_pt),
            preserve_runs=True,
        )


def _clear_body_render_area(slide, manifest: TemplateManifest, protected_shapes=None):
    protected_shapes = list(protected_shapes or [])
    removable = []
    for shape in slide.shapes:
        if any(shape is protected for protected in protected_shapes):
            continue
        if manifest.selectors.body_render_area.matches(shape):
            removable.append(shape)
    for shape in removable:
        element = shape._element
        element.getparent().remove(element)


def apply_theme_title(prs, title: str, manifest: TemplateManifest):
    theme_texts = pick_text_shapes(prs.slides[manifest.slide_roles.theme])
    title_candidates = [shape for shape in theme_texts if manifest.selectors.theme_title.matches(shape)]
    if title_candidates:
        main_title = max(title_candidates, key=lambda shape: shape.width)
        write_text(
            main_title,
            title,
            color=_accent_rgb(manifest),
            size_pt=int(manifest.fonts.theme_title_pt),
            preserve_runs=True,
        )


def fill_body_slide(slide, page: BodyPageSpec, manifest: TemplateManifest):
    validated_payload = validate_body_page_payload(page.pattern_id, page.payload)
    if hasattr(validated_payload, "model_dump"):
        effective_payload = validated_payload.model_dump(mode="json")
    else:
        effective_payload = validated_payload
    page = replace(page, payload=effective_payload)

    texts = sorted(pick_text_shapes(slide), key=lambda shape: (shape.top, shape.left))
    title_candidates = [shape for shape in texts if manifest.selectors.body_title.matches(shape)]
    if title_candidates:
        write_text(
            title_candidates[0],
            page.title,
            color=_accent_rgb(manifest),
            size_pt=choose_title_font_size(page.title),
            preserve_runs=True,
        )
    else:
        add_textbox(
            slide,
            manifest.fallback_boxes.body_title.left,
            manifest.fallback_boxes.body_title.top,
            manifest.fallback_boxes.body_title.width,
            manifest.fallback_boxes.body_title.height,
            page.title,
            color=_accent_rgb(manifest),
            size_pt=choose_title_font_size(page.title),
            bold=True,
        )

    subtitle_candidates = [
        shape
        for shape in texts
        if manifest.selectors.body_subtitle.matches(shape) and shape not in title_candidates
    ]
    if subtitle_candidates:
        write_text(subtitle_candidates[0], page.subtitle, preserve_runs=True)

    protected_shapes = title_candidates + subtitle_candidates[:1]
    _clear_body_render_area(slide, manifest, protected_shapes=protected_shapes)
    renderer = PATTERN_RENDERERS[resolve_render_pattern(page.pattern_id)]
    renderer(slide, page, manifest)


def _render_cards_2x2(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout_for_page(manifest, page, "general_business")
    x0 = int(spec["origin_left"])
    y0 = int(spec["origin_top"])
    card_w = int(spec["card_width"])
    card_h = int(spec["card_height"])
    columns = max(1, int(spec.get("columns", 2)))
    gap_x = int(spec["gap_x"])
    gap_y = int(spec["gap_y"])
    padding = int(spec["textbox_padding"])
    fill_rgb = _rgb(spec["fill_rgb"])
    line_rgb = _rgb(spec["line_rgb"])
    wrap_chars = int(spec["wrap_chars"])
    base_font_pt = int(spec["base_font_pt"])
    max_items = int(spec.get("max_items", 4))

    for i, text in enumerate(page.bullets[:max_items]):
        row, col = divmod(i, columns)
        left = x0 + col * (card_w + gap_x)
        top = y0 + row * (card_h + gap_y)
        card = slide.shapes.add_shape(1, left, top, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*fill_rgb)
        card.line.color.rgb = RGBColor(*line_rgb)
        safe_text = normalize_text_for_body_box(text, wrap_chars)
        font_size = choose_font_size_by_length(safe_text, base_font_pt)
        add_textbox(
            slide,
            left + padding,
            top + padding,
            card_w - padding * 2,
            card_h - padding * 2,
            safe_text,
            size_pt=font_size,
        )


def _render_process_flow(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout_for_page(manifest, page, "process_flow")
    start_x = int(spec["origin_left"])
    start_y = int(spec["origin_top"])
    step_w = int(spec["step_width"])
    step_h = int(spec["step_height"])
    columns = max(1, int(spec.get("columns", int(spec.get("max_items", 4)))))
    gap_x = int(spec["gap_x"])
    gap_y = int(spec.get("gap_y", 0))
    fill_rgb = _rgb(spec["fill_rgb"])
    line_rgb = _rgb(spec["line_rgb"])
    number_box = spec["number_box"]
    title_box = spec.get("title_box", {})
    detail_box = spec.get("detail_box")
    if detail_box is None:
        detail_box = spec["text_box"]
    default_title_left = int(number_box["left_offset"]) + int(number_box["width"]) + 100000
    steps = list(page.payload.get("steps", []))
    if not steps:
        for index, text in enumerate(page.bullets[: int(spec.get("max_items", 4))], start=1):
            title, detail = split_title_detail(text)
            steps.append({"number": f"{index:02d}", "title": title, "detail": detail})

    for i, step in enumerate(steps[: int(spec.get("max_items", 4))]):
        row, col = divmod(i, columns)
        left = start_x + col * (step_w + gap_x)
        top = start_y + row * (step_h + gap_y)
        box = slide.shapes.add_shape(1, left, top, step_w, step_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*fill_rgb)
        box.line.color.rgb = RGBColor(*line_rgb)
        add_textbox(
            slide,
            left + int(number_box["left_offset"]),
            top + int(number_box["top_offset"]),
            int(number_box["width"]),
            int(number_box["height"]),
            str(step.get("number", f"{i + 1:02d}")),
            size_pt=int(number_box["font_pt"]),
            bold=True,
            align=PP_ALIGN.CENTER,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )
        add_textbox(
            slide,
            left + int(title_box.get("left_offset", default_title_left)),
            top + int(title_box.get("top_offset", number_box["top_offset"])),
            step_w - int(title_box.get("width_padding", detail_box["width_padding"])),
            int(title_box.get("height", 220000)),
            str(step.get("title", "")),
            size_pt=int(title_box.get("font_pt", 12)),
            bold=True,
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )
        safe_text = normalize_text_for_body_box(str(step.get("detail", "")), int(detail_box["wrap_chars"]))
        add_textbox(
            slide,
            left + int(detail_box["left_offset"]),
            top + int(detail_box["top_offset"]),
            step_w - int(detail_box["width_padding"]),
            step_h - int(detail_box["height_padding"]),
            safe_text,
            size_pt=int(detail_box["font_pt"]),
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )


def _render_architecture_layers(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "solution_architecture")
    frame = spec["frame"]
    tag_box = spec["tag_box"]
    tag_text = spec["tag_text"]
    body_text = spec["body_text"]
    palette = spec["palette"]

    layer_items = list(page.payload.get("layers", []))
    if layer_items:
        max_layers = int(spec.get("max_items", 4))
        layer_texts = [f"{item.get('title', '')}：{item.get('detail', '')}".strip("：") for item in layer_items[:max_layers]]
        layer_labels = [str(item.get("label", f"L{i + 1:02d}")) for i, item in enumerate(layer_items[:max_layers])]
    else:
        max_layers = int(spec.get("max_items", 4))
        layer_texts = page.bullets[:max_layers]
        layer_labels = [f"L{i + 1:02d}" for i in range(len(layer_texts))]

    for i, text in enumerate(layer_texts):
        y = int(frame["top"]) + i * (int(frame["layer_height"]) + int(frame["gap_y"]))
        left = int(frame["left"])
        width = int(frame["width"])
        layer_h = int(frame["layer_height"])
        colors = palette[i % len(palette)]

        layer = slide.shapes.add_shape(1, left, y, width, layer_h)
        layer.fill.solid()
        layer.fill.fore_color.rgb = RGBColor(*_rgb(colors["fill_rgb"]))
        layer.line.color.rgb = RGBColor(*_rgb(colors["line_rgb"]))

        tag = slide.shapes.add_shape(
            1,
            left + int(tag_box["left_offset"]),
            y + int(tag_box["top_offset"]),
            int(tag_box["width"]),
            layer_h - int(tag_box["height_padding"]),
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(*_rgb(colors["accent_rgb"]))
        tag.line.color.rgb = RGBColor(*_rgb(colors["accent_rgb"]))
        add_textbox(
            slide,
            left + int(tag_text["left_offset"]),
            y + int(tag_text["top_offset"]),
            int(tag_text["width"]),
            int(tag_text["height"]),
            layer_labels[i],
            size_pt=int(tag_text["font_pt"]),
            bold=True,
            color=(255, 255, 255),
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )

        safe_text = normalize_text_for_body_box(text, int(body_text["wrap_chars"]))
        add_textbox(
            slide,
            left + int(body_text["left_offset"]),
            y + int(body_text["top_offset"]),
            width - int(body_text["width_padding"]),
            layer_h - int(body_text["height_padding"]),
            safe_text,
            size_pt=int(body_text["font_pt"]),
        )

    banner = spec["banner"]
    banner_shape = slide.shapes.add_shape(
        1,
        int(banner["left"]),
        int(banner["top"]),
        int(banner["width"]),
        int(banner["height"]),
    )
    banner_shape.fill.solid()
    banner_shape.fill.fore_color.rgb = RGBColor(*_rgb(banner["fill_rgb"]))
    banner_shape.line.color.rgb = RGBColor(*_rgb(banner["fill_rgb"]))
    add_textbox(
        slide,
        int(banner["text_left"]),
        int(banner["text_top"]),
        int(banner["text_width"]),
        int(banner["text_height"]),
        str(page.payload.get("banner_text", banner["text"])),
        size_pt=int(banner["font_pt"]),
        bold=True,
        color=_rgb(banner["text_rgb"]),
    )

    chip_spec = spec["chips"]
    system_tags = _extract_system_tags(page.bullets)
    if system_tags:
        for idx, system in enumerate(system_tags[: int(chip_spec.get("max_items", 4))]):
            left = int(chip_spec["start_left"]) + idx * (int(chip_spec["width"]) + int(chip_spec["gap_x"]))
            chip = slide.shapes.add_shape(
                1,
                left,
                int(chip_spec["top"]),
                int(chip_spec["width"]),
                int(chip_spec["height"]),
            )
            chip.fill.solid()
            chip.fill.fore_color.rgb = RGBColor(*_rgb(chip_spec["fill_rgb"]))
            chip.line.color.rgb = RGBColor(*_rgb(chip_spec["line_rgb"]))
            add_textbox(
                slide,
                left,
                int(chip_spec["top"]) + int(chip_spec["text_top_offset"]),
                int(chip_spec["width"]),
                int(chip_spec["height"]),
                system,
                size_pt=int(chip_spec["font_pt"]),
                bold=True,
                align=PP_ALIGN.CENTER,
            )


def _render_governance_grid(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout_for_page(manifest, page, "org_governance")
    grid = spec["grid"]
    label_box = spec["label_box"]
    label_text = spec["label_text"]
    body_text = spec["body_text"]

    cards = list(page.payload.get("cards", []))
    if not cards:
        cards = []
        max_items = int(spec.get("max_items", 4))
        for index, text in enumerate(page.bullets[:max_items], start=1):
            title, detail = split_title_detail(text)
            cards.append({"label": title or f"{page.payload.get('label_prefix', label_text['label_prefix'])} {index}", "detail": detail})

    max_cards = int(spec.get("max_items", 4))
    for i, card_data in enumerate(cards[:max_cards]):
        row, col = divmod(i, 2)
        left = int(grid["left"]) + col * (int(grid["card_width"]) + int(grid["gap_x"]))
        top = int(grid["top"]) + row * (int(grid["card_height"]) + int(grid["gap_y"]))
        card = slide.shapes.add_shape(1, left, top, int(grid["card_width"]), int(grid["card_height"]))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["card_fill_rgb"]))
        card.line.color.rgb = RGBColor(*_rgb(spec["card_line_rgb"]))

        label = slide.shapes.add_shape(
            1,
            left + int(label_box["left_offset"]),
            top + int(label_box["top_offset"]),
            int(label_box["width"]),
            int(label_box["height"]),
        )
        label.fill.solid()
        label.fill.fore_color.rgb = RGBColor(*_rgb(label_box["fill_rgb"]))
        label.line.color.rgb = RGBColor(*_rgb(label_box["line_rgb"]))
        label_prefix = str(page.payload.get("label_prefix", label_text["label_prefix"]))
        add_textbox(
            slide,
            left + int(label_text["left_offset"]),
            top + int(label_text["top_offset"]),
            int(label_text["width"]),
            int(label_text["height"]),
            str(card_data.get("label", f"{label_prefix} {i + 1}")),
            size_pt=int(label_text["font_pt"]),
            bold=True,
            color=_rgb(label_text["text_rgb"]),
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )

        safe_text = normalize_text_for_body_box(str(card_data.get("detail", "")), int(body_text["wrap_chars"]))
        add_textbox(
            slide,
            left + int(body_text["left_offset"]),
            top + int(body_text["top_offset"]),
            int(grid["card_width"]) - int(body_text["width_padding"]),
            int(grid["card_height"]) - int(body_text["height_padding"]),
            safe_text,
            size_pt=int(body_text["font_pt"]),
            margin_left=0,
            margin_right=0,
            margin_top=0,
            margin_bottom=0,
        )

    footer = spec["footer_bar"]
    footer_bar = slide.shapes.add_shape(
        1,
        int(footer["left"]),
        int(footer["top"]),
        int(footer["width"]),
        int(footer["height"]),
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(*_rgb(footer["fill_rgb"]))
    footer_bar.line.color.rgb = RGBColor(*_rgb(footer["line_rgb"]))
    add_textbox(
        slide,
        int(footer["text_left"]),
        int(footer["text_top"]),
        int(footer["text_width"]),
        int(footer["text_height"]),
        str(page.payload.get("footer_text", page.subtitle or footer["text"])),
        size_pt=int(footer["font_pt"]),
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def _render_comparison_upgrade(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "comparison_upgrade")
    payload = page.payload
    left_panel = spec["left_panel"]
    right_panel = spec["right_panel"]
    center_panel = spec["center_panel"]
    card_spec = spec["card"]

    for panel, label_key, cards_key, title_color in (
        (left_panel, "left_label", "left_cards", _rgb(left_panel["title_rgb"])),
        (right_panel, "right_label", "right_cards", _rgb(right_panel["title_rgb"])),
    ):
        add_textbox(
            slide,
            int(panel["left"]),
            int(panel["top"]) - int(panel["header_gap_y"]),
            int(panel["width"]),
            int(panel["header_height"]),
            str(payload.get(label_key, "")),
            size_pt=int(panel["header_font_pt"]),
            bold=True,
            color=title_color,
            align=PP_ALIGN.CENTER,
        )
        cards = list(payload.get(cards_key, []))[: int(panel["max_cards"])]
        for idx, card in enumerate(cards):
            top = int(panel["top"]) + idx * (int(card_spec["height"]) + int(card_spec["gap_y"]))
            shape = slide.shapes.add_shape(1, int(panel["left"]), top, int(panel["width"]), int(card_spec["height"]))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*_rgb(panel["fill_rgb"]))
            shape.line.color.rgb = RGBColor(*_rgb(panel["line_rgb"]))
            add_textbox(
                slide,
                int(panel["left"]) + int(card_spec["padding_x"]),
                top + int(card_spec["title_top_offset"]),
                int(panel["width"]) - int(card_spec["padding_x"]) * 2,
                int(card_spec["title_height"]),
                str(card.get("title", "")),
                size_pt=int(card_spec["title_font_pt"]),
                bold=True,
                color=title_color,
            )
            add_textbox(
                slide,
                int(panel["left"]) + int(card_spec["padding_x"]),
                top + int(card_spec["detail_top_offset"]),
                int(panel["width"]) - int(card_spec["padding_x"]) * 2,
                int(card_spec["detail_height"]),
                normalize_text_for_body_box(str(card.get("detail", "")), int(card_spec["detail_wrap_chars"])),
                size_pt=int(card_spec["detail_font_pt"]),
                color=_rgb(panel["detail_rgb"]),
            )

    center_shape = slide.shapes.add_shape(
        1,
        int(center_panel["left"]),
        int(center_panel["top"]),
        int(center_panel["width"]),
        int(center_panel["height"]),
    )
    center_shape.fill.solid()
    center_shape.fill.fore_color.rgb = RGBColor(*_rgb(center_panel["fill_rgb"]))
    center_shape.line.color.rgb = RGBColor(*_rgb(center_panel["line_rgb"]))

    kicker = spec["center_kicker"]
    kicker_shape = slide.shapes.add_shape(
        1,
        int(kicker["left"]),
        int(kicker["top"]),
        int(kicker["width"]),
        int(kicker["height"]),
    )
    kicker_shape.fill.solid()
    kicker_shape.fill.fore_color.rgb = RGBColor(*_rgb(kicker["fill_rgb"]))
    kicker_shape.line.color.rgb = RGBColor(*_rgb(kicker["fill_rgb"]))
    add_textbox(
        slide,
        int(kicker["text_left"]),
        int(kicker["text_top"]),
        int(kicker["text_width"]),
        int(kicker["text_height"]),
        str(payload.get("center_kicker", kicker["text"])),
        size_pt=int(kicker["font_pt"]),
        bold=True,
        color=_rgb(kicker["text_rgb"]),
        align=PP_ALIGN.CENTER,
    )

    center_text = spec["center_text"]
    add_textbox(
        slide,
        int(center_text["title_left"]),
        int(center_text["title_top"]),
        int(center_text["title_width"]),
        int(center_text["title_height"]),
        str(payload.get("center_title", page.title)),
        size_pt=int(center_text["title_font_pt"]),
        bold=True,
        color=_rgb(center_text["title_rgb"]),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        int(center_text["subtitle_left"]),
        int(center_text["subtitle_top"]),
        int(center_text["subtitle_width"]),
        int(center_text["subtitle_height"]),
        str(payload.get("center_subtitle", page.subtitle)),
        size_pt=int(center_text["subtitle_font_pt"]),
        color=_rgb(center_text["subtitle_rgb"]),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        int(center_text["footer_left"]),
        int(center_text["footer_top"]),
        int(center_text["footer_width"]),
        int(center_text["footer_height"]),
        str(payload.get("center_bottom_footer", center_text["footer_text"])),
        size_pt=int(center_text["footer_font_pt"]),
        bold=True,
        color=_rgb(center_text["footer_rgb"]),
        align=PP_ALIGN.CENTER,
    )


def _render_capability_ring(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "capability_ring")
    payload = page.payload
    items = list(payload.get("items", []))[:7]
    grid = spec["grid"]
    title_box = spec["title_box"]

    add_textbox(
        slide,
        int(title_box["left"]),
        int(title_box["top"]),
        int(title_box["width"]),
        int(title_box["height"]),
        str(payload.get("headline", page.subtitle or title_box["text"])),
        size_pt=int(title_box["font_pt"]),
        bold=True,
        color=_rgb(title_box["text_rgb"]),
        align=PP_ALIGN.CENTER,
    )

    for idx, item in enumerate(items):
        row, col = divmod(idx, int(grid["columns"]))
        left = int(grid["left"]) + col * (int(grid["card_width"]) + int(grid["gap_x"]))
        top = int(grid["top"]) + row * (int(grid["card_height"]) + int(grid["gap_y"]))
        card = slide.shapes.add_shape(1, left, top, int(grid["card_width"]), int(grid["card_height"]))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["card_fill_rgb"]))
        card.line.color.rgb = RGBColor(*_rgb(spec["card_line_rgb"]))
        add_textbox(
            slide,
            left + int(spec["card_title"]["left_offset"]),
            top + int(spec["card_title"]["top_offset"]),
            int(grid["card_width"]) - int(spec["card_title"]["width_padding"]),
            int(spec["card_title"]["height"]),
            str(item.get("title", "")),
            size_pt=int(spec["card_title"]["font_pt"]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + int(spec["card_detail"]["left_offset"]),
            top + int(spec["card_detail"]["top_offset"]),
            int(grid["card_width"]) - int(spec["card_detail"]["width_padding"]),
            int(spec["card_detail"]["height"]),
            normalize_text_for_body_box(str(item.get("detail", "")), int(spec["card_detail"]["wrap_chars"])),
            size_pt=int(spec["card_detail"]["font_pt"]),
            align=PP_ALIGN.CENTER,
        )


def _render_five_phase_path(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "five_phase_path")
    payload = page.payload
    intro_box = spec["intro_box"]
    add_textbox(
        slide,
        int(intro_box["left"]),
        int(intro_box["top"]),
        int(intro_box["width"]),
        int(intro_box["height"]),
        str(payload.get("intro", page.subtitle)),
        size_pt=int(intro_box["font_pt"]),
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    stages = list(payload.get("stages", []))[: int(spec["stage_count"])]
    while len(stages) < int(spec["stage_count"]):
        stages.append({"header": "", "tasks": []})

    for idx, stage in enumerate(stages):
        left = int(spec["origin_left"]) + idx * (int(spec["stage_width"]) + int(spec["gap_x"]))
        header = slide.shapes.add_shape(
            1,
            left,
            int(spec["origin_top"]),
            int(spec["stage_width"]),
            int(spec["header_height"]),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*_rgb(spec["header_fill_rgb"]))
        header.line.color.rgb = RGBColor(*_rgb(spec["header_fill_rgb"]))
        add_textbox(
            slide,
            left,
            int(spec["origin_top"]) + int(spec["header_text_top_offset"]),
            int(spec["stage_width"]),
            int(spec["header_height"]),
            str(stage.get("header", "")),
            size_pt=int(spec["header_font_pt"]),
            bold=True,
            color=_rgb(spec["header_text_rgb"]),
            align=PP_ALIGN.CENTER,
        )

        body = slide.shapes.add_shape(
            1,
            left,
            int(spec["origin_top"]) + int(spec["body_top_offset"]),
            int(spec["stage_width"]),
            int(spec["body_height"]),
        )
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(*_rgb(spec["body_fill_rgb"]))
        body.line.color.rgb = RGBColor(*_rgb(spec["body_line_rgb"]))

        tasks = list(stage.get("tasks", []))[: int(spec["max_tasks_per_stage"])]
        for task_idx, task in enumerate(tasks):
            top = (
                int(spec["origin_top"])
                + int(spec["body_top_offset"])
                + int(spec["task_top_offset"])
                + task_idx * int(spec["task_gap_y"])
            )
            add_textbox(
                slide,
                left + int(spec["task_left_offset"]),
                top,
                int(spec["stage_width"]) - int(spec["task_width_padding"]),
                int(spec["task_height"]),
                normalize_text_for_body_box(str(task), int(spec["task_wrap_chars"])),
                size_pt=int(spec["task_font_pt"]),
                align=PP_ALIGN.CENTER,
            )


def _render_pain_cards(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "pain_cards")
    payload = page.payload
    lead_box = spec["lead_box"]
    add_textbox(
        slide,
        int(lead_box["left"]),
        int(lead_box["top"]),
        int(lead_box["width"]),
        int(lead_box["height"]),
        str(payload.get("lead", page.subtitle)),
        size_pt=int(lead_box["font_pt"]),
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    cards = list(payload.get("cards", []))[:3]
    grid = spec["grid"]
    for idx, card_data in enumerate(cards):
        left = int(grid["left"]) + idx * (int(grid["card_width"]) + int(grid["gap_x"]))
        top = int(grid["top"])
        card = slide.shapes.add_shape(1, left, top, int(grid["card_width"]), int(grid["card_height"]))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["card_fill_rgb"]))
        card.line.color.rgb = RGBColor(*_rgb(spec["card_line_rgb"]))

        add_textbox(
            slide,
            left + int(spec["title_box"]["left_offset"]),
            top + int(spec["title_box"]["top_offset"]),
            int(grid["card_width"]) - int(spec["title_box"]["width_padding"]),
            int(spec["title_box"]["height"]),
            str(card_data.get("title", "")),
            size_pt=int(spec["title_box"]["font_pt"]),
            bold=True,
            color=_accent_rgb(manifest),
        )
        add_textbox(
            slide,
            left + int(spec["detail_box"]["left_offset"]),
            top + int(spec["detail_box"]["top_offset"]),
            int(grid["card_width"]) - int(spec["detail_box"]["width_padding"]),
            int(spec["detail_box"]["height"]),
            normalize_text_for_body_box(str(card_data.get("detail", "")), int(spec["detail_box"]["wrap_chars"])),
            size_pt=int(spec["detail_box"]["font_pt"]),
        )
        points = list(card_data.get("points", []))[:3]
        for point_idx, point in enumerate(points):
            point_top = top + int(spec["bullet_box"]["top_offset"]) + point_idx * int(spec["bullet_box"]["gap_y"])
            add_textbox(
                slide,
                left + int(spec["bullet_box"]["left_offset"]),
                point_top,
                int(grid["card_width"]) - int(spec["bullet_box"]["width_padding"]),
                int(spec["bullet_box"]["height"]),
                str(point),
                size_pt=int(spec["bullet_box"]["font_pt"]),
            )

    banner = spec["banner"]
    banner_shape = slide.shapes.add_shape(
        1,
        int(banner["left"]),
        int(banner["top"]),
        int(banner["width"]),
        int(banner["height"]),
    )
    banner_shape.fill.solid()
    banner_shape.fill.fore_color.rgb = RGBColor(*_rgb(banner["fill_rgb"]))
    banner_shape.line.color.rgb = RGBColor(*_rgb(banner["fill_rgb"]))
    add_textbox(
        slide,
        int(banner["text_left"]),
        int(banner["text_top"]),
        int(banner["text_width"]),
        int(banner["text_height"]),
        str(payload.get("bottom_banner", banner["text"])),
        size_pt=int(banner["font_pt"]),
        bold=True,
        color=_rgb(banner["text_rgb"]),
        align=PP_ALIGN.CENTER,
    )


def _render_roadmap_timeline(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "roadmap_timeline")
    payload = page.payload
    add_textbox(
        slide,
        int(spec["headline_box"]["left"]),
        int(spec["headline_box"]["top"]),
        int(spec["headline_box"]["width"]),
        int(spec["headline_box"]["height"]),
        str(payload.get("headline", page.subtitle)),
        size_pt=int(spec["headline_box"]["font_pt"]),
        bold=True,
        color=_rgb(spec["period_box"]["text_rgb"]),
        align=PP_ALIGN.CENTER,
    )

    timeline = slide.shapes.add_shape(
        1,
        int(spec["timeline"]["left"]),
        int(spec["timeline"]["top"]),
        int(spec["timeline"]["width"]),
        int(spec["timeline"]["height"]),
    )
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = RGBColor(*_rgb(spec["timeline"]["fill_rgb"]))
    timeline.line.color.rgb = RGBColor(*_rgb(spec["timeline"]["fill_rgb"]))

    stages = list(payload.get("stages", []))[: int(spec["grid"]["columns"])]
    for idx, stage in enumerate(stages):
        left = int(spec["grid"]["left"]) + idx * (int(spec["grid"]["card_width"]) + int(spec["grid"]["gap_x"]))
        card_top = int(spec["grid"]["top"]) + int(spec["card_style"]["top_offset"])
        add_textbox(
            slide,
            left + int(spec["period_box"]["left_offset"]),
            int(spec["grid"]["top"]) + int(spec["period_box"]["top_offset"]),
            int(spec["period_box"]["width"]),
            int(spec["period_box"]["height"]),
            str(stage.get("period", f"阶段{idx + 1}")),
            size_pt=int(spec["period_box"]["font_pt"]),
            bold=True,
            color=_rgb(spec["period_box"]["text_rgb"]),
            align=PP_ALIGN.CENTER,
        )

        connector = slide.shapes.add_shape(
            1,
            left + int(spec["milestone"]["center_x_offset"]) - int(spec["timeline"]["height"]) // 2,
            int(spec["grid"]["top"]) + int(spec["period_box"]["height"]),
            int(spec["timeline"]["height"]),
            int(spec["milestone"]["line_height"]),
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = RGBColor(*_rgb(spec["milestone"]["line_rgb"]))
        connector.line.color.rgb = RGBColor(*_rgb(spec["milestone"]["line_rgb"]))

        milestone = slide.shapes.add_shape(
            1,
            left + int(spec["milestone"]["center_x_offset"]) - int(spec["milestone"]["size"]) // 2,
            int(spec["milestone"]["center_y"]) - int(spec["milestone"]["size"]) // 2,
            int(spec["milestone"]["size"]),
            int(spec["milestone"]["size"]),
        )
        milestone.fill.solid()
        milestone.fill.fore_color.rgb = RGBColor(*_rgb(spec["milestone"]["fill_rgb"]))
        milestone.line.color.rgb = RGBColor(*_rgb(spec["milestone"]["fill_rgb"]))

        card = slide.shapes.add_shape(
            1,
            left,
            card_top,
            int(spec["grid"]["card_width"]),
            int(spec["grid"]["card_height"]),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["card_style"]["fill_rgb"]))
        card.line.color.rgb = RGBColor(*_rgb(spec["card_style"]["line_rgb"]))

        add_textbox(
            slide,
            left + int(spec["title_box"]["left_offset"]),
            card_top + int(spec["title_box"]["top_offset"]),
            int(spec["grid"]["card_width"]) - int(spec["title_box"]["width_padding"]),
            int(spec["title_box"]["height"]),
            str(stage.get("title", "")),
            size_pt=int(spec["title_box"]["font_pt"]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + int(spec["detail_box"]["left_offset"]),
            card_top + int(spec["detail_box"]["top_offset"]),
            int(spec["grid"]["card_width"]) - int(spec["detail_box"]["width_padding"]),
            int(spec["detail_box"]["height"]),
            normalize_text_for_body_box(str(stage.get("detail", "")), int(spec["detail_box"]["wrap_chars"])),
            size_pt=int(spec["detail_box"]["font_pt"]),
            align=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        int(spec["footer_box"]["left"]),
        int(spec["footer_box"]["top"]),
        int(spec["footer_box"]["width"]),
        int(spec["footer_box"]["height"]),
        str(payload.get("footer", page.subtitle)),
        size_pt=int(spec["footer_box"]["font_pt"]),
        color=_rgb(spec["period_box"]["text_rgb"]),
        align=PP_ALIGN.CENTER,
    )


def _render_kpi_dashboard(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "kpi_dashboard")
    payload = page.payload
    accent_rgb = _accent_rgb(manifest)
    sub_rgb = _inactive_rgb(manifest)
    add_textbox(
        slide,
        int(spec["headline_box"]["left"]),
        int(spec["headline_box"]["top"]),
        int(spec["headline_box"]["width"]),
        int(spec["headline_box"]["height"]),
        str(payload.get("headline", page.subtitle)),
        size_pt=int(spec["headline_box"]["font_pt"]),
        bold=True,
        color=sub_rgb,
        align=PP_ALIGN.CENTER,
    )

    metrics = list(payload.get("metrics", []))[: int(spec["metrics_grid"]["columns"])]
    for idx, metric in enumerate(metrics):
        left = int(spec["metrics_grid"]["left"]) + idx * (
            int(spec["metrics_grid"]["card_width"]) + int(spec["metrics_grid"]["gap_x"])
        )
        top = int(spec["metrics_grid"]["top"])
        card = slide.shapes.add_shape(
            1,
            left,
            top,
            int(spec["metrics_grid"]["card_width"]),
            int(spec["metrics_grid"]["card_height"]),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["metric_card"]["fill_rgb"]))
        card.line.color.rgb = RGBColor(*_rgb(spec["metric_card"]["line_rgb"]))

        add_textbox(
            slide,
            left + int(spec["metric_label_box"]["left_offset"]),
            top + int(spec["metric_label_box"]["top_offset"]),
            int(spec["metrics_grid"]["card_width"]) - int(spec["metric_label_box"]["width_padding"]),
            int(spec["metric_label_box"]["height"]),
            str(metric.get("label", f"KPI {idx + 1}")),
            size_pt=int(spec["metric_label_box"]["font_pt"]),
            bold=True,
            color=sub_rgb,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + int(spec["metric_value_box"]["left_offset"]),
            top + int(spec["metric_value_box"]["top_offset"]),
            int(spec["metrics_grid"]["card_width"]) - int(spec["metric_value_box"]["width_padding"]),
            int(spec["metric_value_box"]["height"]),
            str(metric.get("value", "--")),
            size_pt=int(spec["metric_value_box"]["font_pt"]),
            bold=True,
            color=_rgb(spec["metric_value_box"]["text_rgb"]) or accent_rgb,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + int(spec["metric_detail_box"]["left_offset"]),
            top + int(spec["metric_detail_box"]["top_offset"]),
            int(spec["metrics_grid"]["card_width"]) - int(spec["metric_detail_box"]["width_padding"]),
            int(spec["metric_detail_box"]["height"]),
            normalize_text_for_body_box(str(metric.get("detail", "")), int(spec["metric_detail_box"]["wrap_chars"])),
            size_pt=int(spec["metric_detail_box"]["font_pt"]),
            color=sub_rgb,
            align=PP_ALIGN.CENTER,
        )

    insight_band = slide.shapes.add_shape(
        1,
        int(spec["insight_band"]["left"]),
        int(spec["insight_band"]["top"]),
        int(spec["insight_band"]["width"]),
        int(spec["insight_band"]["height"]),
    )
    insight_band.fill.solid()
    insight_band.fill.fore_color.rgb = RGBColor(*_rgb(spec["insight_band"]["fill_rgb"]))
    insight_band.line.color.rgb = RGBColor(*_rgb(spec["insight_band"]["line_rgb"]))
    add_textbox(
        slide,
        int(spec["insight_band"]["text_left"]),
        int(spec["insight_band"]["text_top"]),
        int(spec["insight_band"]["text_width"]),
        int(spec["insight_band"]["text_height"]),
        str(payload.get("band_text", spec["insight_band"]["text"])),
        size_pt=int(spec["insight_band"]["font_pt"]),
        bold=True,
        color=_rgb(spec["insight_band"]["text_rgb"]),
        align=PP_ALIGN.CENTER,
    )

    insights = list(payload.get("insights", []))[: int(spec["insights_grid"]["columns"])]
    for idx, insight in enumerate(insights):
        left = int(spec["insights_grid"]["left"]) + idx * (
            int(spec["insights_grid"]["card_width"]) + int(spec["insights_grid"]["gap_x"])
        )
        top = int(spec["insights_grid"]["top"])
        card = slide.shapes.add_shape(
            1,
            left,
            top,
            int(spec["insights_grid"]["card_width"]),
            int(spec["insights_grid"]["card_height"]),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["insight_card"]["fill_rgb"]))
        card.line.color.rgb = RGBColor(*_rgb(spec["insight_card"]["line_rgb"]))
        add_textbox(
            slide,
            left + int(spec["insight_text_box"]["left_offset"]),
            top + int(spec["insight_text_box"]["top_offset"]),
            int(spec["insights_grid"]["card_width"]) - int(spec["insight_text_box"]["width_padding"]),
            int(spec["insight_text_box"]["height"]),
            normalize_text_for_body_box(str(insight), int(spec["insight_text_box"]["wrap_chars"])),
            size_pt=int(spec["insight_text_box"]["font_pt"]),
            align=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        int(spec["footer_box"]["left"]),
        int(spec["footer_box"]["top"]),
        int(spec["footer_box"]["width"]),
        int(spec["footer_box"]["height"]),
        str(payload.get("footer", page.subtitle)),
        size_pt=int(spec["footer_box"]["font_pt"]),
        color=sub_rgb,
        align=PP_ALIGN.CENTER,
    )


def _render_risk_matrix(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "risk_matrix")
    payload = page.payload
    sub_rgb = _inactive_rgb(manifest)
    add_textbox(
        slide,
        int(spec["headline_box"]["left"]),
        int(spec["headline_box"]["top"]),
        int(spec["headline_box"]["width"]),
        int(spec["headline_box"]["height"]),
        str(payload.get("headline", page.subtitle)),
        size_pt=int(spec["headline_box"]["font_pt"]),
        bold=True,
        color=sub_rgb,
        align=PP_ALIGN.CENTER,
    )

    matrix = slide.shapes.add_shape(
        1,
        int(spec["matrix"]["left"]),
        int(spec["matrix"]["top"]),
        int(spec["matrix"]["width"]),
        int(spec["matrix"]["height"]),
    )
    matrix.fill.background()
    matrix.line.color.rgb = RGBColor(*_rgb(spec["matrix"]["line_rgb"]))

    vertical = slide.shapes.add_shape(
        1,
        int(spec["matrix"]["left"]) + int(spec["matrix"]["width"]) // 2,
        int(spec["matrix"]["top"]),
        40000,
        int(spec["matrix"]["height"]),
    )
    vertical.fill.solid()
    vertical.fill.fore_color.rgb = RGBColor(*_rgb(spec["matrix"]["line_rgb"]))
    vertical.line.color.rgb = RGBColor(*_rgb(spec["matrix"]["line_rgb"]))

    horizontal = slide.shapes.add_shape(
        1,
        int(spec["matrix"]["left"]),
        int(spec["matrix"]["top"]) + int(spec["matrix"]["height"]) // 2,
        int(spec["matrix"]["width"]),
        40000,
    )
    horizontal.fill.solid()
    horizontal.fill.fore_color.rgb = RGBColor(*_rgb(spec["matrix"]["line_rgb"]))
    horizontal.line.color.rgb = RGBColor(*_rgb(spec["matrix"]["line_rgb"]))

    quadrant_positions = {
        "high_high": (1500000, 2260000),
        "low_high": (6150000, 2260000),
        "high_low": (1500000, 3310000),
        "low_low": (6150000, 3310000),
    }
    labels = list(spec["quadrants"]["labels"])
    for label in labels:
        add_textbox(
            slide,
            int(label["left"]),
            int(label["top"]),
            int(spec["label_box"]["width"]),
            int(spec["label_box"]["height"]),
            str(label["name"]),
            size_pt=int(spec["label_box"]["font_pt"]),
            bold=True,
            color=_rgb(spec["label_box"]["text_rgb"]),
        )

    items = list(payload.get("items", []))[:4]
    for item in items:
        quadrant = str(item.get("quadrant", "high_high"))
        left, top = quadrant_positions.get(quadrant, quadrant_positions["high_high"])
        is_high = quadrant == "high_high"
        card = slide.shapes.add_shape(
            1,
            left,
            top,
            int(spec["risk_card"]["width"]),
            int(spec["risk_card"]["height"]),
        )
        fill_key = "high_fill_rgb" if is_high else "card_fill_rgb"
        line_key = "high_line_rgb" if is_high else "card_line_rgb"
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["quadrants"][fill_key]))
        card.line.color.rgb = RGBColor(*_rgb(spec["quadrants"][line_key]))
        add_textbox(
            slide,
            left + int(spec["risk_title_box"]["left_offset"]),
            top + int(spec["risk_title_box"]["top_offset"]),
            int(spec["risk_card"]["width"]) - int(spec["risk_title_box"]["width_padding"]),
            int(spec["risk_title_box"]["height"]),
            str(item.get("title", "")),
            size_pt=int(spec["risk_title_box"]["font_pt"]),
            bold=True,
        )
        add_textbox(
            slide,
            left + int(spec["risk_detail_box"]["left_offset"]),
            top + int(spec["risk_detail_box"]["top_offset"]),
            int(spec["risk_card"]["width"]) - int(spec["risk_detail_box"]["width_padding"]),
            int(spec["risk_detail_box"]["height"]),
            normalize_text_for_body_box(str(item.get("detail", "")), int(spec["risk_detail_box"]["wrap_chars"])),
            size_pt=int(spec["risk_detail_box"]["font_pt"]),
        )

    add_textbox(
        slide,
        int(spec["footer_box"]["left"]),
        int(spec["footer_box"]["top"]),
        int(spec["footer_box"]["width"]),
        int(spec["footer_box"]["height"]),
        str(payload.get("footer", page.subtitle)),
        size_pt=int(spec["footer_box"]["font_pt"]),
        color=sub_rgb,
        align=PP_ALIGN.CENTER,
    )


def _render_claim_breakdown(slide, page: BodyPageSpec, manifest: TemplateManifest):
    spec = _layout(manifest, "claim_breakdown")
    payload = page.payload
    sub_rgb = _inactive_rgb(manifest)
    add_textbox(
        slide,
        int(spec["headline_box"]["left"]),
        int(spec["headline_box"]["top"]),
        int(spec["headline_box"]["width"]),
        int(spec["headline_box"]["height"]),
        str(payload.get("headline", page.subtitle)),
        size_pt=int(spec["headline_box"]["font_pt"]),
        bold=True,
        color=sub_rgb,
        align=PP_ALIGN.CENTER,
    )

    claims = list(payload.get("claims", []))[: int(spec["claims_grid"]["columns"])]
    for idx, claim in enumerate(claims):
        left = int(spec["claims_grid"]["left"]) + idx * (
            int(spec["claims_grid"]["card_width"]) + int(spec["claims_grid"]["gap_x"])
        )
        top = int(spec["claims_grid"]["top"])
        card = slide.shapes.add_shape(
            1,
            left,
            top,
            int(spec["claims_grid"]["card_width"]),
            int(spec["claims_grid"]["card_height"]),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(*_rgb(spec["claim_card"]["fill_rgb"]))
        card.line.color.rgb = RGBColor(*_rgb(spec["claim_card"]["line_rgb"]))

        add_textbox(
            slide,
            left + int(spec["claim_label_box"]["left_offset"]),
            top + int(spec["claim_label_box"]["top_offset"]),
            int(spec["claims_grid"]["card_width"]) - int(spec["claim_label_box"]["width_padding"]),
            int(spec["claim_label_box"]["height"]),
            str(claim.get("label", f"项目{idx + 1}")),
            size_pt=int(spec["claim_label_box"]["font_pt"]),
            bold=True,
            color=sub_rgb,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + int(spec["claim_value_box"]["left_offset"]),
            top + int(spec["claim_value_box"]["top_offset"]),
            int(spec["claims_grid"]["card_width"]) - int(spec["claim_value_box"]["width_padding"]),
            int(spec["claim_value_box"]["height"]),
            str(claim.get("value", "--")),
            size_pt=int(spec["claim_value_box"]["font_pt"]),
            bold=True,
            color=_rgb(spec["claim_value_box"]["text_rgb"]),
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + int(spec["claim_detail_box"]["left_offset"]),
            top + int(spec["claim_detail_box"]["top_offset"]),
            int(spec["claims_grid"]["card_width"]) - int(spec["claim_detail_box"]["width_padding"]),
            int(spec["claim_detail_box"]["height"]),
            normalize_text_for_body_box(str(claim.get("detail", "")), int(spec["claim_detail_box"]["wrap_chars"])),
            size_pt=int(spec["claim_detail_box"]["font_pt"]),
            color=sub_rgb,
            align=PP_ALIGN.CENTER,
        )

    summary_band = slide.shapes.add_shape(
        1,
        int(spec["summary_band"]["left"]),
        int(spec["summary_band"]["top"]),
        int(spec["summary_band"]["width"]),
        int(spec["summary_band"]["height"]),
    )
    summary_band.fill.solid()
    summary_band.fill.fore_color.rgb = RGBColor(*_rgb(spec["summary_band"]["fill_rgb"]))
    summary_band.line.color.rgb = RGBColor(*_rgb(spec["summary_band"]["line_rgb"]))
    add_textbox(
        slide,
        int(spec["summary_band"]["text_left"]),
        int(spec["summary_band"]["text_top"]),
        int(spec["summary_band"]["text_width"]),
        int(spec["summary_band"]["text_height"]),
        str(payload.get("band_text", spec["summary_band"]["text"])),
        size_pt=int(spec["summary_band"]["font_pt"]),
        bold=True,
        color=_rgb(spec["summary_band"]["text_rgb"]),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        int(spec["summary_box"]["left"]),
        int(spec["summary_box"]["top"]),
        int(spec["summary_box"]["width"]),
        int(spec["summary_box"]["height"]),
        normalize_text_for_body_box(str(payload.get("summary", page.subtitle)), int(spec["summary_box"]["wrap_chars"])),
        size_pt=int(spec["summary_box"]["font_pt"]),
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        int(spec["footer_box"]["left"]),
        int(spec["footer_box"]["top"]),
        int(spec["footer_box"]["width"]),
        int(spec["footer_box"]["height"]),
        str(payload.get("footer", page.subtitle)),
        size_pt=int(spec["footer_box"]["font_pt"]),
        color=sub_rgb,
        align=PP_ALIGN.CENTER,
    )


PATTERN_RENDERERS = {
    "general_business": _render_cards_2x2,
    "process_flow": _render_process_flow,
    "solution_architecture": _render_architecture_layers,
    "org_governance": _render_governance_grid,
    "comparison_upgrade": _render_comparison_upgrade,
    "capability_ring": _render_capability_ring,
    "five_phase_path": _render_five_phase_path,
    "pain_cards": _render_pain_cards,
    "roadmap_timeline": _render_roadmap_timeline,
    "kpi_dashboard": _render_kpi_dashboard,
    "risk_matrix": _render_risk_matrix,
    "claim_breakdown": _render_claim_breakdown,
}


def resolve_render_pattern(pattern_id: str) -> str:
    return pattern_id if pattern_id in PATTERN_RENDERERS else "general_business"

