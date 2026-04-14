from copy import deepcopy
import logging

from pptx import Presentation

LOGGER = logging.getLogger(__name__)


def clone_slide_after(prs: Presentation, source_idx: int, insert_after_idx: int, keep_rel_ids: bool = True):
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in list(new_slide.shapes):
        element = shape._element
        element.getparent().remove(element)
    for shape in source.shapes:
        new_element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            if keep_rel_ids:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
            else:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target)
        except Exception as exc:
            LOGGER.warning("slide rel copy failed: %s", exc)
    slide_id_list = prs.slides._sldIdLst
    new_id = slide_id_list[-1]
    del slide_id_list[-1]
    slide_id_list.insert(insert_after_idx + 1, new_id)
    return prs.slides[insert_after_idx + 1]


def remove_slide(prs: Presentation, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    prs.part.drop_rel(slide_id.rId)
    del slide_id_list[slide_index]


def ensure_last_slide(prs: Presentation, slide_id: int):
    slide_id_list = prs.slides._sldIdLst
    target = None
    for item in slide_id_list:
        if int(item.id) == int(slide_id):
            target = item
            break
    if target is None:
        return
    slide_id_list.remove(target)
    slide_id_list.append(target)


__all__ = ["clone_slide_after", "remove_slide", "ensure_last_slide"]
