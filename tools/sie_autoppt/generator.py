import datetime
import gc
import re
import shutil
import warnings
from pathlib import Path

from pptx import Presentation

from .body_renderers import apply_theme_title, fill_body_slide, fill_directory_slide, resolve_render_pattern
from .config import DEFAULT_MIN_TEMPLATE_SLIDES, DEFAULT_OUTPUT_DIR
from .models import BodyPageSpec, DeckPlan, DeckRenderTrace, GenerationArtifacts, PageRenderTrace
from .pipeline import plan_deck_from_html, plan_deck_from_json
from .planning.deck_planner import build_directory_window
from .reference_styles import build_reference_import_plan, populate_reference_body_pages
from .slide_ops import (
    clone_slide_after,
    copy_slide_xml_assets,
    ensure_last_slide,
    import_slides_from_presentation,
    remove_slide,
    slide_assets_preserved,
)
from .template_manifest import TemplateManifest, load_template_manifest


LEGACY_CLONE_DEPRECATION_MESSAGE = (
    "Legacy runtime slide cloning is deprecated. Please migrate this template to a preallocated slide pool."
)


def build_output_path(output_dir: Path, output_prefix: str) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    safe_prefix = re.sub(r'[<>:"/\\\\|?*]+', "_", output_prefix).strip(" ._") or "SIE_AutoPPT"
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")[:-3]
    return output_dir / f"{safe_prefix}_{timestamp}.pptx"


def _directory_lines_for_page(body_pages: list[BodyPageSpec], active_index: int) -> tuple[list[str], int]:
    return build_directory_window(body_pages, active_index)


def validate_slide_pool_configuration(manifest: TemplateManifest, body_page_count: int, slide_count: int):
    if not manifest.slide_pools:
        return

    directory_pool = list(manifest.slide_pools.directory)
    body_pool = list(manifest.slide_pools.body)
    if len(directory_pool) != len(body_pool):
        raise ValueError("Template manifest slide pool is invalid: directory/body pool lengths do not match.")
    if manifest.slide_pools.ending is None:
        raise ValueError("Template manifest slide pool is invalid: ending slide is not configured.")
    if len(directory_pool) < body_page_count or len(body_pool) < body_page_count:
        raise ValueError(
            "Template manifest slide pool is too small for this deck: "
            f"requested {body_page_count} body pages, "
            f"but only {min(len(directory_pool), len(body_pool))} preallocated pairs are configured in "
            f"{manifest.manifest_path}."
        )

    used_indices = directory_pool[:body_page_count] + body_pool[:body_page_count] + [manifest.slide_pools.ending]
    if any(index < 0 for index in used_indices):
        raise ValueError("Template manifest slide pool is invalid: negative slide index detected.")
    if max(used_indices, default=-1) >= slide_count:
        raise ValueError("Template manifest slide pool is invalid: slide index exceeds the template slide count.")
    if len(set(used_indices)) != len(used_indices):
        raise ValueError("Template manifest slide pool is invalid: duplicate slide index detected.")


def _reference_import_unavailable_reason(body_pages: list[BodyPageSpec], reference_body_path: Path | None) -> str:
    if not any(page.reference_style_id for page in body_pages):
        return ""
    if reference_body_path is None or not reference_body_path.exists():
        return "reference body slide library is unavailable"
    return ""


def _build_page_render_traces(
    body_pages: list[BodyPageSpec],
    reference_import_applied: bool,
    reference_import_reason: str,
) -> list[PageRenderTrace]:
    traces = []
    for page in body_pages:
        actual_pattern_id = resolve_render_pattern(page.pattern_id)
        if page.reference_style_id:
            if reference_import_applied:
                render_route = f"reference_import:{page.reference_style_id}"
                fallback_reason = ""
            else:
                render_route = f"native_fallback:{actual_pattern_id}"
                fallback_reason = reference_import_reason or "reference style import was not applied"
        else:
            render_route = f"native_renderer:{actual_pattern_id}"
            fallback_reason = ""
        traces.append(
            PageRenderTrace(
                page_key=page.page_key,
                title=page.title,
                requested_pattern_id=page.pattern_id,
                actual_pattern_id=actual_pattern_id,
                reference_style_id=page.reference_style_id,
                render_route=render_route,
                fallback_reason=fallback_reason,
            )
        )
    return traces


def _apply_reference_body_slides(
    pptx_path: Path,
    reference_body_path: Path | None,
    body_pages: list[BodyPageSpec],
    manifest: TemplateManifest,
) -> bool:
    import_plan = build_reference_import_plan(body_pages, reference_body_path=reference_body_path, manifest=manifest)
    if not import_plan:
        return False
    if reference_body_path is None or not reference_body_path.exists():
        return False
    try:
        if import_slides_from_presentation(pptx_path, reference_body_path, import_plan):
            return True
    except Exception as exc:
        warnings.warn(
            f"Reference body slide import failed during native package merge: {exc}",
            stacklevel=2,
        )
    return False


def _render_with_preallocated_pool(
    prs: Presentation,
    body_pages: list[BodyPageSpec],
    chapter_lines: list[str],
    active_start: int,
    manifest: TemplateManifest,
):
    if not manifest.slide_pools:
        raise ValueError("Template manifest does not define slide pools.")

    used_directory_indices = list(manifest.slide_pools.directory[: len(body_pages)])
    used_body_indices = list(manifest.slide_pools.body[: len(body_pages)])
    unused_pairs = list(
        zip(
            manifest.slide_pools.directory[len(body_pages):],
            manifest.slide_pools.body[len(body_pages):],
            strict=False,
        )
    )

    for directory_idx, body_idx in sorted(unused_pairs, reverse=True):
        for slide_index in sorted((directory_idx, body_idx), reverse=True):
            if slide_index < len(prs.slides):
                remove_slide(prs, slide_index)

    directory_slides = [prs.slides[index] for index in used_directory_indices]
    body_slides = [prs.slides[index] for index in used_body_indices]

    for chapter_idx, directory_slide in enumerate(directory_slides):
        window_lines, window_active_index = _directory_lines_for_page(body_pages, active_start + chapter_idx)
        fill_directory_slide(directory_slide, window_lines, window_active_index, manifest)
    for page, body_slide in zip(body_pages, body_slides):
        fill_body_slide(body_slide, page, manifest)


def _render_with_legacy_clone(
    prs: Presentation,
    body_pages: list[BodyPageSpec],
    chapter_lines: list[str],
    active_start: int,
    manifest: TemplateManifest,
):
    warnings.warn(LEGACY_CLONE_DEPRECATION_MESSAGE, DeprecationWarning, stacklevel=2)
    directory_idx = manifest.slide_roles.directory
    body_template_idx = manifest.slide_roles.body_template
    directory_slides = [prs.slides[directory_idx]]
    body_slides = [prs.slides[body_template_idx]]
    insert_after = body_template_idx
    for _ in body_pages[1:]:
        new_directory = clone_slide_after(prs, directory_idx, insert_after, keep_rel_ids=True)
        directory_slides.append(new_directory)
        insert_after += 1

        new_body = clone_slide_after(prs, body_template_idx, insert_after, keep_rel_ids=False)
        body_slides.append(new_body)
        insert_after += 1

    for chapter_idx, directory_slide in enumerate(directory_slides):
        window_lines, window_active_index = _directory_lines_for_page(body_pages, active_start + chapter_idx)
        fill_directory_slide(directory_slide, window_lines, window_active_index, manifest)
    for page, body_slide in zip(body_pages, body_slides):
        fill_body_slide(body_slide, page, manifest)


def _refresh_legacy_directory_clones(
    pptx_path: Path,
    body_pages: list[BodyPageSpec],
    active_start: int,
    body_page_count: int,
    manifest: TemplateManifest,
):
    warnings.warn(LEGACY_CLONE_DEPRECATION_MESSAGE, DeprecationWarning, stacklevel=2)
    directory_idx = manifest.slide_roles.directory
    targets = [directory_idx + 1 + i * 2 for i in range(1, body_page_count)]
    if not targets:
        return True

    source_idx = directory_idx + 1
    for _ in range(3):
        if not copy_slide_xml_assets(pptx_path, source_idx=source_idx, target_indices=targets):
            continue

        prs_reloaded = Presentation(str(pptx_path))
        window_lines, window_active_index = _directory_lines_for_page(body_pages, active_start)
        fill_directory_slide(prs_reloaded.slides[directory_idx], window_lines, window_active_index, manifest)
        for offset, directory_slide_no in enumerate(targets, start=1):
            slide_index = directory_slide_no - 1
            if slide_index < len(prs_reloaded.slides):
                window_lines, window_active_index = _directory_lines_for_page(body_pages, active_start + offset)
                fill_directory_slide(prs_reloaded.slides[slide_index], window_lines, window_active_index, manifest)
        prs_reloaded.save(str(pptx_path))
        prs_reloaded = None
        gc.collect()

        if slide_assets_preserved(pptx_path, source_idx=source_idx, target_indices=targets):
            return True

    raise RuntimeError(
        "Legacy clone 路径图片资源复制失败（已重试 3 次），"
        f"源目录页索引 {source_idx}，目标页索引 {targets}。请迁移到 preallocated pool 模板路径。"
    )


def _warn_if_reference_import_disabled(body_pages: list[BodyPageSpec], reference_body_path: Path | None):
    if not any(page.reference_style_id for page in body_pages):
        return
    if reference_body_path is None or not reference_body_path.exists():
        warnings.warn(
            "Reference style library is unavailable; using native fallback renderers for reference-style pages.",
            stacklevel=2,
        )


def _refresh_preallocated_directory_assets(pptx_path: Path, body_page_count: int, manifest: TemplateManifest) -> bool:
    if not manifest.slide_pools:
        return True
    target_indices = [index + 1 for index in manifest.slide_pools.directory[1:body_page_count]]
    if not target_indices:
        return True
    source_idx = manifest.slide_pools.directory[0] + 1
    if not copy_slide_xml_assets(pptx_path, source_idx=source_idx, target_indices=target_indices):
        return False
    return slide_assets_preserved(pptx_path, source_idx=source_idx, target_indices=target_indices)


def _generate_ppt_from_plan(
    deck_plan: DeckPlan,
    input_kind: str,
    template_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    manifest = load_template_manifest(template_path=template_path)
    deck = deck_plan.deck
    body_pages = deck.body_pages
    chapter_lines = deck_plan.chapter_lines

    final_output_dir = output_dir or DEFAULT_OUTPUT_DIR
    out = build_output_path(final_output_dir, output_prefix)
    shutil.copy2(template_path, out)

    prs = Presentation(str(out))
    if len(prs.slides) < DEFAULT_MIN_TEMPLATE_SLIDES:
        raise ValueError(
            f"模板页数不足，至少需要 {DEFAULT_MIN_TEMPLATE_SLIDES} 页，实际为 {len(prs.slides)} 页。"
        )

    validate_slide_pool_configuration(manifest, len(body_pages), len(prs.slides))
    apply_theme_title(prs, deck.cover_title, manifest)
    if manifest.slide_pools and manifest.slide_pools.ending is not None and manifest.slide_pools.ending < len(prs.slides):
        thanks_slide_id = int(prs.slides._sldIdLst[manifest.slide_pools.ending].id)
    else:
        thanks_slide_id = int(prs.slides._sldIdLst[len(prs.slides) - 1].id)

    used_preallocated_pool = bool(manifest.slide_pools)
    if used_preallocated_pool:
        _render_with_preallocated_pool(prs, body_pages, chapter_lines, active_start, manifest)
    else:
        warnings.warn(
            "Template does not provide a preallocated slide pool; using deprecated legacy runtime cloning for this deck.",
            stacklevel=2,
        )
        _render_with_legacy_clone(prs, body_pages, chapter_lines, active_start, manifest)

    ensure_last_slide(prs, thanks_slide_id)
    prs.save(str(out))
    prs = None
    gc.collect()

    reference_import_reason = _reference_import_unavailable_reason(body_pages, reference_body_path)
    _warn_if_reference_import_disabled(body_pages, reference_body_path)
    reference_import_applied = _apply_reference_body_slides(out, reference_body_path, body_pages, manifest)
    if reference_import_applied:
        populate_reference_body_pages(out, body_pages, manifest=manifest)
        reference_import_reason = ""
    elif any(page.reference_style_id for page in body_pages) and not reference_import_reason:
        reference_import_reason = "reference style import failed during native package merge"

    if used_preallocated_pool and not _refresh_preallocated_directory_assets(out, len(body_pages), manifest):
        warnings.warn(
            "Preallocated directory slide assets could not be fully refreshed after save; manual review is recommended.",
            stacklevel=2,
        )

    if not used_preallocated_pool:
        _refresh_legacy_directory_clones(out, body_pages, active_start, len(body_pages), manifest)

    render_trace = DeckRenderTrace(
        input_kind=input_kind,
        body_render_mode="preallocated_pool" if used_preallocated_pool else "legacy_clone",
        reference_import_applied=reference_import_applied,
        reference_import_reason=reference_import_reason,
        page_traces=_build_page_render_traces(body_pages, reference_import_applied, reference_import_reason),
    )
    return GenerationArtifacts(
        output_path=out,
        deck_plan=deck_plan,
        render_trace=render_trace,
    )


def generate_ppt_artifacts_from_html(
    template_path: Path,
    html_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    chapters: int | None,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    if not html_path.exists():
        raise FileNotFoundError(f"HTML not found: {html_path}")
    deck_plan = plan_deck_from_html(html_path, chapters)
    return _generate_ppt_from_plan(
        deck_plan=deck_plan,
        input_kind="html",
        template_path=template_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        active_start=active_start,
        output_dir=output_dir,
    )


def generate_ppt_artifacts_from_deck_spec(
    template_path: Path,
    deck_spec_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    if not deck_spec_path.exists():
        raise FileNotFoundError(f"Deck JSON not found: {deck_spec_path}")
    deck_plan = plan_deck_from_json(deck_spec_path)
    return _generate_ppt_from_plan(
        deck_plan=deck_plan,
        input_kind="deck_spec_json",
        template_path=template_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        active_start=active_start,
        output_dir=output_dir,
    )


def generate_ppt_artifacts_from_deck_plan(
    deck_plan: DeckPlan,
    input_kind: str,
    template_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    active_start: int,
    output_dir: Path | None = None,
) -> GenerationArtifacts:
    return _generate_ppt_from_plan(
        deck_plan=deck_plan,
        input_kind=input_kind,
        template_path=template_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        active_start=active_start,
        output_dir=output_dir,
    )


def generate_ppt(
    template_path: Path,
    html_path: Path,
    reference_body_path: Path | None,
    output_prefix: str,
    chapters: int | None,
    active_start: int,
    output_dir: Path | None = None,
):
    artifacts = generate_ppt_artifacts_from_html(
        template_path=template_path,
        html_path=html_path,
        reference_body_path=reference_body_path,
        output_prefix=output_prefix,
        chapters=chapters,
        active_start=active_start,
        output_dir=output_dir,
    )
    return artifacts.output_path, artifacts.deck_plan.pattern_ids, artifacts.deck_plan.chapter_lines
