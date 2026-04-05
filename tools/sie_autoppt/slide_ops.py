import time
import zipfile
from copy import deepcopy
from pathlib import Path
from posixpath import basename, dirname, join, normpath, relpath, splitext
from xml.etree import ElementTree

from pptx import Presentation

PACKAGE_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CONTENT_TYPES_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
PRESENTATION_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
SKIPPED_RELATIONSHIP_SUFFIXES = (
    "/notesSlide",
    "/comments",
    "/commentAuthors",
)

ElementTree.register_namespace("", PACKAGE_REL_NS)
ElementTree.register_namespace("", CONTENT_TYPES_NS)


def clone_slide_after(prs: Presentation, source_idx: int, insert_after_idx: int, keep_rel_ids: bool = True):
    source = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for shape in list(new_slide.shapes):
        element = shape._element
        element.getparent().remove(element)
    for shape in source.shapes:
        new_element = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            if keep_rel_ids:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
            else:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target)
        except Exception:
            pass
    slide_id_list = prs.slides._sldIdLst
    new_id = slide_id_list[-1]
    del slide_id_list[-1]
    slide_id_list.insert(insert_after_idx + 1, new_id)
    return prs.slides[insert_after_idx + 1]


def remove_slide(prs: Presentation, slide_index: int):
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    prs.part.drop_rel(slide_id.rId)
    del slide_id_list[slide_index]


def ensure_last_slide(prs: Presentation, slide_id: int):
    slide_id_list = prs.slides._sldIdLst
    target = None
    for item in slide_id_list:
        if int(item.id) == int(slide_id):
            target = item
            break
    if target is None:
        return
    slide_id_list.remove(target)
    slide_id_list.append(target)


def slide_image_targets(pptx_path: Path, slide_no: int) -> set[str]:
    rel_path = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    targets: set[str] = set()
    with zipfile.ZipFile(pptx_path) as package:
        try:
            root = ElementTree.fromstring(package.read(rel_path))
        except KeyError:
            return targets
    for rel in root:
        rel_type = rel.attrib.get("Type", "")
        if rel_type.endswith("/image"):
            target = rel.attrib.get("Target", "")
            if target:
                targets.add(target)
    return targets


def slide_assets_preserved(pptx_path: Path, source_idx: int, target_indices: list[int]) -> bool:
    source_targets = slide_image_targets(pptx_path, source_idx)
    if not source_targets:
        return True
    for slide_no in target_indices:
        target_assets = slide_image_targets(pptx_path, slide_no)
        if not source_targets.issubset(target_assets):
            return False
    return True


def _normalize_package_name(name: str) -> str:
    normalized = name.replace("\\", "/").lstrip("/")
    return normpath(normalized)


def _rels_part_name(part_name: str) -> str:
    part_name = _normalize_package_name(part_name)
    part_dir = dirname(part_name)
    if part_dir:
        return join(part_dir, "_rels", f"{basename(part_name)}.rels")
    return join("_rels", f"{basename(part_name)}.rels")


def _resolve_internal_target(base_part_name: str, target: str) -> str:
    if target.startswith("/"):
        return _normalize_package_name(target)
    return _normalize_package_name(join(dirname(_normalize_package_name(base_part_name)), target))


def _build_relative_target(base_part_name: str, target_part_name: str) -> str:
    base_dir = dirname(_normalize_package_name(base_part_name)) or "."
    return relpath(_normalize_package_name(target_part_name), start=base_dir).replace("\\", "/")


def _parse_content_types(package: zipfile.ZipFile):
    root = ElementTree.fromstring(package.read("[Content_Types].xml"))
    defaults: dict[str, str] = {}
    overrides: dict[str, str] = {}
    for child in root:
        local_name = child.tag.rsplit("}", 1)[-1]
        if local_name == "Default":
            extension = child.attrib.get("Extension", "").lower()
            content_type = child.attrib.get("ContentType", "")
            if extension and content_type:
                defaults[extension] = content_type
        elif local_name == "Override":
            part_name = child.attrib.get("PartName", "")
            content_type = child.attrib.get("ContentType", "")
            if part_name and content_type:
                overrides[_normalize_package_name(part_name)] = content_type
    return root, defaults, overrides


def _ensure_content_type_entry(
    content_types_root,
    target_defaults: dict[str, str],
    target_overrides: dict[str, str],
    source_defaults: dict[str, str],
    source_overrides: dict[str, str],
    source_part_name: str,
    target_part_name: str,
):
    source_part_name = _normalize_package_name(source_part_name)
    target_part_name = _normalize_package_name(target_part_name)
    source_override = source_overrides.get(source_part_name)
    if source_override:
        if target_part_name not in target_overrides:
            ElementTree.SubElement(
                content_types_root,
                f"{{{CONTENT_TYPES_NS}}}Override",
                {
                    "PartName": f"/{target_part_name}",
                    "ContentType": source_override,
                },
            )
            target_overrides[target_part_name] = source_override
        return

    extension = splitext(source_part_name)[1].lstrip(".").lower()
    if not extension:
        return
    source_default = source_defaults.get(extension)
    if source_default and extension not in target_defaults:
        ElementTree.SubElement(
            content_types_root,
            f"{{{CONTENT_TYPES_NS}}}Default",
            {
                "Extension": extension,
                "ContentType": source_default,
            },
        )
        target_defaults[extension] = source_default


def _select_target_part_name(
    target_package: zipfile.ZipFile,
    new_parts: dict[str, bytes],
    occupied_names: set[str],
    desired_name: str,
    source_bytes: bytes,
) -> str:
    desired_name = _normalize_package_name(desired_name)
    if desired_name not in occupied_names:
        occupied_names.add(desired_name)
        return desired_name

    existing_bytes = new_parts.get(desired_name)
    if existing_bytes is None:
        try:
            existing_bytes = target_package.read(desired_name)
        except KeyError:
            existing_bytes = None
    if existing_bytes == source_bytes:
        return desired_name

    part_dir = dirname(desired_name)
    part_stem, part_ext = splitext(basename(desired_name))
    suffix = 1
    while True:
        candidate = join(part_dir, f"{part_stem}_import{suffix}{part_ext}") if part_dir else f"{part_stem}_import{suffix}{part_ext}"
        existing_bytes = new_parts.get(candidate)
        if existing_bytes is None:
            try:
                existing_bytes = target_package.read(candidate)
            except KeyError:
                existing_bytes = None
        if candidate not in occupied_names:
            occupied_names.add(candidate)
            return candidate
        if existing_bytes == source_bytes:
            return candidate
        suffix += 1


def _copy_related_part_graph(
    source_package: zipfile.ZipFile,
    target_package: zipfile.ZipFile,
    source_defaults: dict[str, str],
    source_overrides: dict[str, str],
    content_types_root,
    target_defaults: dict[str, str],
    target_overrides: dict[str, str],
    source_part_name: str,
    new_parts: dict[str, bytes],
    rel_updates: dict[str, bytes],
    copied_parts: dict[str, str],
    occupied_names: set[str],
) -> str:
    source_part_name = _normalize_package_name(source_part_name)
    existing_target_name = copied_parts.get(source_part_name)
    if existing_target_name:
        return existing_target_name

    source_bytes = source_package.read(source_part_name)
    target_part_name = _select_target_part_name(
        target_package=target_package,
        new_parts=new_parts,
        occupied_names=occupied_names,
        desired_name=source_part_name,
        source_bytes=source_bytes,
    )
    copied_parts[source_part_name] = target_part_name
    new_parts[target_part_name] = source_bytes
    _ensure_content_type_entry(
        content_types_root,
        target_defaults,
        target_overrides,
        source_defaults,
        source_overrides,
        source_part_name,
        target_part_name,
    )

    source_rel_name = _rels_part_name(source_part_name)
    if source_rel_name not in source_package.namelist():
        return target_part_name

    source_rel_root = ElementTree.fromstring(source_package.read(source_rel_name))
    target_rel_root = ElementTree.Element(source_rel_root.tag, source_rel_root.attrib)
    for rel in source_rel_root:
        rel_type = rel.attrib.get("Type", "")
        if any(rel_type.endswith(suffix) for suffix in SKIPPED_RELATIONSHIP_SUFFIXES):
            continue
        rel_copy = deepcopy(rel)
        target_mode = rel.attrib.get("TargetMode")
        rel_target = rel.attrib.get("Target", "")
        if target_mode == "External" or not rel_target:
            target_rel_root.append(rel_copy)
            continue

        child_source_part = _resolve_internal_target(source_part_name, rel_target)
        if child_source_part not in source_package.namelist():
            target_rel_root.append(rel_copy)
            continue
        child_target_part = _copy_related_part_graph(
            source_package=source_package,
            target_package=target_package,
            source_defaults=source_defaults,
            source_overrides=source_overrides,
            content_types_root=content_types_root,
            target_defaults=target_defaults,
            target_overrides=target_overrides,
            source_part_name=child_source_part,
            new_parts=new_parts,
            rel_updates=rel_updates,
            copied_parts=copied_parts,
            occupied_names=occupied_names,
        )
        rel_copy.attrib["Target"] = _build_relative_target(target_part_name, child_target_part)
        target_rel_root.append(rel_copy)

    rel_updates[_rels_part_name(target_part_name)] = ElementTree.tostring(
        target_rel_root,
        encoding="utf-8",
        xml_declaration=True,
    )
    return target_part_name


def _build_imported_slide_relationships(
    source_package: zipfile.ZipFile,
    target_package: zipfile.ZipFile,
    source_defaults: dict[str, str],
    source_overrides: dict[str, str],
    content_types_root,
    target_defaults: dict[str, str],
    target_overrides: dict[str, str],
    source_slide_name: str,
    target_slide_name: str,
    new_parts: dict[str, bytes],
    rel_updates: dict[str, bytes],
    copied_parts: dict[str, str],
    occupied_names: set[str],
    target_slide_rel_bytes: bytes | None,
) -> bytes:
    source_rel_name = _rels_part_name(source_slide_name)
    target_rel_root = ElementTree.Element(f"{{{PACKAGE_REL_NS}}}Relationships")
    if target_slide_rel_bytes:
        existing_target_rel_root = ElementTree.fromstring(target_slide_rel_bytes)
        for rel in existing_target_rel_root:
            rel_type = rel.attrib.get("Type", "")
            if rel_type.endswith("/slideLayout"):
                target_rel_root.append(deepcopy(rel))

    if source_rel_name not in source_package.namelist():
        return ElementTree.tostring(target_rel_root, encoding="utf-8", xml_declaration=True)

    source_rel_root = ElementTree.fromstring(source_package.read(source_rel_name))
    for rel in source_rel_root:
        rel_type = rel.attrib.get("Type", "")
        if any(rel_type.endswith(suffix) for suffix in SKIPPED_RELATIONSHIP_SUFFIXES):
            continue
        if rel_type.endswith("/slideLayout"):
            continue
        rel_copy = deepcopy(rel)
        target_mode = rel.attrib.get("TargetMode")
        rel_target = rel.attrib.get("Target", "")
        if target_mode == "External" or not rel_target:
            target_rel_root.append(rel_copy)
            continue

        source_part_name = _resolve_internal_target(source_slide_name, rel_target)
        if source_part_name not in source_package.namelist():
            target_rel_root.append(rel_copy)
            continue
        target_part_name = _copy_related_part_graph(
            source_package=source_package,
            target_package=target_package,
            source_defaults=source_defaults,
            source_overrides=source_overrides,
            content_types_root=content_types_root,
            target_defaults=target_defaults,
            target_overrides=target_overrides,
            source_part_name=source_part_name,
            new_parts=new_parts,
            rel_updates=rel_updates,
            copied_parts=copied_parts,
            occupied_names=occupied_names,
        )
        rel_copy.attrib["Target"] = _build_relative_target(target_slide_name, target_part_name)
        target_rel_root.append(rel_copy)

    return ElementTree.tostring(
        target_rel_root,
        encoding="utf-8",
        xml_declaration=True,
    )


def _replace_package_with_retries(source_path: Path, rebuilt_path: Path):
    last_error = None
    for _ in range(10):
        try:
            rebuilt_path.replace(source_path)
            return
        except PermissionError as exc:
            last_error = exc
            time.sleep(0.5)
    if rebuilt_path.exists():
        rebuilt_path.unlink(missing_ok=True)
    if last_error is not None:
        raise last_error


def import_slides_from_presentation(target_pptx: Path, source_pptx: Path, mappings: list[tuple[int, int]]) -> bool:
    if not mappings:
        return True

    rebuilt_path = target_pptx.with_name(target_pptx.stem + "_rebuilt.pptx")
    with zipfile.ZipFile(target_pptx, "r") as target_package, zipfile.ZipFile(source_pptx, "r") as source_package:
        content_types_root, target_defaults, target_overrides = _parse_content_types(target_package)
        _, source_defaults, source_overrides = _parse_content_types(source_package)
        occupied_names = set(target_package.namelist())
        slide_updates: dict[str, bytes] = {}
        rel_updates: dict[str, bytes] = {}
        new_parts: dict[str, bytes] = {}
        copied_parts: dict[str, str] = {}

        for target_slide_no, source_slide_no in mappings:
            source_slide_name = f"ppt/slides/slide{source_slide_no}.xml"
            target_slide_name = f"ppt/slides/slide{target_slide_no}.xml"
            if source_slide_name not in source_package.namelist() or target_slide_name not in target_package.namelist():
                return False

            target_slide_rel_name = _rels_part_name(target_slide_name)
            target_slide_rel_bytes = (
                target_package.read(target_slide_rel_name)
                if target_slide_rel_name in target_package.namelist()
                else None
            )
            slide_updates[target_slide_name] = source_package.read(source_slide_name)
            _ensure_content_type_entry(
                content_types_root,
                target_defaults,
                target_overrides,
                source_defaults,
                source_overrides,
                source_slide_name,
                target_slide_name,
            )
            rel_updates[_rels_part_name(target_slide_name)] = _build_imported_slide_relationships(
                source_package=source_package,
                target_package=target_package,
                source_defaults=source_defaults,
                source_overrides=source_overrides,
                content_types_root=content_types_root,
                target_defaults=target_defaults,
                target_overrides=target_overrides,
                source_slide_name=source_slide_name,
                target_slide_name=target_slide_name,
                new_parts=new_parts,
                rel_updates=rel_updates,
                copied_parts=copied_parts,
                occupied_names=occupied_names,
                target_slide_rel_bytes=target_slide_rel_bytes,
            )

        content_types_bytes = ElementTree.tostring(
            content_types_root,
            encoding="utf-8",
            xml_declaration=True,
        )

        with zipfile.ZipFile(rebuilt_path, "w", zipfile.ZIP_DEFLATED) as rebuilt:
            for info in target_package.infolist():
                data = target_package.read(info.filename)
                if info.filename == "[Content_Types].xml":
                    data = content_types_bytes
                elif info.filename in slide_updates:
                    data = slide_updates[info.filename]
                elif info.filename in rel_updates:
                    data = rel_updates[info.filename]
                elif info.filename in new_parts:
                    data = new_parts[info.filename]
                rebuilt.writestr(info, data)

            for collection in (slide_updates, rel_updates, new_parts):
                for filename, data in collection.items():
                    if filename not in target_package.namelist():
                        rebuilt.writestr(filename, data)

    _replace_package_with_retries(target_pptx, rebuilt_path)
    return True


def set_slide_metadata_names(pptx_path: Path, name_by_slide_no: dict[int, str]) -> bool:
    if not name_by_slide_no:
        return True

    rebuilt_path = pptx_path.with_name(pptx_path.stem + "_rebuilt.pptx")
    slide_updates: dict[str, bytes] = {}

    with zipfile.ZipFile(pptx_path, "r") as package:
        for slide_no, slide_name in name_by_slide_no.items():
            filename = f"ppt/slides/slide{slide_no}.xml"
            if filename not in package.namelist():
                return False
            root = ElementTree.fromstring(package.read(filename))
            c_sld = root.find(f"{{{PRESENTATION_NS}}}cSld")
            if c_sld is None:
                continue
            c_sld.set("name", slide_name)
            slide_updates[filename] = ElementTree.tostring(
                root,
                encoding="utf-8",
                xml_declaration=True,
            )

        with zipfile.ZipFile(rebuilt_path, "w", zipfile.ZIP_DEFLATED) as rebuilt:
            for info in package.infolist():
                data = package.read(info.filename)
                if info.filename in slide_updates:
                    data = slide_updates[info.filename]
                rebuilt.writestr(info, data)

    _replace_package_with_retries(pptx_path, rebuilt_path)
    return True


def copy_slide_xml_assets(pptx_path: Path, source_idx: int, target_indices: list[int]) -> bool:
    if not target_indices:
        return True

    source_slide_name = f"ppt/slides/slide{source_idx}.xml"
    source_rel_name = f"ppt/slides/_rels/slide{source_idx}.xml.rels"
    target_slide_names = {f"ppt/slides/slide{target}.xml" for target in target_indices}
    target_rel_names = {f"ppt/slides/_rels/slide{target}.xml.rels" for target in target_indices}
    rebuilt_path = pptx_path.with_name(pptx_path.stem + "_rebuilt.pptx")

    with zipfile.ZipFile(pptx_path, "r") as source_package:
        if source_slide_name not in source_package.namelist():
            return False

        slide_root = ElementTree.fromstring(source_package.read(source_slide_name))
        slide_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        source_sp_tree = slide_root.find(f".//{slide_ns}spTree")
        source_pics = [deepcopy(pic) for pic in source_sp_tree.findall(f"{slide_ns}pic")] if source_sp_tree is not None else []
        rel_bytes = source_package.read(source_rel_name) if source_rel_name in source_package.namelist() else None
        source_image_rels = []
        if rel_bytes is not None:
            source_rel_root = ElementTree.fromstring(rel_bytes)
            source_image_rels = [
                deepcopy(rel)
                for rel in source_rel_root
                if rel.attrib.get("Type", "").endswith("/image")
            ]

        slide_replacements: dict[str, bytes] = {}
        rel_replacements: dict[str, bytes] = {}
        rel_root_tag = (
            ElementTree.fromstring(rel_bytes).tag
            if rel_bytes is not None
            else "{http://schemas.openxmlformats.org/package/2006/relationships}Relationships"
        )

        for target_slide_name in target_slide_names:
            if target_slide_name not in source_package.namelist():
                continue
            target_root = ElementTree.fromstring(source_package.read(target_slide_name))
            target_sp_tree = target_root.find(f".//{slide_ns}spTree")
            if target_sp_tree is not None and not target_sp_tree.findall(f"{slide_ns}pic") and source_pics:
                insert_at = next(
                    (index for index, child in enumerate(list(target_sp_tree)) if child.tag == f"{slide_ns}extLst"),
                    len(target_sp_tree),
                )
                for pic in source_pics:
                    target_sp_tree.insert(insert_at, deepcopy(pic))
                    insert_at += 1
            slide_replacements[target_slide_name] = ElementTree.tostring(
                target_root,
                encoding="utf-8",
                xml_declaration=True,
            )

        for target_rel_name in target_rel_names:
            if target_rel_name in source_package.namelist():
                target_rel_root = ElementTree.fromstring(source_package.read(target_rel_name))
            else:
                target_rel_root = ElementTree.Element(rel_root_tag)
            existing_image_targets = {
                rel.attrib.get("Target", "")
                for rel in target_rel_root
                if rel.attrib.get("Type", "").endswith("/image")
            }
            for image_rel in source_image_rels:
                if image_rel.attrib.get("Target", "") in existing_image_targets:
                    continue
                target_rel_root.append(deepcopy(image_rel))
            rel_replacements[target_rel_name] = ElementTree.tostring(
                target_rel_root,
                encoding="utf-8",
                xml_declaration=True,
            )

        with zipfile.ZipFile(rebuilt_path, "w", zipfile.ZIP_DEFLATED) as rebuilt:
            for info in source_package.infolist():
                data = source_package.read(info.filename)
                if info.filename in slide_replacements:
                    data = slide_replacements[info.filename]
                elif info.filename in rel_replacements:
                    data = rel_replacements[info.filename]
                rebuilt.writestr(info, data)

            for target_slide_name, data in slide_replacements.items():
                if target_slide_name not in source_package.namelist():
                    rebuilt.writestr(target_slide_name, data)
            for target_rel_name, data in rel_replacements.items():
                if target_rel_name not in source_package.namelist():
                    rebuilt.writestr(target_rel_name, data)

    _replace_package_with_retries(pptx_path, rebuilt_path)
    return True
