from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..theme_loader import ThemeSpec


@dataclass(frozen=True)
class RenderContext:
    prs: Any
    theme: ThemeSpec
    log: Any
    slide_number: int
    total_slides: int


TIMELINE_PREFIX_PATTERN = re.compile(
    r"^(第[一二三四五六七八九十0-9]+阶段|Q\d{1,2}|H[12]|Phase\s*\d+|Step\s*\d+|[^\s，,:：]{1,8}阶段)\s*",
    re.IGNORECASE,
)
ARCHITECTURE_KEYWORDS = ("架构", "框架", "平台", "architecture", "framework")
MAP_KEYWORDS = ("地图", "场景", "map")


def rgb(hex_color: str) -> RGBColor:
    hex_value = hex_color.strip().lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_background(slide, theme: ThemeSpec, color_hex: str | None = None) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        Inches(theme.page.width),
        Inches(theme.page.height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color_hex or theme.colors.bg)
    shape.line.fill.background()


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    font_name: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    vertical_anchor=MSO_ANCHOR.TOP,
    margin_left: float = 0.06,
    margin_right: float = 0.06,
    margin_top: float = 0.04,
    margin_bottom: float = 0.04,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical_anchor
    frame.margin_left = Inches(margin_left)
    frame.margin_right = Inches(margin_right)
    frame.margin_top = Inches(margin_top)
    frame.margin_bottom = Inches(margin_bottom)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color_hex)
    return box


def add_card(slide, left: float, top: float, width: float, height: float, theme: ThemeSpec):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = rgb(theme.colors.card_bg)
    card.line.color.rgb = rgb(theme.colors.line)
    card.line.width = Pt(1.1)
    return card


def add_bullet_list(
    slide,
    items: list[str],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
    font_size: int,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Inches(0.1)
    frame.margin_right = Inches(0.1)
    frame.margin_top = Inches(0.06)
    frame.margin_bottom = Inches(0.06)

    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(6)
        paragraph.font.name = theme.fonts.body
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = rgb(theme.colors.text_main)
    return box


def add_page_number(slide, slide_number: int, total_slides: int, theme: ThemeSpec) -> None:
    add_textbox(
        slide,
        left=theme.page.width - theme.spacing.page_margin_right - 0.7,
        top=theme.page.height - theme.spacing.page_margin_bottom - 0.22,
        width=0.65,
        height=0.18,
        text=f"{slide_number}/{total_slides}",
        font_name=theme.fonts.body,
        font_size=theme.font_sizes.small,
        color_hex=theme.colors.text_sub,
        align=PP_ALIGN.RIGHT,
    )


def resolve_body_font_size(theme: ThemeSpec, item_count: int) -> int:
    base = theme.font_sizes.body
    if item_count <= 4:
        return base
    if item_count == 5:
        return max(theme.font_sizes.small + 2, base - 1)
    if item_count == 6:
        return max(theme.font_sizes.small + 1, base - 2)
    return max(theme.font_sizes.small, base - 3)


def split_label_detail(text: str) -> tuple[str, str]:
    normalized = re.sub(r"\s+", " ", str(text).strip())
    for sep in ("：", ":", "，", ",", " - ", " | ", "；"):
        if sep in normalized:
            left, right = normalized.split(sep, 1)
            left = left.strip()
            right = right.strip()
            if left and right:
                return left, right
    return normalized, ""


def parse_timeline_items(items: list[str]) -> list[tuple[str, str]] | None:
    if not 3 <= len(items) <= 5:
        return None

    parsed: list[tuple[str, str]] = []
    for item in items:
        normalized = re.sub(r"\s+", " ", str(item).strip())
        if not normalized:
            return None

        match = TIMELINE_PREFIX_PATTERN.match(normalized)
        if match:
            label = match.group(1).strip()
            detail = normalized[match.end() :].lstrip(" ：:，,、-")
            if not detail:
                return None
            parsed.append((label, detail))
            continue

        label, detail = split_label_detail(normalized)
        if label != normalized and TIMELINE_PREFIX_PATTERN.match(label):
            parsed.append((label, detail))
            continue

        return None

    return parsed


def should_render_comparison_table(left_heading: str, right_heading: str, left_items: list[str], right_items: list[str]) -> bool:
    if len(left_items) != len(right_items):
        return False
    if not 2 <= len(left_items) <= 4:
        return False
    return bool(left_heading.strip() and right_heading.strip())


def classify_placeholder_visual(title: str, caption: str | None, content: list[str]) -> str:
    combined = f"{title} {caption or ''}".lower()
    if len(content) >= 2 and any(keyword in combined for keyword in ARCHITECTURE_KEYWORDS):
        return "architecture"
    if len(content) >= 2 and any(keyword in combined for keyword in MAP_KEYWORDS):
        return "map"
    return "placeholder"


def add_timeline_flow(
    slide,
    stages: list[tuple[str, str]],
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
) -> None:
    count = len(stages)
    gap = 0.18
    card_width = max(1.8, (width - gap * (count - 1)) / count)
    line_y = top + 0.8
    line_left = left + 0.28
    line_width = width - 0.56

    track = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(line_left),
        Inches(line_y),
        Inches(line_width),
        Inches(0.05),
    )
    track.fill.solid()
    track.fill.fore_color.rgb = rgb(theme.colors.line)
    track.line.fill.background()

    for index, (label, detail) in enumerate(stages):
        card_left = left + index * (card_width + gap)
        center_x = card_left + card_width / 2

        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(center_x - 0.11),
            Inches(line_y - 0.085),
            Inches(0.22),
            Inches(0.22),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(theme.colors.primary)
        dot.line.color.rgb = rgb(theme.colors.primary)

        connector = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(center_x - 0.01),
            Inches(line_y + 0.12),
            Inches(0.02),
            Inches(0.28),
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = rgb(theme.colors.line)
        connector.line.fill.background()

        badge = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(card_left + 0.2),
            Inches(top),
            Inches(card_width - 0.4),
            Inches(0.34),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = rgb(theme.colors.primary)
        badge.line.color.rgb = rgb(theme.colors.primary)
        add_textbox(
            slide,
            left=card_left + 0.2,
            top=top + 0.02,
            width=card_width - 0.4,
            height=0.28,
            text=label,
            font_name=theme.fonts.body,
            font_size=theme.font_sizes.small + 1,
            color_hex=theme.colors.card_bg,
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
            margin_left=0.02,
            margin_right=0.02,
            margin_top=0.0,
            margin_bottom=0.0,
        )

        add_card(slide, card_left, top + 1.18, card_width, height - 1.18, theme)
        add_textbox(
            slide,
            left=card_left + 0.12,
            top=top + 1.34,
            width=card_width - 0.24,
            height=height - 1.5,
            text=detail,
            font_name=theme.fonts.body,
            font_size=theme.font_sizes.small + 1,
            color_hex=theme.colors.text_main,
            bold=False,
            align=PP_ALIGN.LEFT,
            margin_left=0.04,
            margin_right=0.04,
            margin_top=0.02,
            margin_bottom=0.02,
        )


def add_comparison_table(
    slide,
    *,
    left_heading: str,
    right_heading: str,
    left_items: list[str],
    right_items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
) -> None:
    add_card(slide, left, top, width, height, theme)
    header_height = 0.48
    col_gap = 0.1
    col_width = (width - col_gap) / 2
    row_count = len(left_items)
    row_gap = 0.08
    row_height = (height - header_height - 0.22 - row_gap * row_count) / row_count

    for idx, (heading, col_left) in enumerate(((left_heading, left), (right_heading, left + col_width + col_gap))):
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(col_left + 0.08),
            Inches(top + 0.08),
            Inches(col_width - 0.16),
            Inches(header_height),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(theme.colors.primary if idx == 0 else theme.colors.secondary)
        band.line.fill.background()
        add_textbox(
            slide,
            left=col_left + 0.12,
            top=top + 0.16,
            width=col_width - 0.24,
            height=0.28,
            text=heading,
            font_name=theme.fonts.title,
            font_size=theme.font_sizes.subtitle,
            color_hex=theme.colors.card_bg,
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
            margin_left=0.02,
            margin_right=0.02,
            margin_top=0.0,
            margin_bottom=0.0,
        )

    cell_top = top + header_height + 0.16
    for row_index, (left_text, right_text) in enumerate(zip(left_items, right_items)):
        y = cell_top + row_index * (row_height + row_gap)
        for col_index, (col_left, text) in enumerate(((left, left_text), (left + col_width + col_gap, right_text))):
            cell = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(col_left + 0.08),
                Inches(y),
                Inches(col_width - 0.16),
                Inches(row_height),
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(theme.colors.card_bg if col_index == 0 else theme.colors.bg)
            cell.line.color.rgb = rgb(theme.colors.line)
            add_textbox(
                slide,
                left=col_left + 0.16,
                top=y + 0.08,
                width=col_width - 0.32,
                height=row_height - 0.14,
                text=text,
                font_name=theme.fonts.body,
                font_size=max(theme.font_sizes.small + 1, theme.font_sizes.body - 1),
                color_hex=theme.colors.text_main,
                bold=False,
                align=PP_ALIGN.LEFT,
                margin_left=0.03,
                margin_right=0.03,
                margin_top=0.01,
                margin_bottom=0.01,
            )


def add_architecture_placeholder(
    slide,
    *,
    caption: str | None,
    layers: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
) -> None:
    add_textbox(
        slide,
        left=left + 0.15,
        top=top + 0.08,
        width=width - 0.3,
        height=0.28,
        text=caption or "Architecture Overview",
        font_name=theme.fonts.body,
        font_size=theme.font_sizes.small + 1,
        color_hex=theme.colors.text_sub,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    layer_count = len(layers)
    gap = 0.14
    layer_height = min(0.82, (height - 0.72 - gap * (layer_count - 1)) / max(layer_count, 1))
    palette = [theme.colors.card_bg, theme.colors.bg, theme.colors.card_bg, theme.colors.bg]

    for index, layer_text in enumerate(layers[:4]):
        y = top + 0.48 + index * (layer_height + gap)
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left + 0.35),
            Inches(y),
            Inches(width - 0.7),
            Inches(layer_height),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = rgb(palette[index % len(palette)])
        box.line.color.rgb = rgb(theme.colors.line)

        tag = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left + 0.52),
            Inches(y + 0.12),
            Inches(0.7),
            Inches(layer_height - 0.24),
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = rgb(theme.colors.primary)
        tag.line.fill.background()
        add_textbox(
            slide,
            left=left + 0.52,
            top=y + 0.16,
            width=0.7,
            height=layer_height - 0.32,
            text=f"L{index + 1}",
            font_name=theme.fonts.body,
            font_size=theme.font_sizes.small,
            color_hex=theme.colors.card_bg,
            bold=True,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
            margin_left=0.0,
            margin_right=0.0,
            margin_top=0.0,
            margin_bottom=0.0,
        )
        add_textbox(
            slide,
            left=left + 1.35,
            top=y + 0.1,
            width=width - 1.75,
            height=layer_height - 0.18,
            text=layer_text,
            font_name=theme.fonts.body,
            font_size=theme.font_sizes.small + 1,
            color_hex=theme.colors.text_main,
            bold=True if index == 0 else False,
            align=PP_ALIGN.LEFT,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
            margin_left=0.02,
            margin_right=0.02,
            margin_top=0.0,
            margin_bottom=0.0,
        )


def add_capability_map(
    slide,
    *,
    caption: str | None,
    nodes: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
) -> None:
    center = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left + width / 2 - 0.95),
        Inches(top + 0.38),
        Inches(1.9),
        Inches(1.05),
    )
    center.fill.solid()
    center.fill.fore_color.rgb = rgb(theme.colors.primary)
    center.line.fill.background()
    add_textbox(
        slide,
        left=left + width / 2 - 0.83,
        top=top + 0.67,
        width=1.66,
        height=0.36,
        text=caption or "Map",
        font_name=theme.fonts.body,
        font_size=theme.font_sizes.small + 1,
        color_hex=theme.colors.card_bg,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
        margin_left=0.0,
        margin_right=0.0,
        margin_top=0.0,
        margin_bottom=0.0,
    )

    positions = [
        (left + 0.2, top + 1.9),
        (left + width / 2 - 1.15, top + 2.55),
        (left + width - 2.5, top + 1.9),
        (left + width / 2 - 1.15, top + 1.7),
    ]
    for index, node in enumerate(nodes[:4]):
        card_left, card_top = positions[index]
        connector = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(card_left + 1.1),
            Inches(top + 1.22 if index != 1 else top + 1.42),
            Inches(0.02),
            Inches(max(0.35, abs(card_top - (top + 1.28)))),
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = rgb(theme.colors.line)
        connector.line.fill.background()

        add_card(slide, card_left, card_top, 2.3, 0.92, theme)
        add_textbox(
            slide,
            left=card_left + 0.12,
            top=card_top + 0.16,
            width=2.06,
            height=0.56,
            text=node,
            font_name=theme.fonts.body,
            font_size=theme.font_sizes.small + 1,
            color_hex=theme.colors.text_main,
            bold=False,
            align=PP_ALIGN.CENTER,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
            margin_left=0.02,
            margin_right=0.02,
            margin_top=0.0,
            margin_bottom=0.0,
        )


def add_local_image_or_placeholder(
    slide,
    *,
    image_mode: str,
    image_path: str | None,
    caption: str | None,
    left: float,
    top: float,
    width: float,
    height: float,
    theme: ThemeSpec,
    log,
) -> None:
    path = Path(image_path) if image_path else None
    if image_mode == "local_path" and path and path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))
        return

    if image_mode == "local_path":
        log.warn(f"image file not found, using placeholder instead: {image_path}")
    placeholder = add_card(slide, left, top, width, height, theme)
    placeholder.fill.fore_color.rgb = rgb(theme.colors.bg)
    placeholder.line.color.rgb = rgb(theme.colors.line)
    add_textbox(
        slide,
        left=left + 0.2,
        top=top + height / 2 - 0.3,
        width=width - 0.4,
        height=0.6,
        text=caption or "Image Placeholder",
        font_name=theme.fonts.body,
        font_size=theme.font_sizes.body,
        color_hex=theme.colors.text_sub,
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )
