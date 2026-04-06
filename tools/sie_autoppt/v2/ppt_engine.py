from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from .content_rewriter import rewrite_deck, write_rewrite_log
from .io import RenderLog, write_deck_document
from .layout_router import render_slide
from .quality_checks import ContentWarning, QualityGateResult, quality_gate, write_quality_gate_result
from .schema import DeckDocument, ValidatedDeck
from .theme_loader import load_theme


@dataclass(frozen=True)
class RenderArtifacts:
    output_path: Path
    log_path: Path | None
    warnings_path: Path | None
    rewrite_log_path: Path | None
    deck_path: Path | None
    slide_count: int
    quality_gate: QualityGateResult
    final_deck: DeckDocument
    warnings: tuple[str, ...] = ()
    content_warnings: tuple[ContentWarning, ...] = ()
    warning_count: int = 0
    high_count: int = 0
    error_count: int = 0
    review_required: bool = False
    auto_score: int = 100
    auto_level: str = "优秀"
    rewrite_applied: bool = False
def generate_ppt(
    deck_data: DeckDocument | ValidatedDeck | dict[str, object],
    output_path: str | Path,
    theme_name: str | None = None,
    log_path: str | Path | None = None,
    deck_output_path: str | Path | None = None,
    rewrite_log_path: str | Path | None = None,
    max_errors: int = 0,
) -> RenderArtifacts:
    """
    Generate a PowerPoint presentation from a deck specification.

    Args:
        deck_data: Deck specification (DeckDocument, ValidatedDeck, or dict)
        output_path: Path where the PPTX file will be saved
        theme_name: Optional theme name override
        log_path: Optional path to write the generation log
        deck_output_path: Optional path to write the validated or rewritten deck JSON
        rewrite_log_path: Optional path to write rewrite decisions
        max_errors: Maximum number of error-level quality issues allowed (default: 0)
                   If error count exceeds this threshold, generation will be aborted.

    Returns:
        RenderArtifacts containing output paths and quality warnings

    Raises:
        ValueError: If error count exceeds max_errors threshold
    """
    log = RenderLog()
    warnings_json_path = (Path(log_path).parent / "warnings.json") if log_path else (Path(output_path).parent / "warnings.json")
    final_rewrite_log_path = Path(rewrite_log_path) if rewrite_log_path else (
        (Path(log_path).parent / "rewrite_log.json") if log_path else (Path(output_path).parent / "rewrite_log.json")
    )
    final_deck_path = Path(deck_output_path) if deck_output_path else None

    quality_result = quality_gate(deck_data)
    rewrite_result = rewrite_deck(quality_result.validated_deck, quality_result)
    write_rewrite_log(rewrite_result, final_rewrite_log_path)
    quality_result = rewrite_result.final_quality_gate
    write_quality_gate_result(quality_result, warnings_json_path)

    validated = quality_result.validated_deck
    if validated is None:
        for issue in quality_result.errors:
            log.error(issue.to_log_line())
        error_msg = "Schema validation failed; renderer was skipped."
        log.error(error_msg)
        final_log_path = Path(log_path) if log_path else None
        if final_log_path is not None:
            log.write(final_log_path)
        raise ValueError(error_msg)

    deck = validated.deck
    if final_deck_path is not None:
        write_deck_document(deck, final_deck_path)

    theme = load_theme(theme_name or deck.meta.theme)
    log.info(f"deck title: {deck.meta.title}")
    log.info(f"theme: {theme.theme_name}")
    log.extend(validated.warnings)
    if rewrite_result.applied:
        log.info(f"content rewrite applied: {len(rewrite_result.actions)} change(s)")

    content_warnings = quality_result.all_issues()
    for warning in content_warnings:
        if warning.warning_level == "error":
            log.error(warning.to_log_line())
        else:
            log.warn(warning.to_log_line())

    error_count = quality_result.summary["error_count"]
    if not quality_result.passed or error_count > max_errors:
        error_msg = f"Quality gate failed: {error_count} error(s) found. Renderer was skipped."
        log.error(error_msg)
        final_log_path = Path(log_path) if log_path else None
        if final_log_path is not None:
            log.write(final_log_path)
        raise ValueError(error_msg)

    prs = Presentation()
    prs.slide_width = Inches(theme.page.width)
    prs.slide_height = Inches(theme.page.height)

    for index, slide in enumerate(deck.slides, start=1):
        render_slide(prs, slide, theme, log, index, len(deck.slides))

    final_output = Path(output_path)
    final_output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(final_output))

    final_log_path = Path(log_path) if log_path else None
    if final_log_path is not None:
        log.write(final_log_path)

    return RenderArtifacts(
        output_path=final_output,
        log_path=final_log_path,
        warnings_path=warnings_json_path,
        rewrite_log_path=final_rewrite_log_path,
        deck_path=final_deck_path,
        slide_count=len(deck.slides),
        quality_gate=quality_result,
        final_deck=deck,
        warnings=tuple(validated.warnings),
        content_warnings=content_warnings,
        warning_count=quality_result.summary["warning_count"],
        high_count=quality_result.summary["high_count"],
        error_count=error_count,
        review_required=quality_result.review_required,
        auto_score=quality_result.auto_score,
        auto_level=quality_result.auto_level,
        rewrite_applied=rewrite_result.applied,
    )
