import re
from pathlib import Path

from pptx import Presentation


def _slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = re.sub(r"\s+", " ", shape.text_frame.text).strip()
            if text:
                texts.append(text)
    return " ".join(texts)


def _is_ending_slide(slide) -> bool:
    text = _slide_text(slide).lower()
    if not text and len(slide.shapes) <= 2:
        return True
    return any(keyword in text for keyword in ("thanks", "thank you", "q&a", "谢谢", "感谢"))


def _is_directory_slide(slide, chapter_lines: list[str]) -> bool:
    text = _slide_text(slide)
    if not text:
        return False
    lowered = text.lower()
    if "目录" in text or "content" in lowered:
        return True
    matched = sum(1 for line in chapter_lines if line and line in text)
    return matched >= 2


def write_qa_report(pptx_path: Path, chapter_count: int, pattern_ids=None, chapter_lines=None) -> Path:
    prs = Presentation(str(pptx_path))
    chapter_lines = chapter_lines or []
    lines = []
    lines.append("SIE AutoPPT QA Report")
    lines.append(f"file: {pptx_path}")
    lines.append(f"slides: {len(prs.slides)}")

    has_ending_last = _is_ending_slide(prs.slides[-1])
    lines.append(f"check_ending_last: {'PASS' if has_ending_last else 'WARN'}")

    expected_dirs = [3 + i * 2 for i in range(chapter_count)]
    actual_dirs = [i for i, slide in enumerate(prs.slides, start=1) if _is_directory_slide(slide, chapter_lines)]
    lines.append(f"expected_directory_pages: {expected_dirs}")
    lines.append(f"actual_directory_pages:   {actual_dirs}")
    if pattern_ids:
        lines.append(f"semantic_patterns: {pattern_ids}")

    risk = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = re.sub(r"\s+", " ", shape.text_frame.text).strip()
                if len(text) > 180:
                    risk += 1
    lines.append(f"overflow_risk_boxes: {risk}")
    lines.append("note: overflow_risk_boxes > 0 means manual review recommended.")

    report = pptx_path.with_name(pptx_path.stem + "_QA.txt")
    report.write_text("\n".join(lines), encoding="utf-8")
    return report
