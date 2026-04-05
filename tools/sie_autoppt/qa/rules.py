import re
import zipfile
from collections import Counter
from pathlib import Path
from xml.etree import ElementTree

from pptx import Presentation

from .models import QaChecks, QaMetrics, QaResult
from ..patterns import infer_pattern_details
from ..template_manifest import TemplateManifest, load_template_manifest


def _slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = re.sub(r"\s+", " ", shape.text_frame.text).strip()
            if text:
                texts.append(text)
    return " ".join(texts)


def _is_ending_slide(slide) -> bool:
    text = _slide_text(slide).lower()
    if not text and len(slide.shapes) <= 2:
        return True
    return any(keyword in text for keyword in ("thanks", "thank you", "q&a", "谢谢", "感谢"))


def _is_directory_slide(slide, chapter_lines: list[str], manifest: TemplateManifest | None = None) -> bool:
    if manifest is not None:
        title_boxes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and manifest.selectors.directory_title.matches(shape)
        ]
        required_count = min(len(chapter_lines) or 1, 5)
        if len(title_boxes) >= required_count:
            return True

    text = _slide_text(slide)
    if not text:
        return False
    lowered = text.lower()
    return "目录" in text or "content" in lowered


def _iter_text_run_sizes(shape) -> list[float]:
    sizes = []
    if not getattr(shape, "has_text_frame", False):
        return sizes
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip() and run.font.size is not None:
                sizes.append(round(run.font.size.pt, 1))
    return sizes


def _theme_title_font_ok(prs: Presentation, manifest: TemplateManifest) -> bool:
    if manifest.slide_roles.theme >= len(prs.slides):
        return False
    slide = prs.slides[manifest.slide_roles.theme]
    candidates = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and manifest.selectors.theme_title.matches(shape)
    ]
    if not candidates:
        return False
    title_shape = max(candidates, key=lambda shape: shape.width)
    sizes = _iter_text_run_sizes(title_shape)
    expected = float(manifest.fonts.theme_title_pt)
    return bool(sizes) and all(size == expected for size in sizes)


def _directory_title_font_ok(
    prs: Presentation,
    directory_slides: list[int],
    chapter_lines: list[str],
    manifest: TemplateManifest,
) -> bool:
    if not directory_slides or not chapter_lines:
        return False
    for slide_no in directory_slides:
        slide = prs.slides[slide_no - 1]
        title_boxes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and manifest.selectors.directory_title.matches(shape)
        ]
        title_boxes = sorted(title_boxes, key=lambda shape: (shape.top, shape.left))[: min(len(chapter_lines), 5)]
        if len(title_boxes) < min(len(chapter_lines), 5):
            return False
        for shape in title_boxes:
            sizes = _iter_text_run_sizes(shape)
            expected = float(manifest.fonts.directory_title_pt)
            if not sizes or any(size != expected for size in sizes):
                return False
    return True


def _slide_image_targets(pptx_path: Path, slide_no: int) -> set[str]:
    rel_path = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    targets: set[str] = set()
    with zipfile.ZipFile(pptx_path) as package:
        try:
            root = ElementTree.fromstring(package.read(rel_path))
        except KeyError:
            return targets
    for rel in root:
        rel_type = rel.attrib.get("Type", "")
        if rel_type.endswith("/image"):
            target = rel.attrib.get("Target", "")
            if target:
                targets.add(target)
    return targets


def _directory_assets_preserved(pptx_path: Path, directory_slides: list[int]) -> bool:
    if len(directory_slides) <= 1:
        return True
    source_assets = _slide_image_targets(pptx_path, directory_slides[0])
    if not source_assets:
        return True
    for slide_no in directory_slides[1:]:
        if not source_assets.issubset(_slide_image_targets(pptx_path, slide_no)):
            return False
    return True


def _overflow_risk_boxes(prs: Presentation) -> int:
    risk = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = re.sub(r"\s+", " ", shape.text_frame.text).strip()
                if len(text) > 180:
                    risk += 1
    return risk


def _long_text_shape_count(prs: Presentation) -> int:
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = re.sub(r"\s+", " ", shape.text_frame.text).strip()
                if len(text) > 120:
                    count += 1
    return count


def _repeated_title_count(render_trace) -> int:
    if render_trace is None:
        return 0
    titles = [trace.title.strip() for trace in render_trace.page_traces if trace.title.strip()]
    counts = Counter(titles)
    return sum(count - 1 for count in counts.values() if count > 1)


def _fallback_render_pages(render_trace) -> int:
    if render_trace is None:
        return 0
    return sum(1 for trace in render_trace.page_traces if trace.fallback_reason or trace.render_route.startswith("native_fallback:"))


def _reference_style_page_count(render_trace) -> int:
    if render_trace is None:
        return 0
    return sum(1 for trace in render_trace.page_traces if trace.reference_style_id)


def _low_confidence_pattern_pages(render_trace) -> int:
    if render_trace is None:
        return 0
    return sum(
        1
        for trace in render_trace.page_traces
        if infer_pattern_details(trace.title, []).low_confidence
    )


def _sparse_bullet_pages(render_trace) -> int:
    if render_trace is None:
        return 0
    return sum(1 for trace in render_trace.page_traces if not trace.title.strip())


def _template_pool_mode(manifest: TemplateManifest) -> str:
    return "PASS" if manifest.slide_pools else "WARN"


def _reference_style_coverage(render_trace) -> str:
    if render_trace is None:
        return "PASS"
    for trace in render_trace.page_traces:
        if trace.reference_style_id and trace.render_route.startswith("native_fallback:"):
            return "WARN"
    return "PASS"


def _preflight_status(render_trace) -> str:
    if render_trace is None:
        return "PASS"
    return "WARN" if render_trace.preflight_notes else "PASS"


def _content_density_status(long_text_shape_count: int, overflow_risk_boxes: int) -> str:
    return "WARN" if long_text_shape_count > 0 or overflow_risk_boxes > 0 else "PASS"


def _title_uniqueness_status(repeated_title_count: int) -> str:
    return "WARN" if repeated_title_count > 0 else "PASS"


def build_qa_result(
    pptx_path: Path,
    chapter_count: int,
    pattern_ids=None,
    chapter_lines=None,
    template_path: Path | None = None,
    render_trace=None,
) -> QaResult:
    prs = Presentation(str(pptx_path))
    chapter_lines = chapter_lines or []
    manifest = load_template_manifest(template_path=template_path)

    has_ending_last = _is_ending_slide(prs.slides[-1])
    if manifest.slide_pools:
        expected_dirs = [index + 1 for index in manifest.slide_pools.directory[:chapter_count]]
    else:
        expected_dirs = [manifest.slide_roles.directory + 1 + i * 2 for i in range(chapter_count)]
    if expected_dirs and all(
        slide_no <= len(prs.slides)
        and _is_directory_slide(prs.slides[slide_no - 1], chapter_lines, manifest=manifest)
        for slide_no in expected_dirs
    ):
        actual_dirs = expected_dirs
    else:
        actual_dirs = [
            i
            for i, slide in enumerate(prs.slides, start=1)
            if _is_directory_slide(slide, chapter_lines, manifest=manifest)
        ]
    overflow_risk = _overflow_risk_boxes(prs)
    long_text_shape_count = _long_text_shape_count(prs)
    repeated_title_count = _repeated_title_count(render_trace)
    fallback_render_pages = _fallback_render_pages(render_trace)
    low_confidence_pattern_pages = _low_confidence_pattern_pages(render_trace)
    reference_style_pages = _reference_style_page_count(render_trace)
    sparse_bullet_pages = _sparse_bullet_pages(render_trace)
    preflight_note_count = len(render_trace.preflight_notes) if render_trace is not None else 0

    return QaResult(
        file=str(pptx_path),
        slides=len(prs.slides),
        expected_directory_pages=expected_dirs,
        actual_directory_pages=actual_dirs,
        template_name=manifest.template_name,
        template_manifest_path=manifest.manifest_path,
        template_manifest_version=manifest.version,
        expected_theme_title_font_pt=manifest.fonts.theme_title_pt,
        expected_directory_title_font_pt=manifest.fonts.directory_title_pt,
        semantic_patterns=list(pattern_ids or []),
        chapter_lines=list(chapter_lines),
        render_trace=render_trace,
        checks=QaChecks(
            ending_last="PASS" if has_ending_last else "WARN",
            theme_title_font="PASS" if _theme_title_font_ok(prs, manifest) else "WARN",
            directory_title_font="PASS" if _directory_title_font_ok(prs, actual_dirs, chapter_lines, manifest) else "WARN",
            directory_assets_preserved="PASS" if _directory_assets_preserved(pptx_path, actual_dirs) else "WARN",
            template_pool_mode=_template_pool_mode(manifest),
            reference_style_coverage=_reference_style_coverage(render_trace),
            preflight=_preflight_status(render_trace),
            content_density=_content_density_status(long_text_shape_count, overflow_risk),
            title_uniqueness=_title_uniqueness_status(repeated_title_count),
        ),
        metrics=QaMetrics(
            overflow_risk_boxes=overflow_risk,
            long_text_shape_count=long_text_shape_count,
            sparse_bullet_pages=sparse_bullet_pages,
            repeated_title_count=repeated_title_count,
            fallback_render_pages=fallback_render_pages,
            low_confidence_pattern_pages=low_confidence_pattern_pages,
            reference_style_pages=reference_style_pages,
            preflight_note_count=preflight_note_count,
        ),
        notes=[
            "overflow_risk_boxes > 0 means manual review recommended.",
            "long_text_shape_count > 0 indicates dense text that may need layout polish.",
            "fallback_render_pages > 0 means some pages fell back from requested reference styles.",
            "low_confidence_pattern_pages > 0 means some page titles looked semantically ambiguous.",
            "preflight_note_count > 0 means runtime detected template or content warnings before final QA.",
            "render_trace.page_traces records each page's render route and fallback reason when available.",
        ],
    )
