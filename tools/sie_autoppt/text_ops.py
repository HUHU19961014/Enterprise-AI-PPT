from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .config import FONT_NAME


def _apply_run_style(run, size_pt: int | None, bold: bool | None, color):
    run.font.name = FONT_NAME
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = RGBColor(*color)


def write_text(
    shape,
    text: str,
    *,
    size_pt: int | None = None,
    bold: bool | None = None,
    color=None,
    align=PP_ALIGN.LEFT,
    preserve_runs: bool = False,
):
    if not getattr(shape, "has_text_frame", False):
        return

    text_frame = shape.text_frame
    text_frame.word_wrap = True

    if preserve_runs and text_frame.paragraphs:
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = align
        if paragraph.runs:
            run = paragraph.runs[0]
            run.text = text
            _apply_run_style(run, size_pt, bold, color)
            for extra_run in paragraph.runs[1:]:
                extra_run.text = ""
        else:
            run = paragraph.add_run()
            run.text = text
            _apply_run_style(run, size_pt, bold, color)
        for extra_paragraph in text_frame.paragraphs[1:]:
            for extra_run in extra_paragraph.runs:
                extra_run.text = ""
        return

    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    _apply_run_style(run, size_pt, bold, color)


def add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    *,
    size_pt: int | None = None,
    bold: bool | None = None,
    color=None,
    align=PP_ALIGN.LEFT,
    margin_left: int | None = None,
    margin_right: int | None = None,
    margin_top: int | None = None,
    margin_bottom: int | None = None,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    if margin_left is not None:
        text_frame.margin_left = margin_left
    if margin_right is not None:
        text_frame.margin_right = margin_right
    if margin_top is not None:
        text_frame.margin_top = margin_top
    if margin_bottom is not None:
        text_frame.margin_bottom = margin_bottom
    write_text(
        textbox,
        text,
        size_pt=size_pt,
        bold=bold,
        color=color,
        align=align,
        preserve_runs=False,
    )
    return textbox
