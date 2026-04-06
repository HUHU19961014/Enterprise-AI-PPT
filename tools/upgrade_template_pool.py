import argparse
import sys
from pathlib import Path

from pptx import Presentation


def _bootstrap():
    tools_dir = Path(__file__).resolve().parent
    if str(tools_dir) not in sys.path:
        sys.path.insert(0, str(tools_dir))


_bootstrap()

from sie_autoppt.slide_ops import clone_slide_after, copy_slide_xml_assets, ensure_last_slide, slide_assets_preserved
from sie_autoppt.template_manifest import load_template_manifest
from sie_autoppt.generator import validate_slide_pool_configuration

VALIDATION_MODE = "python-openxml"


def upgrade_template_pool(template_path: Path) -> bool:
    manifest = load_template_manifest(template_path=template_path)
    if not manifest.slide_pools:
        raise ValueError("Template manifest does not define slide_pools.")

    changed = _upgrade_template_pool_with_python(template_path, manifest)
    validate_template_pool(template_path, manifest)
    return changed


def validate_template_pool(template_path: Path, manifest=None) -> dict[str, object]:
    manifest = manifest or load_template_manifest(template_path=template_path)
    if not manifest.slide_pools:
        raise ValueError("Template manifest does not define slide_pools.")

    prs = Presentation(str(template_path))
    required_pairs = len(manifest.slide_pools.body)
    validate_slide_pool_configuration(manifest, body_page_count=required_pairs, slide_count=len(prs.slides))
    if len(prs.slides) != manifest.slide_pools.ending + 1:
        raise RuntimeError("Template slide count does not match the configured ending slide index after pool upgrade.")

    source_idx = manifest.slide_pools.directory[0] + 1
    target_indices = [index + 1 for index in manifest.slide_pools.directory[1:required_pairs]]
    if target_indices and not slide_assets_preserved(template_path, source_idx=source_idx, target_indices=target_indices):
        raise RuntimeError("Template pool validation failed: directory slide assets were not preserved.")

    return {
        "slides": len(prs.slides),
        "required_pairs": required_pairs,
        "ending_slide_no": manifest.slide_pools.ending + 1,
        "directory_asset_targets": len(target_indices),
        "validation_mode": VALIDATION_MODE,
    }


def format_validation_summary(summary: dict[str, object]) -> str:
    return (
        "Validation passed: "
        f"mode={summary['validation_mode']}, "
        f"slides={summary['slides']}, "
        f"required_pairs={summary['required_pairs']}, "
        f"ending_slide_no={summary['ending_slide_no']}, "
        f"directory_asset_targets={summary['directory_asset_targets']}"
    )


def _upgrade_template_pool_with_python(template_path: Path, manifest) -> bool:
    prs = Presentation(str(template_path))
    required_pairs = len(manifest.slide_pools.body)
    if manifest.supports_preallocated_pool(required_pairs, len(prs.slides)):
        source_idx = manifest.slide_pools.directory[0] + 1
        target_indices = [index + 1 for index in manifest.slide_pools.directory[1:required_pairs]]
        if slide_assets_preserved(template_path, source_idx=source_idx, target_indices=target_indices):
            return False

    directory_idx = manifest.slide_roles.directory
    body_template_idx = manifest.slide_roles.body_template
    thanks_slide_id = int(prs.slides._sldIdLst[len(prs.slides) - 1].id)
    insert_after = body_template_idx

    existing_pairs = 1
    while existing_pairs < required_pairs:
        clone_slide_after(prs, directory_idx, insert_after, keep_rel_ids=True)
        insert_after += 1
        clone_slide_after(prs, body_template_idx, insert_after, keep_rel_ids=False)
        insert_after += 1
        existing_pairs += 1

    ensure_last_slide(prs, thanks_slide_id)
    prs.save(str(template_path))

    source_idx = directory_idx + 1
    target_indices = [index + 1 for index in manifest.slide_pools.directory[1:required_pairs]]
    if target_indices:
        copy_slide_xml_assets(template_path, source_idx=source_idx, target_indices=target_indices)
        if not slide_assets_preserved(template_path, source_idx=source_idx, target_indices=target_indices):
            raise RuntimeError("Directory slide assets were not preserved while upgrading the template pool.")
    return True


def main():
    parser = argparse.ArgumentParser(
        description="Expand a template into a preallocated slide pool using the Python/OpenXML path."
    )
    parser.add_argument("template", nargs="?", default="assets/templates/sie_template.pptx", help="Template PPTX path.")
    parser.add_argument(
        "--validate-only",
        action="store_true",
        help="Only validate the existing template pool without modifying the PPTX.",
    )
    args = parser.parse_args()

    template_path = Path(args.template).resolve()
    if args.validate_only:
        print(f"Template validation requested: {template_path}")
    else:
        changed = upgrade_template_pool(template_path)
        if changed:
            print(f"Template upgraded: {template_path}")
        else:
            print(f"Template already pooled: {template_path}")

    summary = validate_template_pool(template_path)
    print(format_validation_summary(summary))


if __name__ == "__main__":
    main()
