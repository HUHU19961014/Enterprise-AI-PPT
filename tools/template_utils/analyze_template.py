from pathlib import Path

from pptx import Presentation

project_root = Path(__file__).resolve().parents[2]
template = project_root / "assets" / "templates" / "sie_template.pptx"
prs = Presentation(str(template))

print(f"slides={len(prs.slides)}")
print(f"size={prs.slide_width}x{prs.slide_height}")

for i, slide in enumerate(prs.slides, start=1):
    print(f"\n[slide {i}] shapes={len(slide.shapes)}")
    for j, shape in enumerate(slide.shapes, start=1):
        has_text = getattr(shape, "has_text_frame", False)
        text = ""
        if has_text:
            text = shape.text_frame.text.replace("\n", " ").strip()
        print(
            f"  - {j:02d} type={shape.shape_type} "
            f"name={shape.name!r} left={shape.left} top={shape.top} "
            f"w={shape.width} h={shape.height} text={text[:40]!r}"
        )

