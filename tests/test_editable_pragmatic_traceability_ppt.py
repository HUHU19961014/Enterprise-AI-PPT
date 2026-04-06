import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.build_editable_pragmatic_traceability_ppt import (
    DEFAULT_HTML,
    build_ppt,
    normalize_text,
    parse_html,
)


class EditablePragmaticTraceabilityPptTests(unittest.TestCase):
    def test_normalize_text_repairs_known_encoding_artifacts(self):
        self.assertEqual(normalize_text("前缀锛?内容"), "前缀：内容")
        self.assertEqual(normalize_text("说明Ŗē补充"), "说明：补充")

    def test_parse_html_extracts_sections_and_points(self):
        spec = parse_html(DEFAULT_HTML)

        self.assertEqual(len(spec.sections), 2)
        self.assertEqual(spec.sections[0].title, "对标头部：业务规则已被IT系统固化")
        self.assertEqual(len(spec.value_items), 3)
        self.assertEqual(spec.points[0].name, "英发睿能")

    def test_build_ppt_outputs_editable_shapes_instead_of_picture_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "editable.pptx"
            build_ppt(DEFAULT_HTML, output)

            prs = Presentation(output)
            slide = prs.slides[0]

            self.assertGreater(len(slide.shapes), 20)
            self.assertFalse(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes))


if __name__ == "__main__":
    unittest.main()
