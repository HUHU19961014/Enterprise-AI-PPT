import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from tools.sie_autoppt.v2.layout_ids import SUPPORTED_LAYOUTS
from tools.sie_autoppt.v2.ppt_engine import generate_ppt


def _slide_for_layout(layout: str) -> dict[str, object]:
    if layout == "section_break":
        return {"slide_id": "s1", "layout": layout, "title": "Section", "subtitle": "Bridge"}
    if layout == "title_only":
        return {"slide_id": "s1", "layout": layout, "title": "Executive Summary"}
    if layout == "title_content":
        return {"slide_id": "s1", "layout": layout, "title": "Context", "content": ["Point A", "Point B"]}
    if layout == "two_columns":
        return {
            "slide_id": "s1",
            "layout": layout,
            "title": "Comparison",
            "left": {"heading": "Current", "items": ["Issue 1", "Issue 2"]},
            "right": {"heading": "Target", "items": ["Action 1", "Action 2"]},
        }
    if layout == "title_image":
        return {
            "slide_id": "s1",
            "layout": layout,
            "title": "Architecture",
            "content": ["Layer 1", "Layer 2"],
            "image": {"mode": "placeholder", "caption": "Architecture"},
        }
    if layout == "timeline":
        return {
            "slide_id": "s1",
            "layout": layout,
            "title": "Roadmap",
            "stages": [{"title": "Q1", "detail": "Align"}, {"title": "Q2", "detail": "Pilot"}],
        }
    if layout == "stats_dashboard":
        return {
            "slide_id": "s1",
            "layout": layout,
            "title": "KPI",
            "metrics": [{"label": "OTD", "value": "95%"}, {"label": "Yield", "value": "98%"}],
        }
    if layout == "matrix_grid":
        return {
            "slide_id": "s1",
            "layout": layout,
            "title": "Risk Matrix",
            "cells": [{"title": "Low-Low", "body": "Monitor"}, {"title": "High-High", "body": "Escalate"}],
        }
    if layout == "cards_grid":
        return {
            "slide_id": "s1",
            "layout": layout,
            "title": "Capabilities",
            "cards": [{"title": "Plan", "body": "Align"}, {"title": "Operate", "body": "Close loop"}],
        }
    raise ValueError(f"unsupported test layout: {layout}")


class RenderIntegrationTests(unittest.TestCase):
    def test_all_supported_layouts_render(self):
        for layout in SUPPORTED_LAYOUTS:
            with self.subTest(layout=layout):
                payload = {
                    "meta": {
                        "title": f"Render {layout}",
                        "theme": "sie_consulting_fixed",
                        "language": "en-US",
                        "author": "AI",
                        "version": "2.0",
                    },
                    "slides": [_slide_for_layout(layout)],
                }
                with tempfile.TemporaryDirectory() as temp_dir:
                    ppt_path = Path(temp_dir) / f"{layout}.pptx"
                    log_path = Path(temp_dir) / f"{layout}.log.txt"
                    result = generate_ppt(payload, output_path=ppt_path, log_path=log_path)
                    self.assertTrue(result.output_path.exists())
                    self.assertEqual(result.slide_count, 1)
                    self.assertEqual(result.final_deck.slides[0].layout, layout)
                    prs = Presentation(str(result.output_path))
                    self.assertGreater(len(prs.slides), 0)


if __name__ == "__main__":
    unittest.main()
