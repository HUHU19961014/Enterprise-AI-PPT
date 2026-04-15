from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.scenario_generators.build_equipment_mfg_digital_solution_deck import build_deck


class EquipmentManufacturingDigitalSolutionDeckTests(unittest.TestCase):
    def test_build_deck_outputs_full_editable_deck(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "equipment_mfg_solution_deck.pptx"
            built = build_deck(output)

            prs = Presentation(str(built))
            self.assertEqual(len(prs.slides), 9)
            self.assertFalse(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in prs.slides[3].shapes))
            self.assertIn("行业判断与转型命题", "\n".join(
                shape.text.strip()
                for shape in prs.slides[1].shapes
                if hasattr(shape, "text") and shape.text.strip()
            ))
            self.assertIn("建议汇报结构", "\n".join(
                shape.text.strip()
                for shape in prs.slides[2].shapes
                if hasattr(shape, "text") and shape.text.strip()
            ))
            self.assertIn("行业核心痛点", "\n".join(
                shape.text.strip()
                for shape in prs.slides[3].shapes
                if hasattr(shape, "text") and shape.text.strip()
            ))

            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        all_text.append(shape.text.strip())

            with zipfile.ZipFile(built) as zf:
                self.assertEqual(len(zf.namelist()), len(set(zf.namelist())))

        joined = "\n".join(all_text)
        self.assertIn("装备制造行业数字化解决方案", joined)
        self.assertIn("转型路径一：业务主链能力", joined)
        self.assertIn("转型路径二：运营执行能力", joined)
        self.assertIn("实施方向与分阶段路线", joined)


if __name__ == "__main__":
    unittest.main()
