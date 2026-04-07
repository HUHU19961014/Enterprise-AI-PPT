import copy
import unittest
from dataclasses import replace

from pptx import Presentation

from tools.sie_autoppt.body_renderers import _layout_for_page, fill_body_slide
from tools.sie_autoppt.config import DEFAULT_TEMPLATE
from tools.sie_autoppt.models import BodyPageSpec
from tools.sie_autoppt.template_manifest import load_template_manifest


class BodyRendererTests(unittest.TestCase):
    def test_layout_variant_overlays_base_layout_for_theme_variants(self):
        manifest = load_template_manifest(template_path=DEFAULT_TEMPLATE.parent / "business_gold" / "template.pptx")
        page = BodyPageSpec(
            page_key="p1",
            title="Test",
            subtitle="Sub",
            bullets=["A", "B", "C"],
            pattern_id="general_business",
            layout_variant="general_business_3",
        )

        spec = _layout_for_page(manifest, page, "general_business")

        self.assertEqual(int(spec["columns"]), 3)
        self.assertEqual(int(spec["max_items"]), 3)
        self.assertEqual(tuple(int(value) for value in spec["fill_rgb"]), (251, 248, 240))

    def test_fill_body_slide_reports_context_when_manifest_field_is_missing(self):
        manifest = load_template_manifest()
        broken_layouts = copy.deepcopy(manifest.render_layouts)
        del broken_layouts["general_business"]["origin_left"]
        broken_manifest = replace(manifest, render_layouts=broken_layouts)

        presentation = Presentation(str(DEFAULT_TEMPLATE))
        slide = presentation.slides[manifest.slide_roles.body_template]
        page = BodyPageSpec(
            page_key="p1",
            title="测试页",
            subtitle="副标题",
            bullets=["要点一", "要点二"],
            pattern_id="general_business",
        )

        with self.assertRaises(KeyError) as ctx:
            fill_body_slide(slide, page, broken_manifest)

        message = str(ctx.exception)
        self.assertIn("origin_left", message)
        self.assertIn("render_layouts.general_business", message)

    def test_fill_body_slide_supports_dashboard_and_roadmap_patterns(self):
        manifest = load_template_manifest()

        dashboard_presentation = Presentation(str(DEFAULT_TEMPLATE))
        dashboard_slide = dashboard_presentation.slides[manifest.slide_roles.body_template]
        dashboard_page = BodyPageSpec(
            page_key="p_dashboard",
            title="经营指标",
            subtitle="核心 KPI 摘要",
            bullets=["收入增长: 18%", "利润提升: 6%"],
            pattern_id="kpi_dashboard",
            payload={
                "headline": "经营表现摘要",
                "footer": "统一指标口径追踪经营成效",
                "metrics": [
                    {"label": "收入增长", "value": "18%", "detail": "同比提升"},
                    {"label": "利润提升", "value": "6%", "detail": "结构优化驱动"},
                    {"label": "准时交付", "value": "96%", "detail": "关键节点稳定"},
                    {"label": "满意度", "value": "92%", "detail": "重点客户反馈"},
                ],
                "insights": ["增长韧性较强", "利润改善明确", "交付稳定性提升"],
            },
        )
        fill_body_slide(dashboard_slide, dashboard_page, manifest)
        self.assertGreater(len(dashboard_slide.shapes), 4)

        roadmap_presentation = Presentation(str(DEFAULT_TEMPLATE))
        roadmap_slide = roadmap_presentation.slides[manifest.slide_roles.body_template]
        roadmap_page = BodyPageSpec(
            page_key="p_roadmap",
            title="年度路线图",
            subtitle="按季度推进",
            bullets=["Q1: 完成现状诊断", "Q2: 搭建能力底座"],
            pattern_id="roadmap_timeline",
            payload={
                "headline": "里程碑安排",
                "footer": "年度推进节奏与关键任务保持一致",
                "stages": [
                    {"period": "Q1", "title": "现状诊断", "detail": "完成现状调研与问题识别"},
                    {"period": "Q2", "title": "平台建设", "detail": "完成关键能力与底座建设"},
                    {"period": "Q3", "title": "试点上线", "detail": "验证闭环与价值"},
                    {"period": "Q4", "title": "复制推广", "detail": "规模化推广应用"},
                ],
            },
        )
        fill_body_slide(roadmap_slide, roadmap_page, manifest)
        self.assertGreater(len(roadmap_slide.shapes), 4)

        risk_presentation = Presentation(str(DEFAULT_TEMPLATE))
        risk_slide = risk_presentation.slides[manifest.slide_roles.body_template]
        risk_page = BodyPageSpec(
            page_key="p_risk",
            title="风险矩阵",
            subtitle="优先识别高影响事项",
            bullets=["汇率风险", "政策风险"],
            pattern_id="risk_matrix",
            payload={
                "headline": "关键风险分布",
                "footer": "优先处理高概率高影响风险",
                "items": [
                    {"title": "汇率风险", "detail": "影响偿债能力", "quadrant": "high_high"},
                    {"title": "政策风险", "detail": "影响谈判预期", "quadrant": "low_high"},
                    {"title": "执行风险", "detail": "影响项目收益", "quadrant": "high_low"},
                    {"title": "声誉风险", "detail": "影响合作空间", "quadrant": "low_low"},
                ],
            },
        )
        fill_body_slide(risk_slide, risk_page, manifest)
        self.assertGreater(len(risk_slide.shapes), 4)

        claim_presentation = Presentation(str(DEFAULT_TEMPLATE))
        claim_slide = claim_presentation.slides[manifest.slide_roles.body_template]
        claim_page = BodyPageSpec(
            page_key="p_claim",
            title="索赔拆解",
            subtitle="金额构成",
            bullets=["电费欠款", "违约利息"],
            pattern_id="claim_breakdown",
            payload={
                "headline": "索赔主项构成",
                "footer": "先识别金额最大主项",
                "summary": "电费欠款和违约利息构成主要金额来源。",
                "claims": [
                    {"label": "电费欠款", "value": "$2.1B", "detail": "历史欠付主项"},
                    {"label": "违约利息", "value": "$1.3B", "detail": "违约累计形成"},
                    {"label": "汇率损失", "value": "$0.8B", "detail": "货币贬值影响"},
                    {"label": "其他费用", "value": "$0.3B", "detail": "杂项成本"},
                ],
            },
        )
        fill_body_slide(claim_slide, claim_page, manifest)
        self.assertGreater(len(claim_slide.shapes), 4)
