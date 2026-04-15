import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.scenario_generators.build_us_supply_chain_compliance_slide import build_slide, self_check_layout


class UsSupplyChainComplianceSlideTests(unittest.TestCase):
    def test_build_slide_outputs_single_editable_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "us_supply_chain_compliance.pptx"
            built, review_path, score_path, score = build_slide(output)

            prs = Presentation(built)
            self.assertEqual(len(prs.slides), 1)
            self.assertFalse(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in prs.slides[0].shapes))
            self.assertTrue(review_path.exists())
            self.assertTrue(score_path.exists())
            self.assertGreaterEqual(score.total, 90)

    def test_self_check_passes_for_generated_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "us_supply_chain_compliance.pptx"
            built, _, _, _ = build_slide(output)

            prs = Presentation(built)
            self.assertEqual(self_check_layout(prs), [])

    def test_score_json_contains_expected_fields(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "us_supply_chain_compliance.pptx"
            _, _, score_path, _ = build_slide(output)

            payload = json.loads(score_path.read_text(encoding="utf-8"))
            self.assertIn("template_fidelity", payload)
            self.assertIn("title_fidelity_to_sie", payload)
            self.assertIn("layout_originality", payload)
            self.assertIn("heuristic_review", payload)
            self.assertIn("content_coverage", payload)


if __name__ == "__main__":
    unittest.main()
