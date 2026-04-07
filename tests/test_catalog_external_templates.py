import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from tools.template_utils.catalog_external_templates import build_catalog, discover_pptx_files, main


class CatalogExternalTemplatesTests(unittest.TestCase):
    def _make_pptx(self, path: Path, title: str, subtitle: str) -> None:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        prs.save(str(path))

    def test_discover_pptx_files_skips_temp_files(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            self._make_pptx(root / "a.pptx", "A", "alpha")
            self._make_pptx(root / "~$ignore.pptx", "B", "beta")

            files = discover_pptx_files(root)

            self.assertEqual([item.name for item in files], ["a.pptx"])

    def test_build_catalog_writes_ranked_outputs(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir) / "templates"
            output_dir = Path(temp_dir) / "catalog"
            root.mkdir()
            self._make_pptx(root / "agenda_like.pptx", "Agenda", "1. Intro\n2. KPI")
            self._make_pptx(root / "cover_only.pptx", "Cover", "Simple")

            catalog = build_catalog(root, output_dir)

            self.assertEqual(catalog["template_count"], 2)
            self.assertTrue((output_dir / "catalog.json").exists())
            self.assertTrue((output_dir / "catalog.md").exists())
            self.assertEqual(len(catalog["entries"]), 2)
            self.assertTrue(catalog["entries"][0]["score"] >= catalog["entries"][1]["score"])
            self.assertNotEqual(catalog["entries"][0]["import_dir"], catalog["entries"][1]["import_dir"])
            self.assertTrue((Path(catalog["entries"][0]["import_dir"]) / "manifest.json").exists())

    def test_main_returns_zero_for_valid_directory(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir) / "templates"
            output_dir = Path(temp_dir) / "catalog"
            root.mkdir()
            self._make_pptx(root / "one.pptx", "Agenda", "Outline")

            exit_code = main([str(root), "--output", str(output_dir)])

            self.assertEqual(exit_code, 0)
