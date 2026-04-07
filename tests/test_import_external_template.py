import base64
import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from tools.template_utils.import_external_pptx_template import build_import_manifest, main


_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO7Z0xQAAAAASUVORK5CYII="
)


class ImportExternalTemplateTests(unittest.TestCase):
    def _build_sample_pptx(self, target: Path) -> None:
        image_path = target.with_suffix(".png")
        image_path.write_bytes(_PNG_1X1)

        prs = Presentation()
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        cover.shapes.title.text = "Strategy Template"
        cover.placeholders[1].text = "Reusable external source"

        agenda = prs.slides.add_slide(prs.slide_layouts[1])
        agenda.shapes.title.text = "Agenda"
        agenda.placeholders[1].text = "1. Market context\n2. KPI dashboard\n3. Risk response"
        agenda.shapes.add_picture(str(image_path), Inches(8.3), Inches(0.8), width=Inches(1.0), height=Inches(1.0))

        content = prs.slides.add_slide(prs.slide_layouts[1])
        content.shapes.title.text = "Operations Metrics"
        text_frame = content.placeholders[1].text_frame
        text_frame.text = "Revenue +18%"
        for bullet in ("Gross margin 34%", "NPS 62", "Launch 3 pilots", "Pipeline 12 accounts"):
            paragraph = text_frame.add_paragraph()
            paragraph.text = bullet
        content.shapes.add_picture(str(image_path), Inches(8.0), Inches(3.8), width=Inches(1.2), height=Inches(1.2))

        if prs.slides:
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        prs.save(str(target))

    def test_build_import_manifest_exports_assets_and_hints(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_path = Path(temp_dir) / "external_source.pptx"
            output_dir = Path(temp_dir) / "imported"
            self._build_sample_pptx(pptx_path)

            manifest = build_import_manifest(pptx_path, output_dir)

            self.assertEqual(manifest["source"]["name"], "external_source.pptx")
            self.assertEqual(manifest["assets"]["exportDir"], "assets")
            self.assertEqual(len(manifest["slides"]), 2)
            self.assertIn("toc_candidate", manifest["pageTypeCandidates"])
            self.assertTrue(manifest["assets"]["allAssets"])
            self.assertTrue(manifest["assets"]["commonAssets"])
            self.assertTrue(manifest["fusionHints"]["themeReady"])
            self.assertTrue(manifest["fusionHints"]["referenceStyleCandidates"])
            self.assertTrue((output_dir / "analysis.md").exists())
            self.assertTrue((output_dir / "assets" / manifest["assets"]["allAssets"][0]).exists())

    def test_main_writes_manifest_json(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            pptx_path = Path(temp_dir) / "external_source.pptx"
            output_dir = Path(temp_dir) / "imported"
            self._build_sample_pptx(pptx_path)

            exit_code = main([str(pptx_path), "--output", str(output_dir)])

            self.assertEqual(exit_code, 0)
            manifest_path = output_dir / "manifest.json"
            self.assertTrue(manifest_path.exists())
            data = json.loads(manifest_path.read_text(encoding="utf-8"))
            self.assertIn("fusionHints", data)
            self.assertIn("recommendedActions", data["fusionHints"])
