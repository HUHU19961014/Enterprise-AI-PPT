import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from tools.sie_autoppt.v2.ppt_engine import generate_ppt
from tools.sie_autoppt.v2.schema import validate_deck_payload


class V2RenderTests(unittest.TestCase):
    def test_generate_ppt_renders_staged_content_as_timeline(self):
        payload = {
            "meta": {"title": "Timeline Test", "theme": "business_red", "language": "zh-CN", "author": "AI", "version": "2.0"},
            "slides": [
                {
                    "slide_id": "s1",
                    "layout": "title_content",
                    "title": "实施路径",
                    "content": [
                        "第一阶段统一标准与目录，建立治理底座",
                        "第二阶段接入质量规则与问题闭环流程",
                        "第三阶段推动重点域治理运营与效果评估",
                        "第四阶段扩展到跨业务条线的协同治理",
                    ],
                }
            ],
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "timeline.pptx"
            log_path = Path(temp_dir) / "timeline.log.txt"
            result = generate_ppt(payload, output_path=ppt_path, log_path=log_path)

            self.assertTrue(result.output_path.exists())
            self.assertIn("timeline visualization", log_path.read_text(encoding="utf-8"))

    def test_generate_ppt_renders_placeholder_as_architecture_diagram(self):
        payload = {
            "meta": {"title": "Architecture Test", "theme": "business_red", "language": "zh-CN", "author": "AI", "version": "2.0"},
            "slides": [
                {
                    "slide_id": "s1",
                    "layout": "title_image",
                    "title": "方案架构",
                    "content": [
                        "统一元数据、标准、质量、目录与权限管理",
                        "建设治理规则中心，形成平台化运营入口",
                        "通过治理流程连接业务系统与分析应用",
                    ],
                    "image": {"mode": "placeholder", "caption": "数据治理平台架构图"},
                }
            ],
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "architecture.pptx"
            log_path = Path(temp_dir) / "architecture.log.txt"
            result = generate_ppt(payload, output_path=ppt_path, log_path=log_path)

            self.assertTrue(result.output_path.exists())
            self.assertIn("architecture diagram", log_path.read_text(encoding="utf-8"))

    def test_generate_ppt_renders_balanced_two_columns_as_comparison_table(self):
        payload = {
            "meta": {"title": "Comparison Test", "theme": "business_red", "language": "zh-CN", "author": "AI", "version": "2.0"},
            "slides": [
                {
                    "slide_id": "s1",
                    "layout": "two_columns",
                    "title": "常见误区",
                    "left": {
                        "heading": "典型误区",
                        "items": ["只盯任务，不盯目标是否一致", "把问题记录下来，却没有闭环动作", "等风险变成事故后才上报"],
                    },
                    "right": {
                        "heading": "正确做法",
                        "items": ["用里程碑和结果来校验推进质量", "把问题清单与责任和时限绑定", "把风险前移到周会和关键节点评审"],
                    },
                }
            ],
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "comparison.pptx"
            log_path = Path(temp_dir) / "comparison.log.txt"
            result = generate_ppt(payload, output_path=ppt_path, log_path=log_path)

            self.assertTrue(result.output_path.exists())
            self.assertIn("comparison table", log_path.read_text(encoding="utf-8"))

    def test_generate_ppt_renders_sample_deck(self):
        sample_path = Path("samples/sample_deck_v2.json")
        payload = json.loads(sample_path.read_text(encoding="utf-8"))
        validated = validate_deck_payload(payload)

        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "sample_v2.pptx"
            log_path = Path(temp_dir) / "sample_v2.log.txt"
            result = generate_ppt(validated, output_path=ppt_path, log_path=log_path)

            self.assertTrue(result.output_path.exists())
            self.assertTrue(log_path.exists())
            self.assertIsNotNone(result.warnings_path)
            self.assertTrue(result.warnings_path.exists())
            self.assertIsNotNone(result.rewrite_log_path)
            self.assertTrue(result.rewrite_log_path.exists())

            prs = Presentation(str(result.output_path))
            self.assertEqual(len(prs.slides), len(validated.deck.slides))
            all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
            self.assertIn("HTML 不是 PPT 的天然中间语言", all_text)
            self.assertIn("重构前后对比", all_text)

            warnings_payload = json.loads(result.warnings_path.read_text(encoding="utf-8"))
            self.assertIn("passed", warnings_payload)
            self.assertIn("review_required", warnings_payload)
            self.assertIn("warnings", warnings_payload)
            self.assertIn("high", warnings_payload)
            self.assertIn("errors", warnings_payload)
            self.assertIn("summary", warnings_payload)

            rewrite_payload = json.loads(result.rewrite_log_path.read_text(encoding="utf-8"))
            self.assertIn("attempted", rewrite_payload)
            self.assertIn("applied", rewrite_payload)
            self.assertIn("actions", rewrite_payload)

    def test_generate_ppt_applies_single_rewrite_pass(self):
        payload = {
            "meta": {"title": "Test", "theme": "business_red", "language": "zh-CN", "author": "AI", "version": "2.0"},
            "slides": [
                {
                    "slide_id": "s1",
                    "layout": "title_content",
                    "title": "这是一个明显过长并且需要压缩表达的业务分析标题",
                    "content": [
                        "第一条内容明显过长，需要压缩到更适合页面承载的长度，并保留核心信息。",
                        "第二条内容也非常长，需要继续压缩表达，避免页面密度过高。",
                        "第三条需要保留。",
                        "第四条需要保留。",
                        "第五条需要保留。",
                        "第六条需要保留。",
                        "第七条用于测试自动合并。",
                    ],
                }
            ],
        }

        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "rewritten.pptx"
            log_path = Path(temp_dir) / "rewritten.log.txt"
            deck_path = Path(temp_dir) / "generated_deck.json"
            result = generate_ppt(payload, output_path=ppt_path, log_path=log_path, deck_output_path=deck_path)

            self.assertTrue(result.rewrite_applied)
            self.assertTrue(deck_path.exists())
            self.assertTrue(result.output_path.exists())

            rewritten_deck = json.loads(deck_path.read_text(encoding="utf-8"))
            slide = rewritten_deck["slides"][0]
            self.assertLessEqual(len(slide["content"]), 6)

            rewrite_payload = json.loads(result.rewrite_log_path.read_text(encoding="utf-8"))
            self.assertTrue(rewrite_payload["attempted"])
            self.assertTrue(rewrite_payload["applied"])
            self.assertGreater(rewrite_payload["action_count"], 0)
