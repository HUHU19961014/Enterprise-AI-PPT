import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools.scenario_generators.build_internal_traceability_uat_logic_slide import build_slide, self_check_layout


class InternalTraceabilityUatLogicSlideTests(unittest.TestCase):
    def test_build_slide_outputs_editable_single_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "internal_traceability_uat_logic.pptx"
            build_slide(output, export_review=False)

            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            self.assertFalse(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in prs.slides[0].shapes))

    def test_self_check_passes_for_generated_slide(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "internal_traceability_uat_logic.pptx"
            build_slide(output, export_review=False)

            prs = Presentation(output)
            self.assertEqual(self_check_layout(prs), [])

    def test_build_slide_supports_alternate_layout_preset(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "internal_traceability_uat_logic_dense.pptx"
            build_slide(output, preset_id="decision_oriented", export_review=False)

            prs = Presentation(output)
            self.assertEqual(self_check_layout(prs, preset_id="decision_oriented"), [])


if __name__ == "__main__":
    unittest.main()
