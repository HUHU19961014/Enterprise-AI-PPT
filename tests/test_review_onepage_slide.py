import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from tools.scenario_generators.review_onepage_slide import review_onepage_slide, write_review_report


class ReviewOnepageSlideTests(unittest.TestCase):
    def test_review_passes_clean_generated_slide(self):
        from tools.scenario_generators.build_internal_traceability_uat_logic_slide import build_slide

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "internal_traceability.pptx"
            build_slide(pptx_path, export_review=False)

            report = review_onepage_slide(pptx_path)

            self.assertTrue(report.passed)
            self.assertEqual(report.findings, ())

    def test_review_flags_meta_guidance_phrase(self):
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "meta_phrase.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            textbox = slide.shapes.add_textbox(1000000, 1000000, 5000000, 600000)
            textbox.text_frame.text = "讲解重点：关键用户在UAT前先看懂三件事"
            prs.save(pptx_path)

            report = review_onepage_slide(pptx_path, expected_card_count=0)

            self.assertTrue(any("meta guidance phrase" in finding.message for finding in report.findings))

    def test_write_review_report_outputs_json(self):
        from tools.scenario_generators.build_internal_traceability_uat_logic_slide import build_slide

        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = Path(tmpdir) / "internal_traceability.pptx"
            review_path = Path(tmpdir) / "internal_traceability.review.json"
            build_slide(pptx_path, export_review=False)

            report = review_onepage_slide(pptx_path)
            write_review_report(report, review_path)

            payload = json.loads(review_path.read_text(encoding="utf-8"))
            self.assertTrue(payload["passed"])
            self.assertEqual(payload["findings"], [])


if __name__ == "__main__":
    unittest.main()
