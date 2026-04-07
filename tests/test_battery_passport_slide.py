from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tools import build_sie_battery_passport_slide as slide_builder


class BatteryPassportSlideTests(unittest.TestCase):
    def test_faithful_mode_filters_synthesized_labels(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "battery_passport_flow.pptx"
            built = slide_builder.build_slide(output, layout="flow", narrative_mode="faithful")
            prs = Presentation(str(built))
            slide = prs.slides[0]
            all_text = "\n".join(slide_builder.iter_shape_texts(slide))

        for forbidden in slide_builder.FAITHFUL_FORBIDDEN_TEXTS:
            self.assertNotIn(forbidden, all_text)
        self.assertIn(slide_builder.LEAD_LABEL, all_text)

    def test_output_is_editable_not_picture_based(self) -> None:
        with tempfile.TemporaryDirectory() as tmpdir:
            output = Path(tmpdir) / "battery_passport_decision.pptx"
            built = slide_builder.build_slide(output, layout="decision", narrative_mode="faithful")
            prs = Presentation(str(built))
            slide = prs.slides[0]

        self.assertGreater(len(slide.shapes), 20)
        self.assertFalse(any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes))


if __name__ == "__main__":
    unittest.main()
